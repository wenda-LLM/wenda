"""
funnel
"""

myVersion = "0.1"

__version__ = myVersion

import csv
import io
from md2pptx.rectangle import Rectangle
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor, MSO_THEME_COLOR
from md2pptx.colour import setColour
from md2pptx.symbols import resolveSymbols


def massageFunnelText(text):
    fragment = ""
    for c in text:
        if ord(c) == 236:
            fragment = fragment + "<"

        elif ord(c) == 237:
            fragment = fragment + ">"

        else:
            fragment = fragment + c

    return fragment


class Funnel:
    def __init__(
        self,
    ):
        pass

    def makeFunnel(
        self,
        slide,
        renderingRectangle,
        funnelParts,
        partColours,
        codeType,
        funnelBorderColour,
        funnelTitleColour,
        funnelTextColour,
        funnelLabelsPercent,
        funnelLabelPosition,
        funnelWidest,
    ):
        if funnelWidest in ["left", "right", "pipe", "hpipe"]:
            direction = "horizontal"
        else:
            direction = "vertical"
        
        # Turn label percentage into decimal number to multiply by
        funnelLabelsProportion = funnelLabelsPercent / 100
        
        # Proportion of the stage that the tip - narrowest / shortest
        # part is relative to widest / tallest
        tipProportion = 1 / 3

        # Define labels rectangle then funnel body rectangle
        if direction == "horizontal":
            if funnelLabelPosition == "before":
                # Labels above stages
                funnelLabelsRectangle = Rectangle(
                    renderingRectangle.top,
                    renderingRectangle.left,
                    int(renderingRectangle.height * funnelLabelsProportion),
                    renderingRectangle.width,
                )

                funnelBodyRectangle = Rectangle(
                    renderingRectangle.top
                    + int(renderingRectangle.height * funnelLabelsProportion),
                    renderingRectangle.left,
                    int(renderingRectangle.height * (1 - funnelLabelsProportion)),
                    renderingRectangle.width,
                )
            else:
                # Labels below stages
                funnelLabelsRectangle = Rectangle(
                    renderingRectangle.top
                    + renderingRectangle.height * (1 - funnelLabelsProportion),
                    renderingRectangle.left,
                    int(renderingRectangle.height * funnelLabelsProportion),
                    renderingRectangle.width,
                )

                funnelBodyRectangle = Rectangle(
                    renderingRectangle.top,
                    renderingRectangle.left,
                    int(renderingRectangle.height * (1 - funnelLabelsProportion)),
                    renderingRectangle.width,
                )
        else:
            if funnelLabelPosition == "before":
                # Labels left of stages
                funnelLabelsRectangle = Rectangle(
                    renderingRectangle.top,
                    renderingRectangle.left,
                    renderingRectangle.height,
                    int(renderingRectangle.width * funnelLabelsProportion),
                )

                funnelBodyRectangle = Rectangle(
                    renderingRectangle.top,
                    renderingRectangle.left
                    + int(renderingRectangle.width * funnelLabelsProportion),
                    renderingRectangle.height,
                    int(renderingRectangle.width * (1 - funnelLabelsProportion)),
                )
            else:
                # Labels right of stages
                funnelLabelsRectangle = Rectangle(
                    renderingRectangle.top,
                    renderingRectangle.left
                    + renderingRectangle.width * (1 - funnelLabelsProportion),
                    renderingRectangle.height,
                    int(renderingRectangle.width * funnelLabelsProportion),
                )

                funnelBodyRectangle = Rectangle(
                    renderingRectangle.top,
                    renderingRectangle.left,
                    renderingRectangle.height,
                    int(renderingRectangle.width * (1 - funnelLabelsProportion)),
                )

        if direction == "horizontal":
            # How high the narrowest part of the funnel / pipe is
            tipHeight = funnelBodyRectangle.height * tipProportion
        else:
            # How wide the narrowest part of the funnel / pipe is
            tipWidth = funnelBodyRectangle.width * tipProportion

        partColourCount = len(partColours)

        # Get the underlying rows. Only first 2 cells of each row used
        funnelPartRows = [
            r
            for r in csv.reader(
                io.StringIO(str.join("\n", funnelParts)),
                escapechar="\\",
                skipinitialspace=True,
            )
        ]

        funnelPartCount = len(funnelPartRows)

        # Build lists of labels and Body
        funnelLabels = []
        funnelBody = []

        for row in funnelPartRows:
            cell1 = row[0].strip()
            if len(row) == 0:
                funnelLabels.append("")
                funnelBody.append("")
            else:
                funnelLabels.append(cell1)
                if len(row) == 1:
                    funnelBody.append("")
                else:
                    cell2 = row[1].strip()
                    funnelBody.append(cell2)

        if direction == "horizontal":
            partWidth = renderingRectangle.width / funnelPartCount
        else:
            partHeight = renderingRectangle.height / funnelPartCount

        # Create the labels
        for l, label in enumerate(funnelLabels):
            if direction == "horizontal":
                tb = slide.shapes.add_textbox(
                    funnelLabelsRectangle.left + l * partWidth,
                    funnelLabelsRectangle.top,
                    partWidth,
                    funnelLabelsRectangle.height,
                )
            else:
                tb = slide.shapes.add_textbox(
                    funnelLabelsRectangle.left,
                    funnelLabelsRectangle.top  + l * partHeight,
                    funnelLabelsRectangle.width,
                    partHeight,
                )

            tb.text = massageFunnelText(resolveSymbols(label.replace("<br/>", "\n")))
            for p in tb.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                if funnelTitleColour != ("None", ""):
                    setColour(p.font.color, funnelTitleColour)

        # Create the parts of the funnel

        for b, body in enumerate(funnelBody):
            if direction == "horizontal":
                # Horizontal stages
                
                # Left extremity of stage - both top and bottom
                partLeft = funnelBodyRectangle.left + b * partWidth

                # Right extremity of stage - both top and bottom
                partRight = partLeft + partWidth

                if (
                    ((b == funnelPartCount - 1) & (funnelWidest == "left"))
                    | ((b == 0) & (funnelWidest == "right"))
                    | (funnelWidest in ["pipe", "hpipe"])
                ):
                    # Rectangular / pipe stage

                    # Calculate space to be above and below rectangular stage
                    leftSpaceAboveBelow = (funnelBodyRectangle.height - tipHeight) / 2
                    rightSpaceAboveBelow = leftSpaceAboveBelow
                    
                elif funnelWidest == "left":
                    # Calculate space to be above and below left end of stage
                    leftSpaceAboveBelow = (
                        (funnelBodyRectangle.height - tipHeight)
                        / 2
                        * b
                        / (funnelPartCount - 1)
                    )

                    # Calculate space to be above and below right end of stage
                    rightSpaceAboveBelow = (
                        (funnelBodyRectangle.height - tipHeight)
                        / 2
                        * (b + 1)
                        / (funnelPartCount - 1)
                    )

                else:
                    # Calculate space to be above and below left end of stage
                    leftSpaceAboveBelow = (
                        (funnelBodyRectangle.height - tipHeight)
                        / 2
                        * (funnelPartCount - b)
                        / (funnelPartCount - 1)
                    )

                    # Calculate space to be above and below right end of stage
                    rightSpaceAboveBelow = (
                        (funnelBodyRectangle.height - tipHeight)
                        / 2
                        * (funnelPartCount - b - 1)
                        / (funnelPartCount - 1)
                    )

                # Calculate y coordinate of top left corner
                partTopLeft = funnelBodyRectangle.top + leftSpaceAboveBelow

                # Calculate y coordinate of top right corner
                partTopRight = funnelBodyRectangle.top + rightSpaceAboveBelow

                # Calculate y coordinate of bottom left corner
                partBottomLeft = (
                    funnelBodyRectangle.top
                    + funnelBodyRectangle.height
                    - leftSpaceAboveBelow
                )

                # Calculate y coordinate of bottom right corner
                partBottomRight = (
                    funnelBodyRectangle.top
                    + funnelBodyRectangle.height
                    - rightSpaceAboveBelow
                )

                stagePoints = [
                    (partLeft, partTopLeft),
                    (partLeft, partBottomLeft),
                    (partRight, partBottomRight),
                    (partRight, partTopRight),
                ]
            else:
                # Vertical stages
                
                # Top extremity of stage - both left and right
                partTop = funnelBodyRectangle.top + b * partHeight

                # Bottom extremity of stage - both left and right
                partBottom = partTop + partHeight

                if (
                    ((b == funnelPartCount - 1) & (funnelWidest == "top"))
                    | ((b == 0) & (funnelWidest == "bottom"))
                    | (funnelWidest == "vpipe")
                ):
                    # Rectangular / pipe stage
                    
                    # Calculate space to be left either side of rectangular stage
                    topSpaceLeftRight = (funnelBodyRectangle.width - tipWidth) / 2
                    bottomSpaceLeftRight = topSpaceLeftRight

                elif funnelWidest == "top":
                    # Calculate space to be left at left and right top end of stage
                    topSpaceLeftRight = (
                        (funnelBodyRectangle.width - tipWidth)
                        / 2
                        * b
                        / (funnelPartCount - 1)
                    )

                    # Calculate space to be left at left and right bottom end of stage
                    bottomSpaceLeftRight = (
                        (funnelBodyRectangle.width - tipWidth)
                        / 2
                        * (b + 1)
                        / (funnelPartCount - 1)
                    )

                else:
                    # Calculate space to be left at left and right top end of stage
                    topSpaceLeftRight = (
                        (funnelBodyRectangle.width - tipWidth)
                        / 2
                        * (funnelPartCount - b)
                        / (funnelPartCount - 1)
                    )

                    # Calculate space to be left at left and right bottom end of stage
                    bottomSpaceLeftRight = (
                        (funnelBodyRectangle.width - tipWidth)
                        / 2
                        * (funnelPartCount - b - 1)
                        / (funnelPartCount - 1)
                    )

                # Calculate x coordinate of top left corner
                partTopLeft = funnelBodyRectangle.left + topSpaceLeftRight

                # Calculate x coordinate of bottom left corner
                partBottomLeft = funnelBodyRectangle.left + bottomSpaceLeftRight

                # Calculate x coordinate of top right corner
                partTopRight = (
                    funnelBodyRectangle.left
                    + funnelBodyRectangle.width
                    - topSpaceLeftRight
                )

                # Calculate x coordinate of bottom right corner
                partBottomRight = (
                    funnelBodyRectangle.left
                    + funnelBodyRectangle.width
                    - bottomSpaceLeftRight
                )

                stagePoints = [
                    (partTopLeft, partTop),
                    (partTopRight, partTop),
                    (partBottomRight, partBottom),
                    (partBottomLeft, partBottom),
                ]

            # Start shape builder with first point
            ffBuilder = slide.shapes.build_freeform(*stagePoints[0])

            ffBuilder.add_line_segments(stagePoints[1:],
                close=True,
            )

            s = ffBuilder.convert_to_shape()
            s.text = massageFunnelText(resolveSymbols(body.replace("<br/>", "\n")))

            for p in s.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                if funnelTextColour != ("None", ""):
                    setColour(p.font.color, funnelTextColour)

            s.fill.solid()

            partColourType, partColourValue = partColours[b % partColourCount]
            if partColourType == "Theme":
                s.fill.fore_color.theme_color = partColourValue
            else:
                s.fill.fore_color.rgb = RGBColor.from_string(partColourValue[1:])

            if funnelBorderColour != ("None", ""):
                setColour(s.line.color, funnelBorderColour)
