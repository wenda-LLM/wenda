import re
import sys
import os
import time
import collections
import collections.abc
from pptx import Presentation
from pptx import __version__ as pptx_version
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor, MSO_THEME_COLOR
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.action import PP_ACTION
from pptx.enum.dml import MSO_PATTERN

import imghdr, struct
import datetime
import html.parser
from pptx.oxml.xmlchemy import OxmlElement
from pathlib import Path
import urllib.request
import tempfile
import copy
import platform
import shutil
import socket
from pptx.oxml import parse_xml
import uuid
import md2pptx.funnel as funnel
from md2pptx.rectangle import Rectangle
from md2pptx.colour import *
from md2pptx.symbols import resolveSymbols
import random

from lxml import etree
from lxml.html import fromstring


def readinputfile(input_filename):
    if Path(input_filename).exists():
        input_path = Path(input_filename)
        with input_path.open(encoding='utf-8') as file:
            input_file = file.readlines()
            input_file = [i.strip() for i in input_file if i.strip()]
            return input_file
    else:
        print("Input file specified but does not exist. Terminating.")

def writeMdFile(content):
    now = time.strftime("%Y-%m-%d-%H_%M_%S", time.localtime(time.time()))
    filename =os.getcwd() + os.sep + 'md2pptx' + os.sep + 'data' + os.sep + 'md' + now + '.md'
    # if not os.path.exists(filename):
    #     os.mknod(filename)
    print(os.getcwd())
    with open(filename, 'w', encoding='utf-8') as file:
        file.write(content)
    return filename

def genPPTX(outputfile, inputfile):
    prs = Presentation("CHATOVENS.pptx")
    # slide = prs.slides.add_slide(prs.slide_layouts[4])
    # for shape in slide.placeholders:         # 获取这一页所有的占位符
    #     phf = shape.placeholder_format
    #     print(f'{phf.idx}--{shape.name}--{phf.type}')

    input_filename = inputfile
    input_file = readinputfile(input_filename)
    #匹配# 开头的行 首页
    metadataRegex_header = re.compile("^#\s")
    #匹配## 开头的行 一级标题
    metadataRegex_one = re.compile("^##\s")
    #匹配### 开头的行 二级标题和目录
    metadataRegex_two = re.compile("^###\s")
    #匹配* 开头的行 纯内容
    metadataRegex_content = re.compile("^[\*\s|\s+\s]")
    #匹配纯内容
    metadataRegex_noformat = re.compile("^[^\*\s|^#\s|^##\s|^###\s|^+\s]")
    # 匹配纯标题
    metadataRegex_format = re.compile("^[#\s|##\s|###\s|####\s]")
    #匹配数字开头
    metadataRegex_numberformat = re.compile("^\d.")
    #0就是不生成新页，1就是生成新页
    sildetag = 0
    #0就是不增加新页，1就是增加新页
    mutilslidetag = 0
    #目录的内容
    mutilslidetitle = ''
    #当前页
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    #模板页id
    slideid = 0
    #模板页中shape索引
    slideshapeindex = 1
    #目录数
    contentcolindex = 1
    # 内容数量
    contentnumber = 0
    # 内容
    content_ = ''
    for index, line in enumerate(input_file):
        #内容多少选择不同的模板页面
        if metadataRegex_format.match(line) or index==len(input_file)-1:
            if contentnumber !=0 and contentnumber<=4:
                slide = prs.slides.add_slide(prs.slide_layouts[random.sample(range(4,6), 1)[0]])
                slide.shapes[0].text = mutilslidetitle
                slide.shapes[1].text = content_
                slide.shapes[2].insert_picture('fx.png')
            if contentnumber !=0 and contentnumber>4:
                slide.placeholders[0].text = mutilslidetitle
                slide = prs.slides.add_slide(prs.slide_layouts[3])
                slide.placeholders[13].text = content_
            contentnumber = 0
            content_ = ''
        #匹配首页
        if metadataRegex_header.match(line):
            contentnumber = 0
            slideshapeindex = 1
            # slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.placeholders[0]
            title.text = re.sub(metadataRegex_header, '', line, 1)
            mutilslidetitle = title.text
            slideid = 0
        #内容
        if metadataRegex_content.match(line):
            if mutilslidetag != 0 :
                slide = prs.slides.add_slide(slideid)
            if mutilslidetitle == '目录':
                slide.shapes[slideshapeindex].text = re.sub(metadataRegex_content, '', line, 1)
                slideshapeindex += 1
            else:
                content_ = content_ + re.sub(metadataRegex_content, '', line, 1) + '\n'
                contentnumber += 1
                # slide.placeholders[13].text = slide.placeholders[13].text + re.sub(metadataRegex_content, '', line, 1) + '\n'
        # 匹配主内容区
        if metadataRegex_two.match(line):
            contentnumber = 0
            slideshapeindex = 1
            mutilslidetitle = re.sub(metadataRegex_two, '', line, 1)
            if mutilslidetitle == '目录':
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = slide.placeholders[0]
                title.text = re.sub(metadataRegex_two, '', line, 1)
        #匹配二级目录
        if metadataRegex_one.match(line):
            contentnumber = 0
            slideshapeindex = 1
            mutilslidetitle = re.sub(metadataRegex_one, '', line, 1)
            slide = prs.slides.add_slide(prs.slide_layouts[2])
            title = slide.placeholders[0]
            title.text = re.sub(metadataRegex_numberformat, '', re.sub(metadataRegex_one, '', line, 1), 1)
            title = slide.placeholders[10]
            title.text = '0' + str(contentcolindex)
            contentcolindex += 1
        # if metadataRegex_noformat.match(line):
        #     slide.placeholders[2].text = slide.placeholders[2].text + re.sub(metadataRegex_content, '', line, 1) + '\n'

    slide = prs.slides.add_slide(prs.slide_layouts[12])
    title = slide.placeholders[0]
    title.text = '结束'
    # # 修改标题
    # title = slide.placeholders[0]
    # title.text = "这里是标题"

    # col_name = [['hh', 'this'], ['A', 'B']]
    #
    # table = slide.placeholders[10]
    #
    # rows, cols = 2, 2
    # table0 = table.insert_table(rows, cols).table
    #
    # for row in range(0, rows):
    #     for col in range(0, cols):
    #         table0.cell(row, col).text = col_name[row][col]
    prs.save(outputfile)
    return outputfile

if __name__ == '__main__':
    genPPTX('text.pptx', '../test.md')