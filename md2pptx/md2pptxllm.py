#!/usr/bin/env python3

"""
md2pptx - Converts (a subset of) Markdown to Powerpoint (PPTX)

First argument is file to write to

Reads from stdin
"""

import re
import sys
import os
import time
import collections
import collections.abc
from pptx import Presentation
from pptx import __version__ as pptx_version
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor, MSO_THEME_COLOR
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.action import PP_ACTION
from pptx.enum.dml import MSO_PATTERN

import imghdr, struct
import datetime
import html.parser
from pptx.oxml.xmlchemy import OxmlElement
from pathlib import Path
import urllib.request
import tempfile
import copy
import platform
import shutil
import socket
from pptx.oxml import parse_xml
import uuid
import md2pptx.funnel as funnel
from md2pptx.rectangle import Rectangle
from md2pptx.colour import *
from md2pptx.symbols import resolveSymbols

from lxml import etree
from lxml.html import fromstring

# Try to import CairoSVG - which might not be installed.
# Flag availability or otherwise
try:
    import cairosvg
    from cairosvg import helpers

    have_cairosvg = True

except:
    have_cairosvg = False

# Try to import Pillow - which might not be installed.
# Flag availability or otherwise
try:
    import PIL

    have_pillow = True

except:
    have_pillow = False

# Try to import graphviz - which might not be installed.
# Flag availability or otherwise

try:
    import graphviz

    have_graphviz = True

except:
    have_graphviz = False


md2pptx_level = "4.1.2"
md2pptx_date = "11 July, 2023"

namespaceURL = {
    "mc": "http://schemas.openxmlformats.org/markup-compatibility/2006",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "p14": "http://schemas.microsoft.com/office/powerpoint/2010/main",
    "p15": "http://schemas.microsoft.com/office/powerpoint/2012/main",
}


class SlideInfo:
    def __init__(
        self,
        titleText,
        subtitleText,
        blockType,
        bullets,
        tableRows,
        cards,
        code,
        sequence,
    ):
        self.titleText = titleText
        self.subtitleText = subtitleText
        self.blockType = blockType
        self.bullets = bullets
        self.tableRows = tableRows
        self.cards = cards
        self.code = code
        self.sequence = sequence


# Information about a single table. (A slide might have more than one - or none.)
class TableInfo:
    def __init__(self, tableRows, tableCaption):
        self.tableRows = []
        self.tableCaption = ""


# Information about a video
class AudioVideoInfo:
    def __init__(self, elementString):
        audioVideoElement = fromstring(elementString)

        if "src" in audioVideoElement.attrib.keys():
            self.source = audioVideoElement.attrib["src"]
        else:
            self.source = ""

        if audioVideoElement.tag == "audio":
            # audio element doesn't have width or height attributes so make it square
            self.width = 1024
            self.height = 1024
        else:
            # video element can have width and height attributes
            # Default is 4:3
            if "width" in audioVideoElement.attrib.keys():
                self.width = int(audioVideoElement.attrib["width"])
            else:
                self.width = 1024

            if "height" in audioVideoElement.attrib.keys():
                self.height = int(audioVideoElement.attrib["height"])
            else:
                self.height = 768

        self.aspectRatio = self.width / self.height

        if "poster" in audioVideoElement.attrib.keys():
            self.poster = audioVideoElement.attrib["poster"]
        else:
            self.poster = None


# Note: Options stored with lower case keys
class ProcessingOptions:
    def __init__(self):
        self.defaultOptions = {}
        self.presentationOptions = {}
        self.currentOptions = {}
        self.dynamicallyChangedOptions = {}
        self.hideMetadataStyle = False

    def getDefaultOption(self, optionName):
        return self.defaultOptions[optionName.lower()]

    def setDefaultOption(self, optionName, value):
        self.defaultOptions[optionName.lower()] = value

    def getPresentationOption(self, optionName):
        return self.presentationOptions[optionName.lower()]

    def setPresentationOption(self, optionName, value):
        self.presentationOptions[optionName.lower()] = value

    def getCurrentOption(self, optionName):
        return self.currentOptions[optionName.lower()][-1]

    # Note: Can't pop to an empty stack. Will always have default available
    def popCurrentOption(self, optionName):
        key = optionName.lower()
        if len(self.currentOptions[key]) > 1:
            # Have a non-default value to use
            self.currentOptions[key].pop()

    def setCurrentOption(self, optionName, value):
        key = optionName.lower()

        if key in self.currentOptions:
            # Add new value to existing stack
            self.currentOptions[key].append(value)
        else:
            # Start a new stack for this option
            self.currentOptions[key] = [value]

    def setOptionValues(self, optionName, value):
        key = optionName.lower()
        if key not in self.defaultOptions:
            self.setDefaultOption(optionName, value)

        self.setPresentationOption(optionName, value)

        self.setCurrentOption(optionName, value)

    def setOptionValuesArray(self, optionArray):
        for keyValuePair in optionArray:
            self.setOptionValues(keyValuePair[0], keyValuePair[1])

    def dynamicallySetOption(self, optionName, optionValue, conversion):
        lowerName = optionName.lower()
        if optionValue == "default":
            self.setCurrentOption(lowerName, self.getDefaultOption(lowerName))

        elif optionValue == "pres":
            self.setCurrentOption(lowerName, self.getPresentationOption(lowerName))

        elif optionValue in ["pop", "prev"]:
            self.popCurrentOption(lowerName)

        elif conversion == "":
            self.setCurrentOption(lowerName, optionValue)

        elif conversion == "float":
            self.setCurrentOption(lowerName, float(optionValue))

        elif conversion == "sortednumericlist":
            self.setCurrentOption(lowerName, sortedNumericList(optionValue))

        elif conversion == "int":
            self.setCurrentOption(lowerName, int(optionValue))

        self.dynamicallyChangedOptions[lowerName] = True


# Find the extLst element - if it exists in presentation.xml
def findExtLst(prs):
    for child in prs._element.getchildren():
        if child.tag.endswith("}extLst"):
            return child

    return None


def addSlide(presentation, slideLayout, slideInfo=None):
    slide = presentation.slides.add_slide(slideLayout)
    slide.slideInfo = slideInfo

    return slide


def createSectionsXML(prs):
    sectionSlideLayout = processingOptions.getCurrentOption("SectionSlideLayout")

    xml = '  <p:ext xmlns:p="' + namespaceURL["p"] + '"\n'

    # ext URI has to be {521415D9-36F7-43E2-AB2F-B90AF26B5E84} as it's a registered extension
    xml += '    uri="{521415D9-36F7-43E2-AB2F-B90AF26B5E84}">\n'

    xml += '    <p14:sectionLst xmlns:p14="' + namespaceURL["p14"] + '">\n'

    sectionCount = 0
    for slide in prs.slides:
        slideID = str(slide.slide_id)

        for idx, slide_layout in enumerate(prs.slide_layouts):
            if slide.slide_layout == slide_layout:
                layoutNumber = idx
                break

        if layoutNumber == sectionSlideLayout:
            # Have a section to contribute
            sectionCount += 1

            # Confect section name from first part of section slide title
            title = findTitleShape(slide)
            sectionName = title.text.split("\n")[0]

            # Clean up section name
            sectionName = "".join(
                letter
                for letter in sectionName
                if (
                    (letter.isalnum())
                    | (letter in "&-+")
                    | (letter in "!/*")
                    | (letter == " ")
                )
            )

            sectionName = (
                sectionName.replace("& ", "&amp; ")
                .replace("\r", " ")
                .replace("\n", " ")
            )

            # section URI's just need to be a GUID wrapped in braces
            xml += (
                '      <p14:section name="'
                + sectionName
                + '" id="{'
                + str(uuid.uuid4()).upper()
                + '}">\n'
            )

            # Only the first slide in the section is added - as section will continue until the next section
            # anyway
            xml += "        <p14:sldIdLst>\n"
            xml += '          <p14:sldId id="' + slideID + '" />\n'
            xml += "        </p14:sldIdLst>\n"
            xml += "      </p14:section>\n"

    # Close out the section list
    xml += "    </p14:sectionLst>\n"

    # Close out the sections extension
    xml += "  </p:ext>\n"

    parsed_xml = parse_xml(xml)

    return parsed_xml, sectionCount


def createExpandingSections(prs):
    # Use the slides' layouts to create an XML fragment with sections in
    xmlFragment, sectionCount = createSectionsXML(prs)

    if sectionCount > 0:
        # Have sections to insert as an XML fragment
        if (extLstElement := findExtLst(prs)) is not None:
            # Need to remove the extension list element before adding a new one
            prs._element.remove(extLstElement)

        # Insert a new extension list element
        extLst = OxmlElement("p:extLst")
        prs._element.insert(-1, extLst)

        # Insert the fragment in the extension list in presentation.xml
        extLst.insert(0, xmlFragment)


def deleteSlide(prs, slideNumber):
    rId = prs.slides._sldIdLst[slideNumber].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[slideNumber]


def startswithOneOf(haystack, needleList):
    for needle in needleList:
        if haystack.startswith(needle):
            return True

    return False


# Splits a string into words, converting each word to an integer. Returns them as a
# sorted list
def sortedNumericList(string):
    return sorted(list(map(int, set(string.split()))))


def substituteFooterVariables(footerText, liveFooters):
    # Decide if the footer should be a live link to the section slide
    wantLiveFooter = (
        (prs.lastSectionSlide is not None)
        & (footerText.find("<section") > -1)
        & (liveFooters == "yes")
    )

    # Substitute any section title occurrences
    sectionTitleLines = prs.lastSectionTitle.split("<br/>")

    footerText = footerText.replace("<section>", sectionTitleLines[0])
    footerText = footerText.replace("<section1>", sectionTitleLines[0])

    if len(sectionTitleLines) > 1:
        footerText = footerText.replace("<section2>", sectionTitleLines[1])

    if len(sectionTitleLines) > 2:
        footerText = footerText.replace("<section3>", sectionTitleLines[2])

    # Substitute any presentation title occurrences
    presTitleLines = prs.lastPresTitle.split("<br/>")
    footerText = footerText.replace("<presTitle>", presTitleLines[0])
    footerText = footerText.replace("<presTitle1>", presTitleLines[0])

    if len(presTitleLines) > 1:
        footerText = footerText.replace("<presTitle2>", presTitleLines[1])

    if len(presTitleLines) > 2:
        footerText = footerText.replace("<presTitle3>", presTitleLines[2])

    # Substitute any presentation subtitle occurrences
    presSubtitleLines = prs.lastPresSubtitle.split("<br/>")

    footerText = footerText.replace("<presSubtitle>", presSubtitleLines[0])
    footerText = footerText.replace("<presSubtitle1>", presSubtitleLines[0])

    if len(presSubtitleLines) > 1:
        footerText = footerText.replace("<presSubtitle2>", presSubtitleLines[1])

    if len(presSubtitleLines) > 2:
        footerText = footerText.replace("<presSubtitle3>", presSubtitleLines[2])

    # Make newlines happen
    footerText = footerText.replace("<br/>", "\n")

    return footerText, wantLiveFooter


def _applyCellBorderStyling(
    tcPr, linePosition, lineWidthMultiplier=1, lineCount=1, lineColour="000000"
):
    # How wide, relatively speaking to make the lines
    lineWidth = int(12700 * lineWidthMultiplier)

    # Whether the line should be single or double
    if lineCount == 2:
        lineCountValue = "dbl"
    else:
        lineCountValue = "sng"

    if linePosition == "l":
        elementName = "a:lnL"
    elif linePosition == "r":
        elementName = "a:lnR"
    elif linePosition == "t":
        elementName = "a:lnT"
    else:
        elementName = "a:lnB"

    lnX = OxmlElement(elementName)

    lnX.attrib.update(
        {"w": str(lineWidth), "cap": "flat", "cmpd": lineCountValue, "algn": "ctr"}
    )

    solidFill = OxmlElement("a:solidFill")
    srgbClr = OxmlElement("a:srgbClr")
    srgbClr.attrib.update({"val": lineColour})

    solidFill.append(srgbClr)
    lnX.append(solidFill)

    tcPr.append(lnX)


def applyCellBorderStyling(
    cell, cellBorderStyling, lineWidthMultiplier, lineCount, lineColour
):
    if cellBorderStyling == "":
        # No cell border styling required
        return

    # Get any existing cell properties element - or make one
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    # Draw any cell borders. More than one might apply
    if cellBorderStyling.find("l") > -1:
        _applyCellBorderStyling(tcPr, "l", lineWidthMultiplier, lineCount, lineColour)
    if cellBorderStyling.find("r") > -1:
        _applyCellBorderStyling(tcPr, "r", lineWidthMultiplier, lineCount, lineColour)
    if cellBorderStyling.find("t") > -1:
        _applyCellBorderStyling(tcPr, "t", lineWidthMultiplier, lineCount, lineColour)
    if cellBorderStyling.find("b") > -1:
        _applyCellBorderStyling(tcPr, "b", lineWidthMultiplier, lineCount, lineColour)


# Apply table line styling
def applyTableLineStyling(
    table,
    processingOptions,
):
    wholeTableLineStyling = processingOptions.getCurrentOption("addTableLines")
    linedColumns = processingOptions.getCurrentOption("addTableColumnLines")
    linedRows = processingOptions.getCurrentOption("addTableRowLines")

    lastRow = len(table.rows) - 1

    # Create blank cell styling matrix
    cellStyling = []
    for rowNumber, row in enumerate(table.rows):
        rowStyling = []
        for cell in row.cells:
            rowStyling.append("")
        cellStyling.append(rowStyling)

    # apply any "whole table" styling - from addTableLines
    if wholeTableLineStyling == "box":
        # Line around the table
        for rowNumber, row in enumerate(table.rows):
            # Figure out whether row is top, middle, or bottom
            if rowNumber == 0:
                rowStyling = "t"
            elif rowNumber == lastRow:
                rowStyling = "b"
            else:
                rowStyling = ""

            lastColumn = len(row.cells) - 1

            for columnNumber, cell in enumerate(row.cells):
                if columnNumber == 0:
                    columnStyling = "l"
                elif columnNumber == lastColumn:
                    columnStyling = "r"
                else:
                    columnStyling = ""
                cellStyling[rowNumber][columnNumber] = rowStyling + columnStyling

    elif wholeTableLineStyling == "all":
        # All edges of all cells have lines
        for rowNumber, row in enumerate(table.rows):
            lastColumn = len(row.cells) - 1

            for columnNumber, cell in enumerate(row.cells):
                cellStyling[rowNumber][columnNumber] = "tlbr"

    # Apply any row styling - from addTableColumnLines
    for rowNumber, row in enumerate(table.rows):
        if rowNumber + 1 in linedRows:
            # Line after this row so below
            for columnNumber, cell in enumerate(row.cells):
                cellStyling[rowNumber][columnNumber] = (
                    cellStyling[rowNumber][columnNumber] + "b"
                )

        elif rowNumber in linedRows:
            # Line before this row so above
            for columnNumber, cell in enumerate(row.cells):
                cellStyling[rowNumber][columnNumber] = (
                    cellStyling[rowNumber][columnNumber] + "t"
                )

    # Apply any column styling - from addTableRowLines
    for rowNumber, row in enumerate(table.rows):
        for columnNumber, cell in enumerate(row.cells):
            if columnNumber + 1 in linedColumns:
                # Line after this column so to right
                cellStyling[rowNumber][columnNumber] = (
                    cellStyling[rowNumber][columnNumber] + "r"
                )

            elif columnNumber + 1 in linedColumns:
                # Line after this column so to left
                cellStyling[rowNumber][columnNumber] = (
                    cellStyling[rowNumber][columnNumber] + "r"
                )

    # Apply the styling from the matrix to all cells
    for rowNumber, row in enumerate(table.rows):
        for columnNumber, cell in enumerate(row.cells):
            applyCellBorderStyling(
                cell,
                cellStyling[rowNumber][columnNumber],
                processingOptions.getCurrentOption("addTableLineWidth"),
                processingOptions.getCurrentOption("addTableLineCount"),
                processingOptions.getCurrentOption("addTableLineColour"),
            )


def handleSpanClass(run, className):
    if className in bgcolors:
        run = set_highlight(run, bgcolors[className])

    if className in fgcolors:
        font = run.font
        font.color.rgb = RGBColor.from_string(fgcolors[className])

    if className in emphases:
        if " bold " in " " + emphases[className] + " ":
            font.bold = True
        else:
            font.bold = False
        if " italic " in " " + emphases[className] + " ":
            font.italic = True
        else:
            font.italic = False
        if " underline " in " " + emphases[className] + " ":
            font.underline = True
        else:
            font.underline = False


def handleSpanStyle(run, styleText):
    styleElements = styleText.split(";")

    # Handle the non-empty ones - as the empty one is after the final semicolon
    for styleElement in list(filter(lambda e: e != "", styleElements)):
        styleElementSplit = styleElement.split(":")
        styleElementName = styleElementSplit[0].strip()
        styleElementValue = styleElementSplit[1].strip()

        if styleElementName == "color":
            check, RGBstring = parseRGB(styleElementValue)
            if check:
                run.font.color.rgb = RGBColor.from_string(RGBstring)
            else:
                print(f"Invalid {styleElementName} RGB value {styleElementValue}")

        elif styleElementName == "background-color":
            check, RGBstring = parseRGB(styleElementValue)
            if check:
                set_highlight(run, RGBstring)
            else:
                print(f"Invalid {styleElementName} RGB value {styleElementValue}")

        elif styleElementName == "text-decoration":
            if styleElementValue == "underline":
                run.font.underline = True

        elif styleElementName == "font-weight":
            if styleElementValue == "bold":
                run.font.bold = True

        elif styleElementName == "font-style":
            if styleElementValue == "italic":
                run.font.italic = True


def reportSlideTitle(slideNumber, indent, titleText):
    print(str(slideNumber).rjust(4) + " " + ("  " * indent) + titleText)


def reportGraphicFilenames(leftFilename, rightFilename=""):
    if rightFilename == "":
        print("             ---> " + leftFilename.ljust(30))
    else:
        print("             ---> " + leftFilename.ljust(30) + " , " + rightFilename)


# Given current indenting regime calculate what level the bullet / number is at
def calculateIndentationLevel(firstNonSpace, indentSpaces):
    return int(firstNonSpace / indentSpaces)


# Calculate picture dimensions given its natural height and bounds
def scalePicture(maxPicWidth, maxPicHeight, imageWidth, imageHeight):
    heightIfWidthUsed = maxPicWidth * imageHeight / imageWidth
    widthIfHeightUsed = maxPicHeight * imageWidth / imageHeight

    if heightIfWidthUsed > maxPicHeight:
        # Use the height to scale
        usingHeightToScale = True

        picWidth = widthIfHeightUsed
        picHeight = maxPicHeight

    else:
        # Use the width to scale
        usingHeightToScale = False

        picWidth = maxPicWidth
        picHeight = heightIfWidthUsed
    return (picWidth, picHeight, usingHeightToScale)


def parseMedia(cellString, graphicCount):
    graphicTitle = ""
    HTML = ""
    audioVideoInfo = None
    graphicHref = ""
    GraphicFilename = ""
    printableGraphicFilename = ""

    graphicCount += 1

    if videoRegexMatch := videoRegex.match(cellString):
        # Cell contains a video
        audioVideoInfo = AudioVideoInfo(cellString)
        _, printableGraphicFilename = handleWhateverGraphicType(audioVideoInfo.source)

    if audioRegexMatch := audioRegex.match(cellString):
        # Cell contains an audio
        audioVideoInfo = AudioVideoInfo(cellString)
        _, printableGraphicFilename = handleWhateverGraphicType(audioVideoInfo.source)

    elif clickableGraphicMatch := clickableGraphicRegex.match(cellString):
        # Cell contains a clickable graphic
        graphicTitle = clickableGraphicMatch.group(1)
        GraphicFilename = clickableGraphicMatch.group(2)
        graphicHref = clickableGraphicMatch.group(3)

        (
            GraphicFilename,
            printableGraphicFilename,
        ) = handleWhateverGraphicType(GraphicFilename)

    elif graphicMatch := graphicRegex.match(cellString):
        # Cell contains a non-clickable graphic
        graphicTitle = graphicMatch.group(1)
        GraphicFilename = graphicMatch.group(2)

        (
            GraphicFilename,
            printableGraphicFilename,
        ) = handleWhateverGraphicType(GraphicFilename)

    else:
        # Not a graphic or video
        GraphicFilename = ""
        printableGraphicFilename = ""
        HTML = cellString
        graphicCount -= 1

    return (
        graphicTitle,
        GraphicFilename,
        printableGraphicFilename,
        graphicHref,
        HTML,
        audioVideoInfo,
        graphicCount,
    )


# Send a shape to the back on a slide
def sendToBack(shapes, shape):
    firstShapeElement = shapes[0]._element
    firstShapeElement.addprevious(shape._element)


# Turn a paragraph into a numbered inList item
def makeNumberedListItem(p):
    if (
        p._element.getchildren()[0].tag
        == "{http://schemas.openxmlformats.org/drawingml/2006/main}pPr"
    ):
        pPr = p._element.getchildren()[0]
        if len(pPr.getchildren()) > 0:
            # Remove Default Text Run Properties element - if present
            x = pPr.getchildren()[0]
            if x.tag == "{http://schemas.openxmlformats.org/drawingml/2006/main}defRPr":
                pPr.remove(x)
    else:
        pPr = OxmlElement("a:pPr")
        p._element.insert(0, pPr)

    buFont = OxmlElement("a:buFont")
    buFont.set("typeface", "+mj-lt")
    pPr.append(buFont)

    buAutoNum = OxmlElement("a:buAutoNum")
    buAutoNum.set("type", "arabicPeriod")
    pPr.append(buAutoNum)


# Add a drop shadow to a shape
def createShadow(shape):
    spPr = shape.fill._xPr

    el = OxmlElement("a:effectLst")
    spPr.append(el)

    outerShdw = OxmlElement("a:outerShdw")
    outerShdw.set("algn", "tl")
    outerShdw.set("blurRad", "50800")
    outerShdw.set("dir", "2700000")
    outerShdw.set("dist", "95250")
    outerShdw.set("rotWithShape", "0")

    el.append(outerShdw)

    prstClr = OxmlElement("a:prstClr")
    prstClr.set("val", "black")

    outerShdw.append(prstClr)

    alpha = OxmlElement("a:alpha")
    alpha.set("val", "40000")

    prstClr.append(alpha)


# Clone a shape in a slide and return the new shape.
# (This is a deep copy so the new shape will have the same
# eg bullet style as the source shape)
def addClonedShape(slide, shape1):
    # Clone the element for the shape
    el1 = shape1.element
    el2 = copy.deepcopy(el1)

    # Insert the cloned element into the shape tree
    slide.shapes._spTree.insert_element_before(el2, "p:extLst")

    # Return the shape associated with this new element
    return slide.shapes[-1]


# Following functions are workarounds for python-pptx not having these functions for the font object
def set_subscript(font):
    if font.size is None:
        font._element.set("baseline", "-50000")
        return
    
    if font.size < Pt(24):
        font._element.set("baseline", "-50000")
    else:
        font._element.set("baseline", "-25000")


def set_superscript(font):
    if font.size is None:
        font._element.set("baseline", "60000")
        return
    
    if font.size < Pt(24):
        font._element.set("baseline", "60000")
    else:
        font._element.set("baseline", "30000")


def set_strikethrough(font):
    font._element.set("strike", "sngStrike")


def set_highlight(run, color):
    # get run properties
    rPr = run._r.get_or_add_rPr()

    # Create highlight element
    hl = OxmlElement("a:highlight")

    # Create specify RGB Colour element with color specified
    srgbClr = OxmlElement("a:srgbClr")
    setattr(srgbClr, "val", color)

    # Add colour specification to highlight element
    hl.append(srgbClr)

    # Add highlight element to run properties
    rPr.append(hl)

    return run


def removeBullet(paragraph):
    pPr = paragraph._p.get_or_add_pPr()
    pPr.insert(
        0,
        etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/main}buNone"),
    )


def removeBullets(textFrame):
    for p in textFrame.paragraphs:
        removeBullet(p)


# Get the slide object the run is in
def SlideFromRun(run):
    return run._parent._parent._parent._parent._parent


# Get the slide object the picture is in
def SlideFromPicture(picture):
    return picture._parent._parent


# Creates a hyperlink to another slide and/or a tooltip - for a
# text run
# Note: To get just a tooltip make to_slide be the source slide
#       so it links to itself.
def createRunHyperlinkOrTooltip(run, to_slide, tooltipText=""):
    # Get hold of the shape the run is in
    if run._parent._parent._parent.__class__.__name__ == "_Cell":
        # Run in a table cell has to be handled differently
        shape = (
            run._parent._parent._parent._parent._parent._parent._parent._graphic_frame
        )
    else:
        # Ordinary text run
        shape = run._parent._parent._parent

    if to_slide == None:
        to_slide = SlideFromRun(run)
    hl = run.hyperlink
    sca = shape.click_action
    sca_hl = sca.hyperlink

    # Add a click action to generate an internal hyperlink address
    sca.target_slide = to_slide

    # Use that internal hyperlink address for the run
    hl.address = sca_hl.address

    # Also clone the hyperlink click action
    hl._hlinkClick.action = sca_hl._hlink.action

    if tooltipText != "":
        hl._hlinkClick.set("tooltip", tooltipText)

    # Also clone the hyperlink rId
    hl._hlinkClick.rId = sca_hl._hlink.rId

    # Delete the shape click action
    sca.target_slide = None


# Creates a hyperlink to another slide or a URL and/or a tooltip - for a
# picture
# Note: To get just a tooltip make to_slide be the source slide
#       so it links to itself.
def createPictureHyperlinkOrTooltip(picture, target, tooltipText=""):
    if target == None:
        # If neither a tooltip nor a target slide then return having
        # done nothing
        if tooltipText == "":
            return

        # Tooltip but no target slide
        target = SlideFromPicture(picture)
        picture.click_action.target_slide = target
    elif target.__class__.__name__ == "str":
        # Is a URL
        picture.click_action.hyperlink.address = target
        # URL might be a macro reference
        if target[:11] == "ppaction://":
            # URL is indeed a macro reference, so treat it as such
            picture.click_action.hyperlink._hlink.set("action", target)
    else:
        picture.click_action.target_slide = target

    if tooltipText != "":
        picture.click_action.hyperlink._hlink.set("tooltip", tooltipText)


# If a tooltip has been set return it else return an empty string
def getPictureTooltip(picture):
    if picture.click_action.hyperlink._hlink != None:
        # There is a tooltip
        return picture.click_action.hyperlink._hlink.get("tooltip")
    else:
        # There is no tooltip
        return ""


# Create hyperlink and optional tooltip from a shape eg Chevron
def createShapeHyperlinkAndTooltip(shape, to_slide, tooltipText=""):
    shape.click_action.target_slide = to_slide
    hl = shape.click_action.hyperlink
    hl._hlink.set("tooltip", tooltipText)


def getGraphicDimensions(fname):
    """Determine the image type of fhandle and return its size.
    from draco"""
    try:
        with open(fname, "rb") as fhandle:
            head = fhandle.read(24)
            if len(head) != 24:
                return -1, -1
            if imghdr.what(fname) == "png":
                check = struct.unpack(">i", head[4:8])[0]
                if check != 0x0D0A1A0A:
                    return -1, -1
                width, height = struct.unpack(">ii", head[16:24])
            elif imghdr.what(fname) == "gif":
                width, height = struct.unpack("<HH", head[6:10])
            elif imghdr.what(fname) == "jpeg":
                try:
                    fhandle.seek(0)  # Read 0xff next
                    size = 2
                    ftype = 0
                    while not 0xC0 <= ftype <= 0xCF:
                        fhandle.seek(size, 1)
                        byte = fhandle.read(1)
                        while ord(byte) == 0xFF:
                            byte = fhandle.read(1)
                        ftype = ord(byte)
                        size = struct.unpack(">H", fhandle.read(2))[0] - 2
                    # We are at a SOFn block
                    fhandle.seek(1, 1)  # Skip 'precision' byte.
                    height, width = struct.unpack(">HH", fhandle.read(4))
                except Exception:  # IGNORE:W0703
                    return
            else:
                return -1, -1

            return width, height

    except EnvironmentError:
        return -1, -1


def getVideoDimensionsPlus(audioVideoInfo):
    if audioVideoInfo.source.find("://") > -1:
        # Video would be sourced from the web
        try:
            operUrl = urllib.request.urlopen(audioVideoInfo.source)
        except urllib.error.HTTPError as e:
            return -1, -1, "Web", None

        except socket.error as s:
            return -1, -1, "Web", None

        data = operUrl.read()

        return audioVideoInfo.width, audioVideoInfo.height, "Web", data
    else:
        # Video would be sourced from a local file
        try:
            fhandle = open(audioVideoInfo.source, "rb")
        except EnvironmentError:
            return -1, -1, "Local", None

        return audioVideoInfo.width, audioVideoInfo.height, "Local", None


# Render a list of bullets
def renderText(shape, bullets):
    baseTextDecrement = processingOptions.getCurrentOption("baseTextDecrement")
    baseTextSize = processingOptions.getCurrentOption("baseTextSize")

    tf = shape.text_frame

    for bulletNumber, bullet in enumerate(bullets):
        para0 = tf.paragraphs[0]

        if bulletNumber == 0:
            # Don't need to create paragraph
            p = para0
        else:
            # We need a new paragraph
            p = tf.add_paragraph()

        # Set the paragraph's level - zero-indexed
        p.level = int(bullet[0])

        # Set the paragraph's font size, adjusted for level, if necessary
        if baseTextSize > 0:
            p.font.size = Pt(baseTextSize - p.level * baseTextDecrement)

        addFormattedText(p, bullet[1])

        if bullet[2] == "numbered":
            makeNumberedListItem(p)

    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def findTitleShape(slide):
    if slide.shapes.title == None:
        # Have to use first shape as title
        return slide.shapes[0]
    else:
        return slide.shapes.title


def findBodyShape(slide):
    if len(slide.shapes) > 1:
        return slide.shapes[1]
    elif slide.shapes.title == None:
        return slide.shapes[0]
    else:
        return None


# Returns a top, left, width, height for content to be rendered into
def getContentRect(presentation, slide, topOfContent, margin):
    numbersHeight = processingOptions.getCurrentOption("numbersHeight")
    # Left and right are always defined by the margins
    rectLeft = margin
    rectWidth = presentation.slide_width - 2 * margin
    if topOfContent == 0:
        # There is no title on this slide
        rectTop = margin
        rectHeight = presentation.slide_height - margin - max(margin, numbersHeight)
    else:
        # There is a title on this slide
        rectTop = topOfContent + margin
        rectHeight = presentation.slide_height - rectTop - max(margin, numbersHeight)

    return (rectLeft, rectWidth, rectTop, rectHeight)


# Finds the title and adds the text to it, returning title bottom, title shape, and
# flattened title
def formatTitle(presentation, slide, titleText, titleFontSize, subtitleFontSize):
    marginBase = processingOptions.getCurrentOption("marginBase")
    pageTitleAlign = processingOptions.getCurrentOption("pagetitlealign")

    # Convert page title alignment text value to constant
    if pageTitleAlign == "left":
        titleAlignment = PP_ALIGN.LEFT
    elif pageTitleAlign == "right":
        titleAlignment = PP_ALIGN.RIGHT
    else:
        titleAlignment = PP_ALIGN.CENTER

    # Find title
    title = findTitleShape(slide)

    if titleText == "&nbsp;":
        deleteSimpleShape(title)

        return (marginBase, None, "<No title>")

    if processingOptions.getCurrentOption("adjustTitles"):
        title.top = marginBase
        title.left = marginBase
        title.width = presentation.slide_width - marginBase * 2

    # Figure out how many lines title will need (ignoring overflow)
    titleLineCount = len(titleLines := titleText.split("<br/>"))

    # This will hold the flattened title lines to be printed
    flattenedTitleLines = []

    # Add the first line of the title text to the first paragraph
    firstTitleParagraph = title.text_frame.paragraphs[0]
    flattenedTitleLines.append(addFormattedText(firstTitleParagraph, titleLines[0]))

    # Set anchor to top
    title.text_frame.vertical_anchor = MSO_ANCHOR.TOP

    # Set this paragraph's font size using pageTitleSize
    firstTitleParagraph.font.size = Pt(titleFontSize)

    # No space before the paragraph
    firstTitleParagraph.space_after = Pt(0)

    # Set first title paragraph's alignment
    firstTitleParagraph.alignment = titleAlignment

    if subtitleFontSize == "same":
        subtitleFontSize = titleFontSize

    # If there are additional title lines then add them
    for lineNumber, titleLine in enumerate(titleLines[1:]):
        # Each title line requires a new paragraph
        newPara = title.text_frame.add_paragraph()

        # Use this new paragraph, adding flattened title line to the list
        flattenedTitleLines.append(addFormattedText(newPara, titleLine))

        # No space before the paragraph
        newPara.space_before = Pt(0)

        # Set this paragraph's font size using pageSubtitleSize
        newPara.font.size = Pt(subtitleFontSize)

        # Set this paragraph's alignment
        newPara.alignment = titleAlignment

    # Note: Working off pageTitleSize and pageSubtitleSize
    if processingOptions.getCurrentOption("adjustTitles"):
        title.height = Pt(titleFontSize) + Pt(subtitleFontSize) * (titleLineCount - 1)

    # Massage title line for printing a little
    if titleLineCount > 1:
        flattenedTitleText = flattenedTitleLines[0] + " ..."
    else:
        flattenedTitleText = flattenedTitleLines[0]

    # Return where the next shape below the title would be - vertically
    return (title.top + title.height + Inches(0.1), title, flattenedTitleText)


# Parse the string after the e.g. ### for a displayable title and
# an optional heading reference
def parseTitleText(titleLineString):
    # Get rid of any cruft on the line
    slideTitleWithPossibleHref = titleLineString.strip().rstrip("#").rstrip()

    if hrefMatch := slideHrefRegex.match(slideTitleWithPossibleHref):
        # Use the explicit href
        slideTitle = hrefMatch.group(1)
        href = hrefMatch.group(2)
    else:
        # No href
        href = ""
        slideTitle = slideTitleWithPossibleHref

    return slideTitle, href


def parseText(text):
    textArray = []
    state = "N"
    fragment = ""
    lastChar = ""
    spanState = "None"

    # Replace any "\#" strings with entity reference
    text2 = text.replace("\\#", "&#x23;")

    # Replace any "<br/>" strings with newline single character
    text2 = text2.replace("<br/>", "\n")

    # Replace any escaped asterisk strings with entity reference
    text2 = text2.replace("\\*", "&lowast;")

    # Replace any asterisks with spaces either side with entity reference
    text2 = text2.replace(" * ", " &lowast; ")
    if text2[-2:] == " *":
        text2 = text2[:-2] + " &lowast;"

    # Replace any footnote reference starts with char 238
    text2 = text2.replace("[^", chr(238))

    # Replace any span style starts with char 239
    text2 = re.sub(spanStyleRegex, chr(239), text2)

    # Replace any span class starts with char 240
    text2 = re.sub(spanClassRegex, chr(240), text2)

    # Replace any span ends with char 241
    text2 = text2.replace("</span>", chr(241))

    # Replace any abbreviation starts with char 242
    text2 = text2.replace("<abbr title=", chr(242))

    # Replace any abbreviation ends with char 243
    text2 = text2.replace("</abbr>", chr(243))

    # Replace any \[ with char 244
    text2 = text2.replace("\[", chr(244))

    # Replace any \] with char 245
    text2 = text2.replace("\]", chr(245))

    # 246 is link separator special character. See below

    # Replace any {~~ with char 247
    text2 = text2.replace("{~~", chr(247))

    # Replace any ~~} with char 247
    text2 = text2.replace("~~}", chr(247))

    # Replace any {== with char 248
    text2 = text2.replace("{==", chr(248))

    # Replace any ==} with char 248
    text2 = text2.replace("==}", chr(248))

    # Replace any {>> with char 249
    text2 = text2.replace("{>>", chr(249))

    # Replace any <<} with char 249
    text2 = text2.replace("<<}", chr(249))

    # Replace any {-- with char 250
    text2 = text2.replace("{--", chr(250))

    # Replace any --} with char 250
    text2 = text2.replace("--}", chr(250))

    # Replace any {++ with char 251
    text2 = text2.replace("{++", chr(251))

    # Replace any ++} with char 251
    text2 = text2.replace("++}", chr(251))

    # Replace any <ins> with char 252
    text2 = text2.replace("<ins>", chr(252))

    # Replace any </ins> with char 252
    text2 = text2.replace("</ins>", chr(252))

    # Replace any <del> with char 253
    text2 = text2.replace("<del>", chr(253))

    # Replace any </del> with char 253
    text2 = text2.replace("</del>", chr(253))

    # Replace any <sub> with char 254
    text2 = text2.replace("<sub>", chr(254))

    # Replace any </sub> with char 254
    text2 = text2.replace("</sub>", chr(254))

    # Replace any <sup> with char 255
    text2 = text2.replace("<sup>", chr(255))

    # Replace any </sup> with char 255
    text2 = text2.replace("</sup>", chr(255))

    # Handle escaped underscore
    text2 = text2.replace("\_", "_")

    # Unescape any numeric character references
    text3 = resolveSymbols(text2)

    for c in text3:
        if c == "*":
            # Changing state
            if state == "N":
                # First * potentially starts italic
                textArray.append([state, fragment])
                fragment = ""
                state = "I"

            elif state == "I":
                # Either go to bold or end italic
                if lastChar == "*":
                    # Go to bold
                    state = "B1"

                else:
                    # End italic
                    textArray.append([state, fragment])
                    fragment = ""
                    state = "N"

            elif state == "B1":
                # Starting to close bold bracket
                state = "B2"

            elif lastChar == "*":
                # closing either bold or italic bracket
                textArray.append([state, fragment])
                fragment = ""
                state = "N"

        elif c == "`":
            if state == "N":
                # Going to code
                textArray.append([state, fragment])
                fragment = ""
                state = "C"

            else:
                # exiting code
                textArray.append([state, fragment])
                fragment = ""
                state = "N"

        elif ord(c) == 247:
            # Entering or leaving CriticMarkup replacement
            if state == "N":
                # Going to CriticMarkup replacement
                textArray.append([state, fragment])
                fragment = ""
                state = "CMRep"

            else:
                # exiting CriticMarkup replacement
                textArray.append([state, fragment])
                fragment = ""
                state = "N"

        elif ord(c) == 248:
            # Entering or leaving CriticMarkup highlight
            if state == "N":
                # Going to CriticMarkup highlight
                textArray.append([state, fragment])
                fragment = ""
                state = "CMHig"

            else:
                # exiting CriticMarkup highlight
                textArray.append([state, fragment])
                fragment = ""
                state = "N"

        elif ord(c) == 249:
            # Entering or leaving CriticMarkup comment
            if state == "N":
                # Going to CriticMarkup comment
                textArray.append([state, fragment])
                fragment = ""
                state = "CMCom"

            else:
                # exiting CriticMarkup comment
                textArray.append([state, fragment])
                fragment = ""
                state = "N"

        elif ord(c) == 250:
            # Entering or leaving CriticMarkup deletion
            if state == "N":
                # Going to CriticMarkup deletion
                textArray.append([state, fragment])
                fragment = ""
                state = "CMDel"

            else:
                # exiting CriticMarkup deletion
                textArray.append([state, fragment])
                fragment = ""
                state = "N"

        elif ord(c) == 251:
            # Entering or leaving CriticMarkup addition
            if state == "N":
                # Going to CriticMarkup addition
                textArray.append([state, fragment])
                fragment = ""
                state = "CMAdd"

            else:
                # exiting CriticMarkup addition
                textArray.append([state, fragment])
                fragment = ""
                state = "N"

        elif ord(c) == 252:
            # Entering or leaving underline
            if state == "N":
                # Going to underline
                textArray.append([state, fragment])
                fragment = ""
                state = "Ins"

            else:
                # exiting underline
                textArray.append([state, fragment])
                fragment = ""
                state = "N"

        elif ord(c) == 253:
            # Entering or leaving strikethrough
            if state == "N":
                # Going to strikethrough
                textArray.append([state, fragment])
                fragment = ""
                state = "Del"

            else:
                # exiting strikethrough
                textArray.append([state, fragment])
                fragment = ""
                state = "N"
        elif ord(c) == 254:
            # Entering or leaving subscript
            if state == "N":
                # Going to subscript
                textArray.append([state, fragment])
                fragment = ""
                state = "Sub"

            else:
                # exiting subscript
                textArray.append([state, fragment])
                fragment = ""
                state = "N"

        elif ord(c) == 255:
            # Entering or leaving superscript
            if state == "N":
                # Going to superscript
                textArray.append([state, fragment])
                fragment = ""
                state = "Sup"

            else:
                # exiting superscript
                textArray.append([state, fragment])
                fragment = ""
                state = "N"

        elif c == "[":
            if state == "N":
                # Could be entering a Link
                if fragment != "":
                    textArray.append([state, fragment])

                # The bracket is kept in in case there is no matching ]
                fragment = "["
                state = "LinkText1"

            elif state == "LinkText2":
                # Could be entering an indirect reference
                indLinkText = fragment[:-1]

                # The bracket is kept in in case there is no matching ]
                fragment = "["
                state = "LinkRef1"

        elif c == "]":
            # Could be ending picking up the link text
            if state == "LinkText1":
                # Picked up end of link text
                state = "LinkText2"

                # Remove [ and add a separator to allow for link URL
                fragment = fragment[1:] + chr(246)

            elif state == "fnref":
                # This terminates a footnote reference
                textArray.append([state, fragment])
                state = "N"
                fragment = ""

            elif state == "LinkRef1":
                # Picked up link reference
                reference = fragment[1:]

                # Attempt to look up reference
                foundReference = False
                for indref, indURL in indirectAnchors:
                    if indref == reference:
                        foundReference = True
                        break

                if foundReference:
                    # Append fragment with resolved reference
                    textArray.append(["Link", indLinkText + chr(246) + indURL])
                else:
                    print(f"Reference {reference} not resolved.\n")
                    textArray.append(["N", indLinkText])
                fragment = ""
                state = "N"

            else:
                # This was an ordinary square bracket
                fragment += "]"

        elif c == "(":
            # Could be starting to pick up the link URL
            if state == "LinkText2":
                # Picked up start of link URL
                state = "LinkURL1"
            else:
                fragment = fragment + c

        elif c == ")":
            # Could be ending picking up the link URL
            if state == "LinkURL1":
                # Picked up end of link URL
                textArray.append(["Link", fragment])
                fragment = ""
                state = "N"
            else:
                fragment = fragment + c

        elif ord(c) == 235:
            fragment = fragment + "`"

        elif ord(c) == 236:
            fragment = fragment + "<"

        elif ord(c) == 244:
            fragment = fragment + "["

        elif ord(c) == 245:
            fragment = fragment + "]"

        elif ord(c) == 242:
            if fragment != "":
                textArray.append([state, fragment])
                fragment = ""
            dictEntry = ""

        elif ord(c) == 243:
            dictEntry = fragment.split(">")
            dictAbbrev = dictEntry[1]
            dictFull = dictEntry[0].strip().strip("'").strip('"')
            abbrevDictionary[dictAbbrev] = dictFull
            textArray.append(["Gloss", dictAbbrev, dictAbbrev, dictFull])
            fragment = ""

        elif ord(c) == 241:
            # End of span
            if spanState == "Class":
                # Span with class
                splitting = fragment.split(">")
                spanText = splitting[1]
                className = splitting[0].strip().strip("'").strip('"').lower()
                styleText = ""
                if (
                    (className in bgcolors)
                    | (className in fgcolors)
                    | (className in emphases)
                ):
                    textArray.append(["SpanClass", [className, spanText]])

                    fragment = ""

                else:
                    print(
                        f"{className} is not defined. Ignoring reference to it in <span> element."
                    )

                    fragment = spanText
            else:
                # Span with style
                splitting = fragment.split(">")
                spanText = splitting[1]
                styleText = splitting[0].strip().strip("'").strip('"')
                textArray.append(["SpanStyle", [styleText, spanText]])
                className = ""
                fragment = ""

            spanState = "None"

        elif ord(c) == 240:
            # In span element where we hit the class name
            if fragment != "":
                textArray.append([state, fragment])

                fragment = ""
            spanState = "Class"

        elif ord(c) == 239:
            # In span element where we hit the style text
            if fragment != "":
                textArray.append([state, fragment])

                fragment = ""
                spanState = "Style"
        elif ord(c) == 238:
            if fragment != "":
                textArray.append([state, fragment])

                fragment = ""
                state = "fnref"
        else:
            fragment = fragment + c

        lastChar = c

    if fragment != "":
        textArray.append([state, fragment])
    return textArray


# Calls the tokeniser and then handles the fragments it gets back
def addFormattedText(p, text):
    boldBold = processingOptions.getCurrentOption("boldBold")
    boldColour = processingOptions.getCurrentOption("boldColour")
    italicItalic = processingOptions.getCurrentOption("italicItalic")
    italicColour = processingOptions.getCurrentOption("italicColour")
    monoFont = processingOptions.getCurrentOption("monoFont")

    # Get back parsed text fragments, along with control information on each
    # fragment
    parsedText = parseText(text)

    # Replace chr(237) with > in each Fragment
    for f in range(len(parsedText)):
        if parsedText[f][0] in ["SpanClass", "SpanStyle"]:
            parsedText[f][1][1] = parsedText[f][1][1].replace(chr(237), ">")
        else:
            parsedText[f][-1] = parsedText[f][-1].replace(chr(237), ">")

    # Prime flattened Text
    flattenedText = ""
    for fragment in parsedText:
        if fragment[0] == "Gloss":
            fragType, fragDetail, fragTerm, fragTitle = fragment
        else:
            fragType, fragDetail = fragment

        # Break into subfragments around a newline
        if fragType == "SpanClass":
            className, fragText = fragDetail
            styleText = ""
            subfragments = fragText.split("\n")
        elif fragType == "SpanStyle":
            styleText, fragText = fragDetail
            className = ""
            subfragments = fragText.split("\n")
        else:
            subfragments = fragDetail.split("\n")

        # Process each subfragment
        sfnum = 0
        for subfragment in subfragments:
            if sfnum > 0:
                # Subfragments after the first need to be preceded by a line break
                p.add_line_break()

            sfnum += 1
            # Ensure "\*" is rendered as a literal asterisk
            subfragment = subfragment.replace("&lowast;", "*")

            # Ensure "\#" is rendered as a literal octothorpe
            subfragment = subfragment.replace("&#x23;", "#")

            run = p.add_run()

            if fragType not in ["Link", "fnref", "Gloss"]:
                run.text = subfragment
            elif fragType == "Gloss":
                run.text = fragTerm

            if fragType == "I":
                font = run.font

                if italicItalic == True:
                    font.italic = True

                if italicColour != ("None", ""):
                    setColour(font.color, italicColour)

            elif fragType == "Gloss":
                # Add this run to abbrevRunsDictionary - for Glossary fix ups later
                if fragTerm not in abbrevRunsDictionary:
                    abbrevRunsDictionary[fragTerm] = []
                abbrevRunsDictionary[fragTerm].append(run)
            elif fragType == "fnref":
                font = run.font
                font.size = Pt(16)
                set_superscript(font)
                fnref = fragment[1]
                if fnref in footnoteReferences:
                    footnoteNumber = footnoteReferences.index(fnref)
                    run.text = str(footnoteNumber + 1)
                    footnoteRunsDictionary[footnoteNumber] = run
                else:
                    run.text = "[?]"
                    print("Error: Footnote reference '" + fnref + "' unresolved.")
                linkText = "!"
                fragment = ""

            elif fragType == "SpanClass":
                handleSpanClass(run, className)

            elif fragType == "SpanStyle":
                handleSpanStyle(run, styleText)

            elif fragType == "B2":
                font = run.font

                if boldBold == True:
                    font.bold = True

                if boldColour != ("None", ""):
                    setColour(font.color, boldColour)

            elif fragType == "C":
                font = run.font
                font.name = monoFont
            elif fragType == "CMRep":
                font = run.font
                font.color.rgb = RGBColor(255, 140, 0)
                run.text = "{~~" + subfragment + "~~}"
            elif fragType == "CMHig":
                font = run.font
                font.color.rgb = RGBColor(195, 0, 195)
                run.text = "{==" + subfragment + "==}"
            elif fragType == "CMCom":
                font = run.font
                font.color.rgb = RGBColor(0, 0, 195)
                run.text = "{>>" + subfragment + "<<}"
            elif fragType == "CMDel":
                font = run.font
                font.color.rgb = RGBColor(195, 0, 0)
                run.text = "{--" + subfragment + "--}"
            elif fragType == "CMAdd":
                font = run.font
                font.color.rgb = RGBColor(0, 195, 0)
                run.text = "{++" + subfragment + "++}"
            elif fragType == "Ins":
                font = run.font
                font.underline = True
            elif fragType == "Del":
                font = run.font
                set_strikethrough(font)
            elif fragType == "Sub":
                font = run.font
                set_subscript(font)
            elif fragType == "Sup":
                font = run.font
                set_superscript(font)
            elif fragType == "Link":
                linkArray = subfragment.split(chr(246))
                linkText = linkArray[0]
                linkURL = linkArray[1]
                run.text = linkText
                if linkURL.startswith("#"):
                    # Is an internal Url
                    linkHref = linkURL[1:].strip()
                    href_runs[linkHref] = run
                else:
                    # Not an internal link so create it
                    hlink = run.hyperlink
                    hlink.address = linkURL

                    # URL might be a macro reference
                    if linkURL[:11] == "ppaction://":
                        # URL is indeed a macro reference, so treat it as such
                        hlink._hlinkClick.action = linkURL

            # Add the flattened text from this subfragment
            if fragType == "Link":
                flattenedText = flattenedText + linkText
            else:
                flattenedText = flattenedText + subfragment

    return flattenedText


def addFooter(presentation, slideNumber, slide):
    numbersHeight = processingOptions.getCurrentOption("numbersHeight")
    
    numbersFontSizeSpec = processingOptions.getCurrentOption("numbersFontSize")
    if numbersFontSizeSpec == "":
        numbersFontSize = Pt(12)
    else:
        numbersFontSize = Pt(numbersFontSizeSpec)
    
    shapes = slide.shapes
    footer = shapes.add_textbox(
        Inches(0.1),
        presentation.slide_height - numbersHeight,
        Inches(0.2),
        numbersHeight / 2,
    )
    frame = footer.text_frame
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = str(slideNumber)
    font = run.font
    font.size = numbersFontSize


# Called "Simple" because more complex shapes might not work
def deleteSimpleShape(shape):
    if shape == None:
        return
    shapeElement = shape.element
    shapeElement.getparent().remove(shapeElement)


def createProcessingSummarySlide(presentation, rawMetadata):
    tableMargin = processingOptions.getCurrentOption("tableMargin")
    pageTitleSize = processingOptions.getCurrentOption("pageTitleSize")
    pageSubtitleSize = processingOptions.getCurrentOption("pageSubtitleSize")
    # Use the first slide in the template presentation as the base
    slide = presentation.slides[0]

    # Delete any body shape - other than action buttons
    bodyShape = findBodyShape(slide)
    if bodyShape.name.startswith("Action Button:") is False:
        deleteSimpleShape(bodyShape)

    # Build "run time" text
    now = datetime.datetime.now()
    runTime = now.strftime("%H:%M").lstrip()
    runDate = now.strftime("%e %B, %G").lstrip()
    runDateTime = "Presentation built: " + runTime + " on " + runDate

    # Format title and add title text
    slideTitleBottom, title, flattenedTitle = formatTitle(
        presentation,
        slide,
        banner + "<br/>" + runDateTime,
        pageTitleSize,
        pageSubtitleSize,
    )

    # Work out how many pairs of columns we need
    if processingOptions.hideMetadataStyle:
        # Adjust metadata item count to remove style.
        metadata = []
        for metadataItem in rawMetadata:
            if metadataItem[0].startswith("style.") == False:
                metadata.append(metadataItem)
    else:
        metadata = rawMetadata

    metadataRows = len(metadata)

    maxMetadataRowsPerColumn = 15
    if metadataRows > 4 * maxMetadataRowsPerColumn:
        metadataColumnPairs = 5
    elif metadataRows > 3 * maxMetadataRowsPerColumn:
        metadataColumnPairs = 4
    elif metadataRows > 2 * maxMetadataRowsPerColumn:
        metadataColumnPairs = 3
    elif metadataRows > maxMetadataRowsPerColumn:
        metadataColumnPairs = 2
    else:
        metadataColumnPairs = 1

    columns = metadataColumnPairs * 2
    rows = min(maxMetadataRowsPerColumn, metadataRows)

    # Get the rectangle the content will draw in
    contentLeft, contentWidth, contentTop, contentHeight = getContentRect(
        presentation, slide, slideTitleBottom, tableMargin
    )

    tableHeight = min(contentHeight, Inches(0.25) * rows)

    # Figure out the width of a single-width column
    columnWidthUnit = int(contentWidth / (2 * metadataColumnPairs))

    # Create the table with the above number of rows and columns
    newTable = slide.shapes.add_table(
        rows, columns, tableMargin, contentTop, contentWidth, tableHeight
    ).table

    # Don't want headings
    newTable.first_row = False

    cols = newTable.columns
    for cp in range(metadataColumnPairs):
        cols[2 * cp].width = columnWidthUnit
        cols[2 * cp + 1].width = columnWidthUnit

    row = 0
    column = 0
    for item in metadata:
        key, value = item

        if row == maxMetadataRowsPerColumn:
            # Move to next column
            column += 2
            row = 0

        # Set text of metadata key cell
        newTable.cell(row, column).text = key
        if processingOptions.dynamicallyChangedOptions.get(key) is not None:
            # Set text of metadata value cell - with asterisk
            newTable.cell(row, column + 1).text = value + "*"

            # Colour key cell blue
            p1 = newTable.cell(row, column).text_frame.paragraphs[0]
            p1.font.color.rgb = RGBColor.from_string("0000FF")

            # Colour value cell blue
            p2 = newTable.cell(row, column + 1).text_frame.paragraphs[0]
            p2.font.color.rgb = RGBColor.from_string("0000FF")
        else:
            # Set text of metadata value cell - without asterisk
            newTable.cell(row, column + 1).text = value

        if metadataColumnPairs == 5:
            newTable.cell(row, column).text_frame.paragraphs[0].font.size = Pt(8)
            newTable.cell(row, column + 1).text_frame.paragraphs[0].font.size = Pt(8)
        elif metadataColumnPairs == 4:
            newTable.cell(row, column).text_frame.paragraphs[0].font.size = Pt(10)
            newTable.cell(row, column + 1).text_frame.paragraphs[0].font.size = Pt(10)
        elif metadataColumnPairs == 3:
            newTable.cell(row, column).text_frame.paragraphs[0].font.size = Pt(12)
            newTable.cell(row, column + 1).text_frame.paragraphs[0].font.size = Pt(12)
        elif metadataColumnPairs == 2:
            newTable.cell(row, column).text_frame.paragraphs[0].font.size = Pt(14)
            newTable.cell(row, column + 1).text_frame.paragraphs[0].font.size = Pt(14)
        else:
            newTable.cell(row, column).text_frame.paragraphs[0].font.size = Pt(16)
            newTable.cell(row, column + 1).text_frame.paragraphs[0].font.size = Pt(16)
        row += 1


# Note: This doesn't use formatTitle()
def createTitleOrSectionSlide(
    presentation,
    slideNumber,
    titleText,
    layout,
    titleSize,
    subtitleText,
    subtitleSize,
    notes_text,
):
    marginBase = processingOptions.getCurrentOption("marginBase")

    slide = addSlide(presentation, presentation.slide_layouts[layout], None)

    # Add title
    title = findTitleShape(slide)
    flattenedTitle = addFormattedText(title.text_frame.paragraphs[0], titleText)

    title.text_frame.paragraphs[0].font.size = Pt(titleSize)

    # Add subtitle - if there is one
    if (subtitleText.strip() != "") & (subtitleText[0:2] != "\n\n"):
        # There is a subtitle
        chunks = subtitleText.strip().split("\n\n")
        subtitleText = chunks[0]
        notes_text = "\n\n".join(chunks[1:])
        createSlideNotes(slide, notes_text)

        subtitleShape = findBodyShape(slide)
        if subtitleShape != None:
            addFormattedText(subtitleShape.text_frame.paragraphs[0], subtitleText)
            subtitleShape.text_frame.paragraphs[0].font.size = Pt(subtitleSize)
        else:
            print("Warning: No subtitle shape on this slide to add text to.")
    else:
        # Reformat subtitle shape to be out of the way
        subtitleShape = findBodyShape(slide)
        if subtitleShape != None:
            subtitleShape.top = title.top + title.height + marginBase * 2
            subtitleShape.width = title.width
            subtitleShape.left = title.left
            subtitleShape.height = marginBase * 2

    reportSlideTitle(slideNumber, 1, flattenedTitle)

    if want_numbers_headings is True:
        addFooter(presentation, slideNumber, slide)

    return slide


def handleWhateverGraphicType(GraphicFilename):
    # Handles both physical file and URI file types
    GraphicFilename = GraphicFilename.strip()

    # First ensure we have the data in a file and know if the source was a URI
    if ":" in GraphicFilename:
        # Is a URI - so we have to retrieve it and store it in a temporary file
        is_uri = True

        # Massage the URI into a printable filename
        if len(GraphicFilename) > 50:
            printableGraphicFilename = (
                GraphicFilename[:25] + "..." + GraphicFilename[-25:]
            )
        else:
            printableGraphicFilename = GraphicFilename

        # Get Temporary File Directory - which might be None
        tempDir = processingOptions.getCurrentOption("tempDir")

        # Retrieve the data into a temporary file
        try:
            operUrl = urllib.request.urlopen(GraphicFilename)

        except urllib.error.HTTPError as e:
            print("HTTP error: " + str(e.code))
            return GraphicFilename, printableGraphicFilename

        except socket.error as s:
            print("Socket error. (Web site not found).")
            return GraphicFilename, printableGraphicFilename

        data = operUrl.read()

        # Get Content-Type header - if set
        content_type = ""

        if str(type(operUrl)) == "<class 'http.client.HTTPResponse'>":
            # Can try to get content header
            content_type = operUrl.getheader("content-type")

            if content_type == "":
                # Obtain file extension by searching in the URL
                extensionPos = GraphicFilename.rindex(".")
                lastSlashPos = GraphicFilename.rindex("/")
                if lastSlashPos > extensionPos:
                    fileExt = ""

                else:
                    fileExt = GraphicFilename[extensionPos:]

            else:
                # Set file extension based on Content-Type header_name
                # Note: It's been translated to lower case above
                if content_type in ["image/jpeg", "image/jpg"]:
                    # Note: only the first of these two is legitimate
                    fileExt = "jpg"
                elif content_type == "image/png":
                    fileExt = "png"
                elif content_type in ["image/svg+xml", "image/svg"]:
                    # Note: only the first of these two is legitimate
                    fileExt = "svg"
                elif content_type == "application/postscript":
                    fileExt = "eps"
                else:
                    fileExt = None
        else:
            fileExt = None

        # Store in a temporary file
        try:
            tempGraphicFile = tempfile.NamedTemporaryFile(
                delete=False, suffix=fileExt, dir=tempDir
            )
        except IOError as e:
            print("Couldn't create temporary file. md2pptx terminating")
            exit()

        tempGraphicFile.write(data)
        convertibleFilename = tempGraphicFile.name
        tempGraphicFile.close()

    else:
        is_uri = False

        # Files don't get their names edited
        printableGraphicFilename = GraphicFilename
        convertibleFilename = GraphicFilename

    if is_uri:
        lastSlash = GraphicFilename.rfind("/")
        lastDot = GraphicFilename.rfind(".")

        PNGname = GraphicFilename[lastSlash + 1 : lastDot] + ".PNG"
    else:
        PNGname = GraphicFilename[:-4] + ".PNG"

    # Process the file - whatever the origin - based on file extension
    if ".svg" in GraphicFilename.lower():
        # is an SVG file
        if have_cairosvg:
            # Convert SVG file to temporary PNG
            # Store in a temporary file

            # Get Temporary File Directory - which might be None
            tempDir = processingOptions.getCurrentOption("tempDir")

            try:
                graphicFile = tempfile.NamedTemporaryFile(
                    delete=False, suffix=".PNG", dir=tempDir
                )
            except IOError as e:
                print("Couldn't create temporary file. md2pptx terminating")
                exit()

            cairosvg.svg2png(file_obj=open(convertibleFilename), write_to=graphicFile)

            # Retrieve the temporary file name
            GraphicFilename = graphicFile.name

            if processingOptions.getCurrentOption("exportGraphics"):
                try:
                    shutil.copy(GraphicFilename, PNGname)
                except:
                    print("Copy error: " + PNGname)

        else:
            print("Don't have CairoSVG installed. Terminating.")
            # sys.exit()
    elif ".eps" in GraphicFilename.lower():
        if have_pillow:
            # Get EPS file
            im = PIL.Image.open(GraphicFilename)

            # Get Temporary File Directory - which might be None
            tempDir = processingOptions.getCurrentOption("tempDir")

            # Store in a temporary file
            try:
                graphicFile = tempfile.NamedTemporaryFile(
                    delete=False, suffix=".PNG", dir=tempDir
                )
            except IOError as e:
                print("Couldn't create temporary file. md2pptx terminating")
                exit()

            try:
                im.save(graphicFile)
            except:
                print("Could not convert EPS file. Is Ghostview installed?\n")
                print("Terminating.\n")
                # sys.exit()

            # Retrieve the temporary file name
            GraphicFilename = graphicFile.name
            if processingOptions.getCurrentOption("exportGraphics"):
                try:
                    shutil.copy(GraphicFilename, PNGname)
                except:
                    print("Copy error: " + PNGname)

    else:
        GraphicFilename = convertibleFilename

    return GraphicFilename, printableGraphicFilename


def handleGraphViz(slide, renderingRectangle, codeLines, codeType):
    # Condition GraphViz source
    s = graphviz.Source("\n".join(codeLines), format="png")

    # Invent a temporary filename for the rendered graphic
    dotFile = "md2pptx-temporary-dotfile.png"

    # Render the .dot source as a graphic
    s.render(cleanup=True, outfile=dotFile)

    # Figure out the dimensions of the rendered graphic
    dotGraphicWidth, dotGraphicHeight = getGraphicDimensions(dotFile)

    # Adjust those dimensions with the usual scaling rules
    (dotPicWidth, dotPicHeight, scaledByHeight) = scalePicture(
        renderingRectangle.width,
        renderingRectangle.height,
        dotGraphicWidth,
        dotGraphicHeight,
    )

    # Add the picture to the current slide
    slide.shapes.add_picture(
        dotFile,
        renderingRectangle.left + (renderingRectangle.width - dotPicWidth) / 2,
        renderingRectangle.top + (renderingRectangle.height - dotPicHeight) / 2,
        dotPicWidth,
        dotPicHeight,
    )

    # Delete the temporary graphic file
    os.remove(dotFile)


def handleFunnel(slide, renderingRectangle, codeLines, codeType):
    funnelColours = processingOptions.getCurrentOption("funnelColours")
    funnelBorderColour = processingOptions.getCurrentOption("funnelBorderColour")
    funnelTitleColour = processingOptions.getCurrentOption("funnelTitleColour")
    funnelTextColour = processingOptions.getCurrentOption("funnelTextColour")
    funnelLabelsPercent = processingOptions.getCurrentOption("funnelLabelsPercent")
    funnelLabelPosition = processingOptions.getCurrentOption("funnelLabelPosition")
    funnelWidest = processingOptions.getCurrentOption("funnelWidest")

    f = funnel.Funnel()

    f.makeFunnel(
        slide,
        renderingRectangle,
        codeLines,
        funnelColours,
        codeType,
        funnelBorderColour,
        funnelTitleColour,
        funnelTextColour,
        funnelLabelsPercent,
        funnelLabelPosition,
        funnelWidest,
    )


def createCodeBlock(slideInfo, slide, renderingRectangle, codeBlockNumber):
    monoFont = processingOptions.getCurrentOption("monoFont")
    baseTextSize = processingOptions.getCurrentOption("baseTextSize")
    defaultBaseTextSize = processingOptions.getDefaultOption("baseTextSize")

    # A variable number of newlines appear before the actual code
    codeLines = slideInfo.code[codeBlockNumber]

    # Figure out code slide type
    if codeLines[0].startswith("<pre"):
        codeType = "pre"
    elif codeLines[0].startswith("<code"):
        codeType = "code"
    elif codeLines[0].startswith("```"):
        if codeLines[0] == "```":
            # Plain backticks
            codeType = "backticks"
        else:
            # Some other type - such as GraphViz .dot
            codeType = codeLines[0][3:].strip().lower()

            if codeType.startswith("dot"):
                # GraphViz dot
                if have_graphviz is False:
                    # Don't have GraphViz so warn and treat as basic backticks
                    sys.stderr.write(
                        "GraphViz not installed. Rendering as backticks.\n"
                    )

                    codeType = "backticks"
                else:
                    # Have GraphViz installed so use it ti add a picture
                    handleGraphViz(slide, renderingRectangle, codeLines[1:-1], codeType)

                    return slide

            elif codeType.startswith("funnel"):
                # Built-in Funnel Diagram
                handleFunnel(slide, renderingRectangle, codeLines[1:-1], codeType)

                return slide

            else:
                # Some other type with backticks
                codeType = "backticks"

    else:
        codeType = "indented"

    # Handle any trailing empty lines
    while codeLines[-1] == "":
        codeLines.pop(-1)

    # Handle any leading <pre>, <code>, triple backtick line
    if startswithOneOf(codeLines[0], ["<pre>", "<code>", "```"]):
        codeLines.pop(0)

    # Handle any trailing </pre>, </code>, triple backtick line
    if startswithOneOf(codeLines[-1], ["</pre>", "</code>", "```"]):
        codeLines.pop(-1)

    codeBox = slide.shapes.add_textbox(
        renderingRectangle.left,
        renderingRectangle.top,
        renderingRectangle.width,
        renderingRectangle.height,
    )

    # Try to control text frame but SHAPE_TO_FIT_TEXT doesn't seem to work
    tf = codeBox.text_frame
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    tf.word_wrap = False

    # Fill the code box with background colour - whether explicit or defaulted
    fill = codeBox.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(
        processingOptions.getCurrentOption("codeBackground")
    )

    # Get the sole paragraph
    p = codeBox.text_frame.paragraphs[0]

    # Set the font size slightly smaller than usual
    if len(codeLines) >= 20:
        divisor = 1.5
    else:
        divisor = 1.2
    if baseTextSize > 0:
        p.font.size = int(Pt(baseTextSize) / divisor)
    else:
        p.font.size = int(Pt(defaultBaseTextSize) / divisor)

    # Estimate how wide the code box would need to be at the current font size
    # versus actual width
    codeColumns = processingOptions.getCurrentOption("codeColumns")
    fixedPitchHeightWidthRatio = processingOptions.getCurrentOption(
        "fixedPitchHeightWidthRatio"
    )

    estimatedWidthVersusCodeboxWidth = (
        p.font.size * codeColumns / codeBox.width / fixedPitchHeightWidthRatio
    )
    if estimatedWidthVersusCodeboxWidth > 1:
        # The code is wider so shrink the font so the code fits
        p.font.size = p.font.size / estimatedWidthVersusCodeboxWidth
    else:
        # The code is narrower so shrink the code textbox so the code just fits
        # - assuming declared width is accurate
        codeBox.width = int(p.font.size * codeColumns / fixedPitchHeightWidthRatio)

        # Center the code box - actually don't - 5 October 2021 temporary "fix"
        # codeBox.left = int((presentation.slide_width - codeBox.width) / 2)

    # Use the code foreground colour - whether explicit or defaulted
    p.font.color.rgb = RGBColor.from_string(
        processingOptions.getCurrentOption("codeforeground")
    )

    # Adjust code box height based on lines
    codeBox.height = min(
        len(codeLines) * Pt(baseTextSize + 5), renderingRectangle.height
    )

    # Add code
    if codeType == "pre":
        # span elements want special handling
        for codeLine in codeLines:
            # Resolve eg entity references
            codeLine = resolveSymbols(codeLine)

            # Split the line - and maybe there are spans
            spanFragments = codeLine.split("<span ")
            if len(spanFragments) > 1:
                textArray = []
                # Break line down into what will become runs
                for fragmentNumber, fragment in enumerate(spanFragments):
                    if fragmentNumber > 0:
                        # Find start of span class
                        fragment = "<span " + fragment
                        if spanClassMatch := spanClassRegex.match(fragment):
                            afterSpanTag = fragment[spanClassMatch.end(1) :]
                            className = afterSpanTag[7 : afterSpanTag.index(">") - 1]
                            if (
                                (className in bgcolors)
                                | (className in fgcolors)
                                | (className in emphases)
                            ):
                                afterClosingAngle = afterSpanTag[
                                    afterSpanTag.index(">") + 1 :
                                ]
                                startEnd = afterClosingAngle.index("</span>")
                                afterSpan2 = afterClosingAngle[:startEnd]
                                afterSpan3 = afterClosingAngle[startEnd + 7 :]
                                textArray.append(["SpanClass", [className, afterSpan2]])
                                textArray.append(["Normal", afterSpan3])
                                fragment = ""
                            else:
                                print(
                                    className
                                    + " is not defined. Ignoring reference to it in <span> element."
                                )
                        elif spanStyleMatch := spanStyleRegex.match(fragment):
                            afterSpanTag = fragment[spanStyleMatch.end(1) :]
                            styleText = afterSpanTag[7 : afterSpanTag.index(">") - 1]
                            styleElements = styleText.split(";")
                            afterClosingAngle = afterSpanTag[
                                afterSpanTag.index(">") + 1 :
                            ]
                            startEnd = afterClosingAngle.index("</span>")
                            afterSpan2 = afterClosingAngle[:startEnd]
                            afterSpan3 = afterClosingAngle[startEnd + 7 :]
                            textArray.append(["SpanStyle", [styleText, afterSpan2]])
                            textArray.append(["Normal", afterSpan3])
                            fragment = ""
                    else:
                        textArray.append(["Normal", fragment])

                # Now we have a text array we can add the runs for the line
                for textArrayItem in textArray:
                    textArrayItemType = textArrayItem[0]
                    if textArrayItemType == "Normal":
                        # Is not in a span element bracket
                        className = ""
                        spanStyle = ""
                        spanText = textArrayItem[1]

                    elif textArrayItemType == "SpanClass":
                        # Is in a span class element bracket
                        className = textArrayItem[1][0]
                        spanText = textArrayItem[1][1]
                        spanStyle = ""

                    else:
                        # Is in a span style element bracket
                        spanStyle = textArrayItem[1][0]
                        spanText = textArrayItem[1][1]

                    if spanText != "":
                        run = p.add_run()
                        run.text = spanText
                        font = run.font
                        font.name = monoFont

                    if className != "":
                        # Augment run with whatever the span class calls for
                        handleSpanClass(run, className)

                    if spanStyle != "":
                        # Augment the run with whatever the style calls for
                        handleSpanStyle(run, spanStyle)

                # Add terminating newline
                run = p.add_run()
                run.text = "\n"
                font = run.font
                font.name = monoFont
            else:
                # Line has no spans in
                run = p.add_run()
                run.text = codeLine + "\n"
                font = run.font
                font.name = monoFont

    else:
        # span doesn't need treating specially
        for codeLine in codeLines:
            # Resolve eg entity references
            codeLine = resolveSymbols(codeLine)

            run = p.add_run()
            run.text = codeLine + "\n"
            font = run.font
            font.name = monoFont

    return slide


def createAbstractSlide(presentation, slideNumber, titleText, paragraphs):
    titleOnlyLayout = processingOptions.getCurrentOption("titleOnlyLayout")
    marginBase = processingOptions.getCurrentOption("marginBase")
    pageTitleSize = processingOptions.getCurrentOption("pageTitleSize")
    pageSubtitleSize = processingOptions.getCurrentOption("pageSubtitleSize")

    slide = addSlide(presentation, presentation.slide_layouts[titleOnlyLayout], None)

    shapes = slide.shapes

    # Add title and constrain its size and placement
    slideTitleBottom, title, flattenedTitle = formatTitle(
        presentation, slide, titleText, pageTitleSize, pageSubtitleSize
    )

    reportSlideTitle(slideNumber, 3, "Abstract: " + flattenedTitle)

    # Get the rectangle the content will draw in
    contentLeft, contentWidth, contentTop, contentHeight = getContentRect(
        presentation, slide, slideTitleBottom, marginBase
    )

    # Add abstract text
    abstractBox = slide.shapes.add_textbox(
        contentLeft,
        contentTop,
        contentWidth,
        contentHeight,
    )

    p = abstractBox.text_frame.paragraphs[0]
    tf = abstractBox.text_frame
    f = p.font
    f.size = Pt(22)
    for para, abstractParagraph in enumerate(paragraphs):
        paragraphLevel, paragraphText, paragraphType = abstractParagraph

        if para > 0:
            # Spacer paragraph
            p = tf.add_paragraph()
            f = p.font
            f.size = Pt(22)

            # Content paragraph
            p = tf.add_paragraph()
            f = p.font
            f.size = Pt(22)
        addFormattedText(p, paragraphText)

    tf.word_wrap = True

    if want_numbers_content is True:
        addFooter(presentation, slideNumber, slide)

    return slide


# Unified creation of a table or a code or a content slide
def createContentSlide(presentation, slideNumber, slideInfo):
    titleOnlyLayout = processingOptions.getCurrentOption("titleOnlyLayout")
    contentSlideLayout = processingOptions.getCurrentOption("contentSlideLayout")
    marginBase = processingOptions.getCurrentOption("marginBase")
    pageTitleSize = processingOptions.getCurrentOption("pageTitleSize")
    pageSubtitleSize = processingOptions.getCurrentOption("pageSubtitleSize")

    # slideInfo's body text is only filled in if there is code - and that's
    # where the code - plus preamble and postamble is.
    if slideInfo.code != "":
        haveCode = True
    else:
        haveCode = False

    # Create the slide and check for bullets and/or cards
    if (slideInfo.bullets == []) & (slideInfo.cards == []):
        # No bullets or cards so "title only"
        slideLayout = titleOnlyLayout
        haveBulletsCards = False
    else:
        # Either bullets or cards or both so not "title only"
        slideLayout = contentSlideLayout
        haveBulletsCards = True

    slide = addSlide(presentation, presentation.slide_layouts[slideLayout], slideInfo)

    # Check for table / graphics content
    if slideInfo.tableRows == []:
        haveTableGraphics = False
    else:
        haveTableGraphics = True

    ####################################################################
    # At this point haveCode, haveBulletsCards, haveTableGraphics have #
    # been appropriately set                                           #
    ####################################################################

    # Add slide title
    titleText = slideInfo.titleText

    slideTitleBottom, title, flattenedTitle = formatTitle(
        presentation, slide, titleText, pageTitleSize, pageSubtitleSize
    )

    # Log slide's title
    reportSlideTitle(slideNumber, 3, flattenedTitle)

    ####################################################################
    # Get the dimensions of the content area to place all content in   #
    ####################################################################
    contentLeft, contentWidth, contentTop, contentHeight = getContentRect(
        presentation, slide, slideTitleBottom, marginBase
    )

    ####################################################################
    # Check whether there are too many elements in the sequence to     #
    # render - and warn if there are. Then calculate how many to render#
    ####################################################################
    if len(slideInfo.sequence) > maxBlocks:
        print(f"Too many blocks to render. Only {str(maxBlocks)} will be rendered.")
    blocksToRender = min(maxBlocks, len(slideInfo.sequence))

    ####################################################################
    # Get the dimensions of the rectangles we'll place the graphics in #
    # and their top left corner coordinates
    ####################################################################
    allContentSplit = 0
    contentSplit = processingOptions.getCurrentOption("contentSplit")
    for b in range(blocksToRender):
        allContentSplit = allContentSplit + contentSplit[b]

    verticalCursor = contentTop
    horizontalCursor = contentLeft

    codeBlockNumber = 0
    tableBlockNumber = 0

    for b in range(blocksToRender):
        if processingOptions.getCurrentOption("contentSplitDirection") == "vertical":
            # Height and top
            blockHeight = int(contentHeight * contentSplit[b] / allContentSplit)
            blockTop = verticalCursor
            verticalCursor = verticalCursor + blockHeight

            # Width and left
            blockWidth = contentWidth
            blockLeft = contentLeft
        else:
            # Height and top
            blockHeight = contentHeight
            blockTop = contentTop

            # Width and left
            blockWidth = int(contentWidth * contentSplit[b] / allContentSplit)
            blockLeft = horizontalCursor
            horizontalCursor = horizontalCursor + blockWidth

        renderingRectangle = Rectangle(blockTop, blockLeft, blockHeight, blockWidth)

        if slideInfo.sequence[b] == "table":
            createTableBlock(slideInfo, slide, renderingRectangle, tableBlockNumber)
            tableBlockNumber += 1

        elif slideInfo.sequence[b] == "list":
            createListBlock(slideInfo, slide, renderingRectangle)
        else:
            createCodeBlock(slideInfo, slide, renderingRectangle, codeBlockNumber)
            codeBlockNumber += 1

    if want_numbers_content is True:
        addFooter(presentation, slideNumber, slide)

    return slide


def createListBlock(slideInfo, slide, renderingRectangle):
    horizontalCardGap = processingOptions.getCurrentOption("horizontalcardgap")
    verticalCardGap = processingOptions.getCurrentOption("verticalcardgap")
    cardTitleAlign = processingOptions.getCurrentOption("cardtitlealign")
    cardTitlePosition = processingOptions.getCurrentOption("cardtitleposition")
    cardShape = processingOptions.getCurrentOption("cardshape")
    cardLayout = processingOptions.getCurrentOption("cardlayout")
    cardPercent = processingOptions.getCurrentOption("cardpercent")
    cardShadow = processingOptions.getCurrentOption("cardshadow")
    cardTitleSize = processingOptions.getCurrentOption("cardtitlesize")
    cardBorderWidth = processingOptions.getCurrentOption("cardborderwidth")
    cardBorderColour = processingOptions.getCurrentOption("cardbordercolour")
    cardTitleColour = processingOptions.getCurrentOption("cardtitlecolour")
    cardTitleBackgrounds = processingOptions.getCurrentOption("cardtitlebackground")
    cardColours = processingOptions.getCurrentOption("cardcolour")
    cardDividerColour = processingOptions.getCurrentOption("carddividercolour")
    cardGraphicHeight = processingOptions.getCurrentOption("cardgraphicheight")
    cardGraphicPosition = processingOptions.getCurrentOption("cardgraphicposition")
    marginBase = processingOptions.getCurrentOption("marginBase")
    pageTitleSize = processingOptions.getCurrentOption("pageTitleSize")
    pageSubtitleSize = processingOptions.getCurrentOption("pageSubtitleSize")

    # Get bulleted text shape - either for bullets above cards or first card's body shape
    bulletsShape = findBodyShape(slide)

    # Set bulleted shape top, left, width
    bulletsShape.top = renderingRectangle.top
    bulletsShape.left = renderingRectangle.left
    bulletsShape.width = renderingRectangle.width

    bulletCount = len(slideInfo.bullets)

    # Set bulleted text height - depending on whether there's a card
    # Remainder is card area height - if there are cards
    if slideInfo.cards == []:
        # There are no cards so the bullets shape takes the whole content area
        bulletsShape.height = renderingRectangle.height

        # There are no cards so the card area is zero height
        cardAreaHeight = 0
        cardCount = 0
    else:
        # There are cards
        if bulletCount > 0:
            # Bullets shape vertically shortened
            bulletsShape.height = int(
                renderingRectangle.height * (100 - cardPercent) / 100
            )

            # Card area takes the rest of the content area
            cardAreaHeight = int(renderingRectangle.height) - bulletsShape.height
        else:
            # No bullets so content is all cards
            bulletsShape.height = 0

            cardAreaHeight = renderingRectangle.height

        cardBackgroundShapes = []
        cardTitleShapes = []
        cardGraphicShapes = []
        cardGraphicDimensions = []
        cardBodyShapes = []

        cardLeft = []
        cardTop = []
        cardBodyTop = []
        cardBackgroundTop = []
        cardTitleTop = []

        cardCount = len(slideInfo.cards)

        # card width applies to card title, card graphic, card background, card body
        if cardLayout == "horizontal":
            # Divide horizontal card space up
            cardWidth = int(
                (renderingRectangle.width - Inches(horizontalCardGap) * (cardCount - 1))
                / cardCount
            )
        else:
            # Card takes all the horizontal space
            cardWidth = int(renderingRectangle.width)

        # Calculate title top and height - horizontal layout
        if cardTitleSize > 0:
            # Specified by user. "72" because in points
            cardTitleHeightRaw = Inches(cardTitleSize / 72)
        else:
            # Shrunk to 2/3 of page title height.  "72" because in points
            cardTitleHeightRaw = Inches(int(10000 * pageTitleSize * 2 / 3 / 72) / 10000)

        # Adjust title height to be slightly larger than the text
        cardTitleHeight = cardTitleHeightRaw + Inches(0.1)

        cardGraphicHeightRaw = int(Inches(cardGraphicHeight))

        if bulletCount > 0:
            # Bullets so cards and their titles occupy less than whole height
            cardAreaTop = bulletsShape.height + renderingRectangle.top

        else:
            # No bullets so cards and their titles occupy whole height
            cardAreaTop = renderingRectangle.top

        if cardLayout == "horizontal":
            # Card takes up all the card area, vertically
            cardHeight = cardAreaHeight
        else:
            # Card height is just a proportion of the card area height

            if cardTitlePosition == "above":
                paddingFactor = Inches(verticalCardGap - 0.05)
            else:
                paddingFactor = Inches(verticalCardGap)

            cardHeight = int((cardAreaHeight) / cardCount - paddingFactor)

        # Store slide title shape for cloning
        slideTitleShape = findTitleShape(slide)

        for c in range(cardCount):
            # Calculate a card's vertical position
            if cardLayout == "horizontal":
                cardTop.append(cardAreaTop)
            else:
                cardTop.append(int((cardHeight + paddingFactor) * c + cardAreaTop))

            # Card top and height
            if cardTitlePosition == "above":
                # Card title above card background
                cardBackgroundTop.append(cardTop[c] + cardTitleHeight)
                cardBodyTop.append(cardBackgroundTop[c])
            else:
                # card title inside card background
                cardBackgroundTop.append(cardTop[c])

                # Leave room above the card body for the card title
                if cards[c][0] == "&nbsp;":
                    cardBodyTop.append(cardBackgroundTop[c])
                else:
                    cardBodyTop.append(cardBackgroundTop[c] + cardTitleHeight)

            # Calculate a card's horizontal position
            if cardLayout == "horizontal":
                cardLeft.append(
                    marginBase + c * (cardWidth + Inches(horizontalCardGap))
                )
            else:
                cardLeft.append(marginBase)

            # Card title modeled on slide title - but smaller
            cardTitleShape = addClonedShape(slide, slideTitleShape)

            cardTitleShapes.append(cardTitleShape)

            cardGraphicFilename = cards[c][1]
            if cardGraphicFilename != "":
                cardGraphicDimensions.append(getGraphicDimensions(
                     cardGraphicFilename
                 ))

                # Create card graphic shape - to be resized later
                cardGraphicShape = slide.shapes.add_picture(
                    cardGraphicFilename,
                    Inches(0),
                    Inches(0),
                )
            else:
                cardGraphicShape = None
                cardGraphicDimensions.append(None)

            cardGraphicShapes.append(cardGraphicShape)

            # Clear text from cloned title and add in the title text
            cardTitleShape.text_frame.paragraphs[0].text = ""
            addFormattedText(cardTitleShape.text_frame.paragraphs[0], cards[c][0])

            # Set card title font size
            if cardTitleSize > 0:
                cardTitleShape.text_frame.paragraphs[0].font.size = Pt(cardTitleSize)
            else:
                cardTitleShape.text_frame.paragraphs[0].font.size = Pt(
                    pageTitleSize * 2 / 3
                )

            # Titles are aligned one of three ways
            if cardTitleAlign == "l":
                cardTitleShape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            elif cardTitleAlign == "c":
                cardTitleShape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            else:
                cardTitleShape.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            # Fill in card's background - if necessary
            if (cardTitleBackgrounds[0] != ("None", "")) & (
                cardTitlePosition != "inside"
            ):
                cardTitleBackground = cardTitleBackgrounds[
                    c % len(cardTitleBackgrounds)
                ]

                fill = cardTitleShape.fill
                fill.solid()

                setColour(fill.fore_color, cardTitleBackground)

            # Create card background and make sure it's behind the card body (and maybe card title)
            if cardShape == "rounded":
                # Rounded Rectangle for card
                cardBackgroundShape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(0),
                    Inches(0),
                    Inches(0),
                    Inches(0),
                )

                # Rounding adjustment works better with different values for horizontal and vertical cards
                if cardLayout == "horizontal":
                    # Make the rounding radius small. This is 1/4 the default
                    cardBackgroundShape.adjustments[0] = 0.0416675
                else:
                    # Make the rounding radius smallish. This is 1/2 the default
                    cardBackgroundShape.adjustments[0] = 0.083335
            else:
                # Squared-corner Rectangle for card
                cardBackgroundShape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0), Inches(0)
                )

            cardBackgroundShapes.append(cardBackgroundShape)

            if cardShape == "line":
                cardBackgroundShape.fill.background()

            sendToBack(slide.shapes, cardBackgroundShape)

            # Card shape modeled on bulleted list
            if (bulletCount > 0) | (c > 0):
                cardBodyShape = addClonedShape(slide, bulletsShape)
            else:
                # Co-opt bullets shape as first card shape
                cardBodyShape = bulletsShape
            cardBodyShapes.append(cardBodyShape)

            # Make card's body transparent
            fill = cardBodyShape.fill
            fill.background()

            # Fill in card's background - if necessary
            if (cardColours[0] != ("None", "")) & (cardShape != "line"):
                cardColour = cardColours[c % len(cardColours)]

                fill = cardBackgroundShape.fill
                fill.solid()

                setColour(fill.fore_color, cardColour)

    # Adjust bullets shape height - and calculate verticals for any cards

    bulletsShape.bottom = bulletsShape.top + bulletsShape.height

    # Fill in the main bullets shape
    renderText(bulletsShape, slideInfo.bullets)

    for c in range(cardCount):
        # Get the shapes for this card
        cardBackgroundShape = cardBackgroundShapes[c]
        cardTitleShape = cardTitleShapes[c]
        cardBodyShape = cardBodyShapes[c]
        cardGraphicShape = cardGraphicShapes[c]

        # Get dimensions of any graphic
        if cardGraphicShape is not None:
            cardGraphicNativeWidth, cardGraphicNativeHeight = cardGraphicDimensions[c]
        else:
            cardGraphicNativeWidth, cardGraphicNativeHeight = (0, 0)

        # Set the card shapes' width
        cardBackgroundShape.width = cardWidth
        cardBodyShape.width = cardWidth
        cardTitleShape.width = cardWidth

        # Set the card shapes' left side
        cardBackgroundShape.left = cardLeft[c]
        cardTitleShape.left = cardLeft[c]
        cardBodyShape.left = cardLeft[c]

        # Position card title
        cardTitleShape.top = cardTop[c]
        cardTitleShape.height = cardTitleHeight

        # Colour the title - if cardTitleColour specified
        if cardTitleColour != ("None", ""):
            setColour(
                cardTitleShape.text_frame.paragraphs[0].font.color, cardTitleColour
            )

        # Calculate positions and heights within card background of body
        if cardTitlePosition == "above":
            cardBackgroundShape.top = cardTop[c] + cardTitleHeight
            cardBodyHeight = cardHeight - cardTitleHeight
            cardBackgroundShape.height = cardBodyHeight
        else:
            cardBackgroundShape.top = cardTop[c]
            cardBodyHeight = cardHeight
            cardBackgroundShape.height = cardBodyHeight

        if (c > 0) & (cardShape == "line"):
            if cardLayout == "horizontal":
                dividingLine = createLine(
                    cardBackgroundShape.left - int(Inches(horizontalCardGap / 2)),
                    cardBackgroundShape.top + Inches(0.75),
                    cardBackgroundShape.left - int(Inches(horizontalCardGap / 2)),
                    cardBackgroundShape.top + cardBackgroundShape.height - Inches(0.75),
                    slide.shapes,
                )
            else:
                dividingLine = createLine(
                    cardBackgroundShape.left + Inches(0.75),
                    cardBackgroundShape.top - int(Inches(verticalCardGap / 2)),
                    cardBackgroundShape.left
                    + cardBackgroundShape.width
                    - int(Inches(0.75)),
                    cardBackgroundShape.top - int(Inches(verticalCardGap / 2)),
                    slide.shapes,
                )

            if cardDividerColour != ("None", ""):
                setColour(dividingLine.line.color, cardDividerColour)

            dividingLine.line.width = Pt(2.0)

        # Position card
        if cardGraphicShape == None:
            # No graphic on this card
            cardBodyShape.top = cardBodyTop[c]
            cardBodyShape.height = cardBodyHeight
        else:
            # Make room for graphic
            if cardGraphicPosition == "before":
                # Leave room above card body shape for graphic
                cardBodyShape.top = cardBodyTop[c] + cardGraphicHeightRaw
            else:
                # Leave room below card body shape for graphic
                cardBodyShape.top = cardBodyTop[c]
            
            cardBodyShape.height = cardBodyHeight - cardGraphicHeightRaw
            
            # Scale graphic
            (cardPicWidth, cardPicHeight, scaledByHeight) = scalePicture(
                cardTitleShape.width,
                cardGraphicHeightRaw,
                cardGraphicNativeWidth,
                cardGraphicNativeHeight,
            )
            if cards[c][0] == "&nbsp;":
                if cardGraphicPosition == "before":
                    cardGraphicShape.top = cardTitleShape.top + int((cardGraphicHeightRaw - cardPicHeight) / 2)
                else:
                    cardGraphicShape.top = cardBodyHeight + int((cardGraphicHeightRaw - cardPicHeight) / 2)
            else:
                if cardGraphicPosition == "before":
                    cardGraphicShape.top = cardTitleShape.top + cardTitleShape.height + int((cardGraphicHeightRaw - cardPicHeight) / 2)
                else:
                    cardGraphicShape.top = cardBodyHeight + cardTitleShape.height + int((cardGraphicHeightRaw - cardPicHeight) / 2)

            cardGraphicShape.left = cardTitleShape.left +int((cardTitleShape.width - cardPicWidth) / 2)
            cardGraphicShape.height = int(cardPicHeight)
            cardGraphicShape.width = int(cardPicWidth)

        if len(cards[c]) > 1:
            renderText(cardBodyShape, slideInfo.cards[c][2])

        lf = cardBackgroundShape.line

        if (cardBorderColour != ("None", "")) & (cardShape != "line"):
            setColour(lf.color, cardBorderColour)

        if cardShape == "line":
            lf.fill.background()
        elif cardBorderWidth > 0:
            lf.width = Pt(cardBorderWidth)
        elif cardBorderWidth == 0:
            lf.fill.background()

        if cardShadow:
            createShadow(cardBackgroundShape)

    return slide


def createTableBlock(slideInfo, slide, renderingRectangle, tableBlockNumber):
    tableRows = slideInfo.tableRows[tableBlockNumber]
    tableMargin = processingOptions.getCurrentOption("tableMargin")
    marginBase = processingOptions.getCurrentOption("marginBase")
    baseTextSize = processingOptions.getCurrentOption("baseTextSize")
    tempDir = processingOptions.getCurrentOption("tempDir")

    printableTopLeftGraphicFilename = ""
    printableTopRightGraphicFilename = ""
    printablebottomLeftGraphicFilename = ""
    printableBottomRightGraphicFilename = ""

    # Handle table body
    if (len(tableRows) <= 2) & (len(tableRows[0]) <= 2):
        # This is a table with 1 or 2 rows and 1 or 2 columns
        isGraphicsGrid = True
        gridRows = len(tableRows)
        if gridRows == 1:
            gridColumns = len(tableRows[0])
        else:
            gridColumns = max(len(tableRows[0]), len(tableRows[1]))

        topGraphicCount = 0

        topLeftCellString = tableRows[0][0]
        # Attempt to retrieve media information for left side - top row
        (
            topLeftGraphicTitle,
            topLeftGraphicFilename,
            printableTopLeftGraphicFilename,
            topLeftGraphicHref,
            topLeftHTML,
            topLeftVideo,
            topGraphicCount,
        ) = parseMedia(topLeftCellString, topGraphicCount)

        # Attempt to retrieve filename for right side - top row
        if len(tableRows[0]) == 2:
            topRightCellString = tableRows[0][1]
        else:
            topRightCellString = ""

        (
            topRightGraphicTitle,
            topRightGraphicFilename,
            printableTopRightGraphicFilename,
            topRightGraphicHref,
            topRightHTML,
            topRightVideo,
            topGraphicCount,
        ) = parseMedia(topRightCellString, topGraphicCount)

        if topGraphicCount == 0:
            # Revert to normal table processing as no graphic spec in at least one cell
            isGraphicsGrid = False

        if gridRows == 2:
            # Attempt to retrieve filename for left side - bottom row
            bottomGraphicCount = 0

            bottomLeftCellString = tableRows[1][0]

            # Attempt to retrieve media information for left side - bottom row
            (
                bottomLeftGraphicTitle,
                bottomLeftGraphicFilename,
                printableBottomLeftGraphicFilename,
                bottomLeftGraphicHref,
                bottomLeftHTML,
                bottomLeftVideo,
                bottomGraphicCount,
            ) = parseMedia(bottomLeftCellString, bottomGraphicCount)

            # Attempt to retrieve filename for right side - bottom row
            if gridColumns == 2:
                if len(tableRows[1]) == 1:
                    # There is one cell in bottom row so this is centred "3-up"
                    bottomRightCellString = ""
                else:
                    bottomRightCellString = tableRows[1][1]
            else:
                bottomRightCellString = ""

            (
                bottomRightGraphicTitle,
                bottomRightGraphicFilename,
                printableBottomRightGraphicFilename,
                bottomRightGraphicHref,
                bottomRightHTML,
                bottomRightVideo,
                bottomGraphicCount,
            ) = parseMedia(bottomRightCellString, bottomGraphicCount)

            if bottomGraphicCount == 0:
                # Revert to normal table processing as no graphic spec in at least one cell
                isGraphicsGrid = False

    else:
        # This is a normal table because it has too many rows or columns to be a graphics grid
        isGraphicsGrid = False

    if isGraphicsGrid == True:
        ####################################################################
        # Print the media filenames                                        #
        ####################################################################
        if gridColumns == 2:
            # Doing 1- or 2-row side-by-side graphics slide
            reportGraphicFilenames(
                printableTopLeftGraphicFilename, printableTopRightGraphicFilename
            )
        else:
            # Doing 2 row, single column graphics slide
            reportGraphicFilenames(printableTopLeftGraphicFilename)

        if gridRows == 2:
            # Second row of filenames
            if gridColumns == 2:
                reportGraphicFilenames(
                    printableBottomLeftGraphicFilename,
                    printableBottomRightGraphicFilename,
                )
            else:
                reportGraphicFilenames(printableBottomLeftGraphicFilename)

        ####################################################################
        # Get the media dimensions                                         #
        ####################################################################

        if topLeftGraphicFilename != "":
            topLeftMediaWidth, topLeftMediaHeight = getGraphicDimensions(
                topLeftGraphicFilename
            )
            if topLeftMediaWidth == -1:
                if gridRows == 2:
                    print(
                        "Missing top left image file: "
                        + printableTopLeftGraphicFilename
                    )
                else:
                    print("Missing left image file: " + printableTopLeftGraphicFilename)

                return slide

        elif topLeftVideo is not None:
            (
                topLeftMediaWidth,
                topLeftMediaHeight,
                topLeftVideoType,
                topLeftVideoData,
            ) = getVideoDimensionsPlus(topLeftVideo)

            if topLeftMediaWidth == -1:
                if gridRows == 2:
                    print(
                        "Missing top left video file: "
                        + printableTopLeftGraphicFilename
                    )
                else:
                    print("Missing left video file: " + printableTopLeftGraphicFilename)

                return slide

        if gridColumns == 2:
            # Get top right image dimensions
            if topRightGraphicFilename != "":
                topRightMediaWidth, topRightMediaHeight = getGraphicDimensions(
                    topRightGraphicFilename
                )
                if topRightMediaWidth == -1:
                    if gridRows == 2:
                        print(
                            "Missing top right image file: "
                            + printableTopRightGraphicFilename
                        )
                    else:
                        print(
                            "Missing right image file: "
                            + printableTopRightGraphicFilename
                        )

                    return slide

            elif topRightVideo is not None:
                (
                    topRightMediaWidth,
                    topRightMediaHeight,
                    topRightVideoType,
                    topRightVideoData,
                ) = getVideoDimensionsPlus(topRightVideo)

                if topRightMediaWidth == -1:
                    if gridRows == 2:
                        print(
                            "Missing top right video file: "
                            + printableTopRightGraphicFilename
                        )
                    else:
                        print(
                            "Missing right video file: "
                            + printableTopRightGraphicFilename
                        )

                    return slide

        if gridRows == 2:
            # Get bottom left image dimensions
            if bottomLeftGraphicFilename != "":
                bottomLeftMediaWidth, bottomLeftMediaHeight = getGraphicDimensions(
                    bottomLeftGraphicFilename
                )
                if bottomLeftMediaWidth == -1:
                    print(
                        "Missing bottom left image file: "
                        + printableBottomLeftGraphicFilename
                    )
                    return slide

            elif bottomLeftVideo is not None:
                (
                    bottomLeftMediaWidth,
                    bottomLeftMediaHeight,
                    bottomLeftVideoType,
                    bottomLeftVideoData,
                ) = getVideoDimensionsPlus(bottomLeftVideo)

                if bottomLeftMediaWidth == -1:
                    if gridRows == 2:
                        print(
                            "Missing bottom left video file: "
                            + printableBottomLeftGraphicFilename
                        )
                    else:
                        print(
                            "Missing left video file: "
                            + printableBottomLeftGraphicFilename
                        )

                    return slide

            if gridColumns == 2:
                # Get bottom right image dimensions
                if bottomRightGraphicFilename != "":
                    (
                        bottomRightMediaWidth,
                        bottomRightMediaHeight,
                    ) = getGraphicDimensions(bottomRightGraphicFilename)

                    if bottomRightMediaWidth == -1:
                        print(
                            "Missing bottom right image file: "
                            + printableBottomRightGraphicFilename
                        )

                        return slide

                elif bottomRightVideo is not None:
                    (
                        bottomRightMediaWidth,
                        bottomRightMediaHeight,
                        bottomRightVideoType,
                        bottomRightVideoData,
                    ) = getVideoDimensionsPlus(bottomRightVideo)

                    if bottomRightMediaWidth == -1:
                        if gridRows == 2:
                            print(
                                "Missing bottom right video file: "
                                + printableBottomRightGraphicFilename
                            )
                        else:
                            print(
                                "Missing right video file: "
                                + printableBottomRightGraphicFilename
                            )

                        return slide

        # Calculate maximum picture height on slide
        maxPicHeight = renderingRectangle.height

        if gridRows == 2:
            # Adjusted if two rows
            maxPicHeight = maxPicHeight / 2 + Inches(0.2)

        # Calculate maximum picture width on slide
        maxPicWidth = renderingRectangle.width
        if gridColumns == 2:
            # Adjusted if two columns
            maxPicWidth = maxPicWidth / 2 - marginBase

        # Calculate horizontal middle of graphics space
        midGraphicsSpaceX = renderingRectangle.left + renderingRectangle.width / 2

        ####################################################################
        # Calculate the size of each graphic - scaled by the above rect    #
        ####################################################################

        if (topLeftGraphicFilename != "") | (topLeftVideo is not None):
            (
                topLeftPicWidth,
                topLeftPicHeight,
                usingHeightToScale,
            ) = scalePicture(
                maxPicWidth, maxPicHeight, topLeftMediaWidth, topLeftMediaHeight
            )

            if usingHeightToScale:
                # Calculate horizontal start
                if (gridColumns == 2) and (
                    (topRightGraphicFilename != "") | (topRightVideo is not None)
                ):
                    # Align top left media item to the left
                    topLeftPicX = (
                        renderingRectangle.left
                        + (midGraphicsSpaceX - marginBase - topLeftPicWidth) / 2
                    )
                else:
                    # Center sole top media item
                    topLeftPicX = midGraphicsSpaceX - topLeftPicWidth / 2
            else:
                # Calculate horizontal start
                if (gridColumns == 2) and (
                    (topRightGraphicFilename != "") | (topRightVideo is not None)
                ):
                    # Align top left media item to the left
                    topLeftPicX = renderingRectangle.left
                else:
                    # Center sole top media item
                    topLeftPicX = midGraphicsSpaceX - topLeftPicWidth / 2

            # Calculate vertical start
            topLeftPicY = renderingRectangle.top + (maxPicHeight - topLeftPicHeight) / 2

            if gridRows == 2:
                topLeftPicY -= Inches(0.2)

        if topLeftGraphicFilename != "":
            topLeftPicture = slide.shapes.add_picture(
                topLeftGraphicFilename,
                topLeftPicX,
                topLeftPicY,
                topLeftPicWidth,
                topLeftPicHeight,
            )

            if topLeftGraphicHref == "":
                topLeftGraphicHref = "#XYZZY-None"

            pictureInfos.append(
                (topLeftPicture, topLeftGraphicHref, topLeftGraphicTitle)
            )
        elif topLeftVideo is not None:
            if topLeftVideoType == "Local":
                # Can use local file directly
                topLeftVideo = slide.shapes.add_movie(
                    topLeftVideo.source,
                    topLeftPicX,
                    topLeftPicY,
                    topLeftPicWidth,
                    topLeftPicHeight,
                    topLeftVideo.poster,
                )
            else:
                # First copy video data to temporary file
                tempVideoFile = tempfile.NamedTemporaryFile(
                    delete=False, suffix="mp4", dir=tempDir
                )
                tempVideoFile.write(topLeftVideoData)
                convertibleFilename = tempVideoFile.name
                tempVideoFile.close()

                # Use temporary file to make video
                topLeftVideo = slide.shapes.add_movie(
                    convertibleFilename,
                    topLeftPicX,
                    topLeftPicY,
                    topLeftPicWidth,
                    topLeftPicHeight,
                    topLeftVideo.poster,
                )

        if gridColumns == 2:
            # Top right media item
            if (topRightGraphicFilename != "") | (topRightVideo is not None):
                (
                    topRightPicWidth,
                    topRightPicHeight,
                    usingHeightToScale,
                ) = scalePicture(
                    maxPicWidth, maxPicHeight, topRightMediaWidth, topRightMediaHeight
                )

                if usingHeightToScale:
                    # Calculate horizontal start
                    topRightPicX = (
                        renderingRectangle.width + midGraphicsSpaceX - topRightPicWidth
                    ) / 2
                else:
                    # Calculate horizontal start
                    topRightPicX = (
                        renderingRectangle.width + midGraphicsSpaceX - topRightPicWidth
                    ) / 2

                # Calculate vertical start
                topRightPicY = (
                    renderingRectangle.top + (maxPicHeight - topRightPicHeight) / 2
                )

                if gridRows == 2:
                    topRightPicY -= Inches(0.2)

            if topRightGraphicFilename != "":
                topRightPicture = slide.shapes.add_picture(
                    topRightGraphicFilename,
                    topRightPicX,
                    topRightPicY,
                    topRightPicWidth,
                    topRightPicHeight,
                )

                if topRightGraphicHref == "":
                    topRightGraphicHref = "#XYZZY-None"

                pictureInfos.append(
                    (topRightPicture, topRightGraphicHref, topRightGraphicTitle)
                )

            elif topRightVideo is not None:
                if topRightVideoType == "Local":
                    # Can use local file directly
                    topRightVideo = slide.shapes.add_movie(
                        topRightVideo.source,
                        topRightPicX,
                        topRightPicY,
                        topRightPicWidth,
                        topRightPicHeight,
                        topRightVideo.poster,
                    )
                else:
                    # First copy video data to temporary file
                    tempVideoFile = tempfile.NamedTemporaryFile(
                        delete=False, suffix="mp4", dir=tempDir
                    )
                    tempVideoFile.write(topRightVideoData)
                    convertibleFilename = tempVideoFile.name
                    tempVideoFile.close()

                    # Use temporary file to make video
                    topRightVideo = slide.shapes.add_movie(
                        convertibleFilename,
                        topRightPicX,
                        topRightPicY,
                        topRightPicWidth,
                        topRightPicHeight,
                        topRightVideo.poster,
                    )

        if gridRows == 2:
            # Need second row of media items
            # Bottom left media item
            if (bottomLeftGraphicFilename != "") | (bottomLeftVideo is not None):
                (
                    bottomLeftPicWidth,
                    bottomLeftPicHeight,
                    usingHeightToScale,
                ) = scalePicture(
                    maxPicWidth,
                    maxPicHeight,
                    bottomLeftMediaWidth,
                    bottomLeftMediaHeight,
                )

                if usingHeightToScale:
                    # Calculate horizontal start
                    if (gridColumns == 2) & (
                        (bottomRightGraphicFilename != "")
                        | (bottomRightVideo is not None)
                    ):
                        bottomLeftPicX = (
                            marginBase
                            + (midGraphicsSpaceX - marginBase - bottomLeftPicWidth) / 2
                        )
                    else:
                        bottomLeftPicX = midGraphicsSpaceX - bottomLeftPicWidth / 2
                else:
                    # Calculate horizontal start
                    if (gridColumns == 2) and (bottomRightGraphicFilename != ""):
                        # Align bottom left picture to the left
                        bottomLeftPicX = marginBase
                    else:
                        # Center sole bottom media item
                        bottomLeftPicX = midGraphicsSpaceX - bottomLeftPicWidth / 2

                # Calculate vertical start
                bottomLeftPicY = (
                    renderingRectangle.top + (maxPicHeight + bottomLeftPicHeight) / 2
                )

                if gridRows == 2:
                    bottomLeftPicY -= Inches(0.2)

            if bottomLeftGraphicFilename != "":
                bottomLeftPicture = slide.shapes.add_picture(
                    bottomLeftGraphicFilename,
                    bottomLeftPicX,
                    bottomLeftPicY,
                    bottomLeftPicWidth,
                    bottomLeftPicHeight,
                )

                if bottomLeftGraphicHref == "":
                    bottomLeftGraphicHref = "#XYZZY-None"

                pictureInfos.append(
                    (bottomLeftPicture, bottomLeftGraphicHref, bottomLeftGraphicTitle)
                )

            elif bottomLeftVideo is not None:
                if bottomLeftVideoType == "Local":
                    # Can use local file directly
                    bottomLeftVideo = slide.shapes.add_movie(
                        bottomLeftVideo.source,
                        bottomLeftPicX,
                        bottomLeftPicY,
                        bottomLeftPicWidth,
                        bottomLeftPicHeight,
                        bottomLeftVideo.poster,
                    )
                else:
                    # First copy video data to temporary file
                    tempVideoFile = tempfile.NamedTemporaryFile(
                        delete=False, suffix="mp4", dir=tempDir
                    )
                    tempVideoFile.write(bottomLeftVideoData)
                    convertibleFilename = tempVideoFile.name
                    tempVideoFile.close()

                    # Use temporary file to make video
                    bottomLeftVideo = slide.shapes.add_movie(
                        convertibleFilename,
                        bottomLeftPicX,
                        bottomLeftPicY,
                        bottomLeftPicWidth,
                        bottomLeftPicHeight,
                        bottomLeftVideo.poster,
                    )

            if gridColumns == 2:
                # Bottom right media item
                if (bottomRightGraphicFilename != "") | (bottomRightVideo is not None):
                    (
                        bottomRightPicWidth,
                        bottomRightPicHeight,
                        usingHeightToScale,
                    ) = scalePicture(
                        maxPicWidth,
                        maxPicHeight,
                        bottomRightMediaWidth,
                        bottomRightMediaHeight,
                    )

                    if usingHeightToScale:
                        # Calculate horizontal start
                        bottomRightPicX = (
                            renderingRectangle.width
                            + midGraphicsSpaceX
                            - bottomRightPicWidth
                        ) / 2

                    else:
                        # Use the width to scale
                        # Calculate horizontal start
                        bottomRightPicX = (
                            renderingRectangle.width
                            + midGraphicsSpaceX
                            - bottomRightPicWidth
                        ) / 2

                    # Calculate vertical start
                    bottomRightPicY = (
                        renderingRectangle.top
                        + (maxPicHeight + bottomRightPicHeight) / 2
                    )

                    if gridRows == 2:
                        bottomRightPicY -= Inches(0.2)

                if bottomRightGraphicFilename != "":
                    if bottomRightGraphicFilename != "":
                        bottomRightPicture = slide.shapes.add_picture(
                            bottomRightGraphicFilename,
                            bottomRightPicX,
                            bottomRightPicY,
                            bottomRightPicWidth,
                            bottomRightPicHeight,
                        )

                        if bottomRightGraphicHref == "":
                            bottomRightGraphicHref = "#XYZZY-None"

                        pictureInfos.append(
                            (
                                bottomRightPicture,
                                bottomRightGraphicHref,
                                bottomRightGraphicTitle,
                            )
                        )
                elif bottomRightVideo is not None:
                    if bottomRightVideoType == "Local":
                        # Can use local file directly
                        bottomRightVideo = slide.shapes.add_movie(
                            bottomRightVideo.source,
                            bottomRightPicX,
                            bottomRightPicY,
                            bottomRightPicWidth,
                            bottomRightPicHeight,
                            bottomRightVideo.poster,
                        )
                    else:
                        # First copy video data to temporary file
                        tempVideoFile = tempfile.NamedTemporaryFile(
                            delete=False, suffix="mp4", dir=tempDir
                        )
                        tempVideoFile.write(bottomRightVideoData)
                        convertibleFilename = tempVideoFile.name
                        tempVideoFile.close()

                        # Use temporary file to make video
                        bottomRightVideo = slide.shapes.add_movie(
                            convertibleFilename,
                            bottomRightPicX,
                            bottomRightPicY,
                            bottomRightPicWidth,
                            bottomRightPicHeight,
                            bottomRightVideo.poster,
                        )

    else:
        ################
        #              #
        # Normal table #
        #              #
        ################

        # Calculate maximum number of columns - as this is how wide we'll make the table
        columns = 0
        for row in tableRows:
            columns = max(columns, len(row))

        alignments = []
        widths = []

        # Adjust table if it contains a dash line as it's second line
        if len(tableRows) > 1:
            firstCellSecondRow = tableRows[1][0]
            if (firstCellSecondRow.startswith("-")) | (
                firstCellSecondRow.startswith(":-")
            ):
                haveTableHeading = True
            else:
                haveTableHeading = False
        else:
            haveTableHeading = False

        if haveTableHeading is True:
            # Has table heading
            tableHeadingBlurb = " with heading"

            # Figure out alignments of cells
            for cell in tableRows[1]:
                if cell.startswith(":-"):
                    if cell.endswith("-:"):
                        alignments.append("c")
                    else:
                        alignments.append("l")
                elif cell.endswith("-:"):
                    alignments.append("r")
                else:
                    alignments.append("l")

                widths.append(cell.count("-"))

            # Default any missing columns to left / single width
            if len(tableRows[1]) < columns:
                for _ in range(columns - len(tableRows[1])):
                    alignments.append("l")
                    widths.append(1)

            widths_total = sum(widths)

            # Remove this alignment / widths row from the table
            del tableRows[1]
        else:
            # No table heading
            tableHeadingBlurb = " without heading"

            # Use default width - 1 - and default alignment - l
            for c in range(columns):
                widths.append(1)
                alignments.append("l")

            # We don't know the widths so treat all equal
            widths_total = columns

        # Calculate number of rows
        rows = len(tableRows)
        alignments_count = len(alignments)

        # Create the table with the above number of rows and columns
        newTableShape = slide.shapes.add_table(rows, columns, 0, 0, 0, 0)
        newTable = newTableShape.table

        newTableShape.top = renderingRectangle.top
        newTableShape.left = renderingRectangle.left + tableMargin - marginBase
        newTableShape.height = min(renderingRectangle.height, Inches(0.25) * rows)
        newTableShape.width = renderingRectangle.width - 2 * (tableMargin - marginBase)
        shapeWidth = newTableShape.width

        # Set whether first row is not special
        newTable.first_row = haveTableHeading

        print(
            "           --> "
            + str(rows)
            + " x "
            + str(columns)
            + " table"
            + tableHeadingBlurb
        )

        # Set column widths
        cols = newTable.columns
        for colno in range(columns):
            cols[colno].width = int(shapeWidth * widths[colno] / widths_total)

        # Get options for filling in the cells
        compactTables = processingOptions.getCurrentOption("compactTables")
        spanCells = processingOptions.getCurrentOption("spanCells")
        tableHeadingSize = processingOptions.getCurrentOption("tableHeadingSize")

        # Fill in the cells
        for rowNumber, row in enumerate(tableRows):
            # Add dummy cells to the end of the row so that there are as many
            # cells in the row as there are columns in the table
            cellCount = len(row)

            # Unless there is a non-empty cell there is no anchor cell for this row
            if spanCells == "yes":
                potentialAnchorCell = None

            for c in range(cellCount, columns):
                row.append("")

            for columnNumber, cell in enumerate(row):
                newCell = newTable.cell(rowNumber, columnNumber)

                if spanCells == "yes":
                    if cell != "":
                        potentialAnchorCell = newCell
                    else:
                        if potentialAnchorCell is not None:
                            # Might need to remove previous cell merge
                            if potentialAnchorCell.span_width > 1:
                                potentialAnchorCell.split()

                            # Merge the cells from the anchor up to this one
                            potentialAnchorCell.merge(newCell)

                # For compact table remove the margins around the text
                if compactTables > 0:
                    newCell.margin_top = Pt(0)
                    newCell.margin_bottom = Pt(0)

                newCell.text = ""
                text_frame = newCell.text_frame

                # Set cell's text alignment
                p = text_frame.paragraphs[0]

                # Set cell's text size - if necessary
                if baseTextSize > 0:
                    p.font.size = Pt(baseTextSize)

                # For compact table use specified point size for text
                if compactTables > 0:
                    p.font.size = Pt(compactTables)

                if (rowNumber == 0) & (tableHeadingSize > 0):
                    p.font.size = Pt(tableHeadingSize)

                if columnNumber >= alignments_count:
                    p.alignment = PP_ALIGN.LEFT
                elif alignments[columnNumber] == "r":
                    p.alignment = PP_ALIGN.RIGHT
                elif alignments[columnNumber] == "c":
                    p.alignment = PP_ALIGN.CENTER
                else:
                    p.alignment = PP_ALIGN.LEFT

                addFormattedText(p, cell)

        # Apply table border styling - whether there is any or not
        applyTableLineStyling(
            newTable,
            processingOptions,
        )

    return slide


def createChevron(
    text,
    x,
    y,
    width,
    height,
    filled,
    shapes,
    fontSize,
    wantLink,
    unhighlightedBackground,
):
    global TOCruns

    # Create shape
    shape = shapes.add_shape(MSO_SHAPE.CHEVRON, x, y, width, height)

    # Set shape's text
    shape.text = text

    # Set shape's text attributes
    tf = shape.text_frame
    p = tf.paragraphs[0]
    f = p.font
    f.size = Pt(fontSize)
    f.color.rgb = RGBColor(0, 0, 0)

    # If want link create it from the first run
    if wantLink:
        TOCruns.append(p.runs[0])

    # Set shape's outline attributes
    shape.line.color.rgb = RGBColor(0, 0, 0)
    shape.line.width = Pt(1.0)

    # Potentially fill background
    if filled is False:
        shape.fill.background()
    else:
        if wantLink & (unhighlightedBackground != ""):
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(unhighlightedBackground)


def createOval(
    text,
    x,
    y,
    width,
    height,
    filled,
    shapes,
    fontSize,
    wantLink,
    unhighlightedBackground,
):
    global TOCruns

    # Create shape
    shape = shapes.add_shape(MSO_SHAPE.OVAL, x, y, width, height)

    # Set shape's text
    shape.text = text

    # Set shape's text attributes
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    f = p.font
    f.size = Pt(fontSize)
    f.color.rgb = RGBColor(0, 0, 0)

    # If want link create it from the first run
    if wantLink:
        TOCruns.append(p.runs[0])

    # Set shape's outline attributes
    shape.line.color.rgb = RGBColor(191, 191, 191)
    shape.line.width = Pt(1.0)

    # Potentially fill background
    if filled is False:
        shape.fill.background()
        shape.line.width = Pt(3.0)
    else:
        if wantLink & (unhighlightedBackground != ""):
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(unhighlightedBackground)


def createLine(x0, y0, x1, y1, shapes, colour=("RGB", "#BFBFBF"), width=4.0):
    # Create line
    line = shapes.add_shape(MSO_SHAPE.LINE_INVERSE, x0, y0, x1 - x0, y1 - y0)

    # Set shape's outline attributes
    setColour(line.line.color, colour)

    line.line.width = Pt(width)

    return line


def delinkify(text):
    if linkMatch := linkRegex.match(text):
        linkText = linkMatch.group(1)
        linkURL = linkMatch.group(2)
        return (linkText, linkURL)

    elif linkMatch := indirectReferenceUsageRegex(text):
        print(linkMatch.group(1))
        print(linkMatch.group(2))
        return (text, "")

    else:
        return (text, "")


def createTOCSlide(presentation, slideNumber, titleText, bullets, tocStyle):
    global SectionSlides
    titleOnlyLayout = processingOptions.getCurrentOption("titleOnlyLayout")
    blankLayout = processingOptions.getCurrentOption("blankLayout")
    tocTitle = processingOptions.getCurrentOption("tocTitle")
    marginBase = processingOptions.getCurrentOption("marginBase")
    pageTitleSize = processingOptions.getCurrentOption("pageTitleSize")
    pageSubtitleSize = processingOptions.getCurrentOption("pageSubtitleSize")

    if tocStyle != "plain":
        if titleText == tocTitle:
            reportSlideTitle(
                slideNumber, 3, f'Table Of Contents (Style: "{tocStyle}") {titleText}'
            )

        else:
            reportSlideTitle(slideNumber, 2, titleText)

    if tocStyle == "plain":
        if titleText != tocTitle:
            slide = createTitleOrSectionSlide(
                presentation,
                slideNumber,
                titleText,
                processingOptions.getCurrentOption("sectionSlideLayout"),
                processingOptions.getCurrentOption("sectionTitleSize"),
                slideInfo.subtitleText,
                processingOptions.getCurrentOption("sectionSubtitleSize"),
                notes_text,
            )
        else:
            # Remove the links from the bullets and replace with target slide title
            for bullet in bullets:
                linkMatch = linkRegex.match(bullet[1])
                bullet[1] = linkMatch.group(1)

            # Create the TOC slide - with these neutralised titles
            slide = createContentSlide(
                presentation,
                slideNumber,
                slideInfo,
            )

            # Postprocess slide to pick up runs - for TOC creation
            body = findBodyShape(slide)
            text_frame = body.text_frame
            for p in text_frame.paragraphs:
                TOCruns.append(p.runs[0])

        # Store the new slide in the list of section slides - for fixing up links
        SectionSlides[titleText] = slide

        return slide

    else:
        slide = addSlide(
            presentation, presentation.slide_layouts[titleOnlyLayout], None
        )
        title = findTitleShape(slide)

    SectionSlides[titleText] = slide

    shapes = slide.shapes

    # Add title if TOC slide. Or delete shape if not
    if titleText == tocTitle:
        # Is TOC slide so add title
        slideTitleBottom, title, flattenedTitle = formatTitle(
            presentation, slide, tocTitle, pageTitleSize, pageSubtitleSize
        )
    else:
        # Is not TOC slide so delete title shape and adjust where title bottom
        # would be
        deleteSimpleShape(title)
        slideTitleBottom = marginBase

    # Get the rectangle the content will draw in
    contentLeft, contentWidth, contentTop, contentHeight = getContentRect(
        presentation, slide, slideTitleBottom, marginBase
    )

    # Create global list of TOC entries
    for bullet in bullets:
        bulletLevel, bulletText, bulletType = bullet
        if bulletLevel == 0:
            # Level 0 is top level so create a TOC entry
            linkText, linkHref = delinkify(bulletText)
            TOCEntries.append([linkText, linkHref])

    TOCEntryCount = len(TOCEntries)

    TOCFontSize = processingOptions.getCurrentOption("TOCFontSize")

    TOCItemHeight = processingOptions.getCurrentOption("TOCItemHeight")

    TOCItemColour = processingOptions.getCurrentOption("TOCItemColour")

    height = Inches(TOCItemHeight)

    if tocStyle == "chevron":
        if height == 0:
            height = Inches(1)

        width = height * 2.5

        entryGap = Inches(-0.5 * height / Inches(1))

        if TOCFontSize == 0:
            TOCFontSize = 14

    elif tocStyle == "circle":
        if height == 0:
            height = Inches(1.25)

        width = height

        entryGap = Inches(0.5)

        if TOCFontSize == 0:
            TOCFontSize = 12

    rowGap = Inches(processingOptions.getCurrentOption("TOCRowGap"))

    TOCEntriesPerRow = int(
        (presentation.slide_width - 2 * marginBase) / (width + entryGap)
    )

    rowCount = 1 + TOCEntryCount / TOCEntriesPerRow

    # Calculate actual TOC height so it can be vertically centred
    TOCHeight = (rowCount * height) + ((rowCount - 1) * rowGap)

    # Calculate where top of TOC should be
    TOCtop = slideTitleBottom + (contentHeight - TOCHeight + height) / 2

    # Calculate actual TOC width
    TOCWidth = TOCEntriesPerRow * (width + entryGap)

    # Calculate where the TOC will start
    TOCleft = (presentation.slide_width - TOCWidth + entryGap) / 2

    x = TOCleft
    y = TOCtop

    AbsoluteTOCEntryNumber = 1

    TOCEntryNumber = 1

    for entry in TOCEntries:
        entryText = entry[0]
        entryHref = entry[1]

        if entryText == titleText:
            wantFilled = False
            wantLink = False
        else:
            wantFilled = True
            wantLink = True

        if tocStyle == "chevron":
            createChevron(
                entryText,
                x,
                y,
                width,
                height,
                wantFilled,
                shapes,
                TOCFontSize,
                wantLink,
                TOCItemColour,
            )

        elif tocStyle == "circle":
            # Create the circle
            createOval(
                entryText,
                x,
                y,
                width,
                height,
                wantFilled,
                shapes,
                TOCFontSize,
                wantLink,
                TOCItemColour,
            )

            # Create half connector to next one - if not last
            if AbsoluteTOCEntryNumber < TOCEntryCount:
                connector = createLine(
                    x + width,
                    y + height / 2,
                    x + width + entryGap / 2,
                    y + height / 2,
                    shapes,
                )

            # Create half connector to previous one - if not first
            if AbsoluteTOCEntryNumber > 1:
                # z =1
                connector = createLine(
                    x - entryGap / 2, y + height / 2, x, y + height / 2, shapes
                )

        # Prepare for the next TOC entry - even if there isn't one
        x = x + width + entryGap

        # If beyond end of line the next TOC entry would be at the start of the next line
        AbsoluteTOCEntryNumber = AbsoluteTOCEntryNumber + 1
        TOCEntryNumber = TOCEntryNumber + 1
        if TOCEntryNumber == TOCEntriesPerRow + 1:
            x = TOCleft
            y = y + rowGap + height
            TOCEntryNumber = 1

    if want_numbers_content is True:
        addFooter(presentation, slideNumber, slide)

    return slide


def createSlide(presentation, slideNumber, slideInfo):
    abstractTitle = processingOptions.getCurrentOption("abstractTitle")
    tocTitle = processingOptions.getCurrentOption("tocTitle")
    tocStyle = processingOptions.getCurrentOption("tocStyle")
    sectionTitleSize = processingOptions.getCurrentOption("sectionTitleSize")
    presTitleSize = processingOptions.getCurrentOption("presTitleSize")
    sectionSubtitleSize = processingOptions.getCurrentOption("sectionSubtitleSize")
    presSubtitleSize = processingOptions.getCurrentOption("presSubtitleSize")
    leftFooterText = processingOptions.getCurrentOption("leftFooterText")
    footerFontSizeSpec = processingOptions.getCurrentOption("footerFontSize")
    middleFooterText = processingOptions.getCurrentOption("middleFooterText")
    rightFooterText = processingOptions.getCurrentOption("rightFooterText")
    sectionFooters = processingOptions.getCurrentOption("sectionFooters")
    liveFooters = processingOptions.getCurrentOption("liveFooters")
    transition = processingOptions.getCurrentOption("transition")

    if slideInfo.blockType in ["content", "code", "table"]:
        if (tocStyle != "") & (tocTitle == slideInfo.titleText):
            # This is a Table Of Contents slide
            slide = createTOCSlide(
                presentation,
                slideNumber,
                slideInfo.titleText,
                slideInfo.bullets,
                tocStyle,
            )
        elif (abstractTitle != "") & (abstractTitle == slideInfo.titleText):
            # This is an abstract slide
            slide = createAbstractSlide(
                presentation,
                slideNumber,
                slideInfo.titleText,
                slideInfo.bullets,
            )
        else:
            # This is an ordinary contents slide
            slide = createContentSlide(
                presentation,
                slideNumber,
                slideInfo,
            )

    elif slideInfo.blockType == "section":
        if tocStyle != "":
            # This is a section slide in TOC style
            slide = createTOCSlide(
                presentation,
                slideNumber,
                slideInfo.titleText,
                slideInfo.bullets,
                tocStyle,
            )
        else:
            slide = createTitleOrSectionSlide(
                presentation,
                slideNumber,
                slideInfo.titleText,
                processingOptions.getCurrentOption("sectionSlideLayout"),
                sectionTitleSize,
                slideInfo.subtitleText,
                sectionSubtitleSize,
                notes_text,
            )

    elif slideInfo.blockType == "title":
        slide = createTitleOrSectionSlide(
            presentation,
            slideNumber,
            slideInfo.titleText,
            processingOptions.getCurrentOption("titleSlideLayout"),
            presTitleSize,
            slideInfo.subtitleText,
            presSubtitleSize,
            notes_text,
        )

    if footerFontSizeSpec == "":
        footerFontSize = Pt(8.0)
    else:
        footerFontSize = Pt(footerFontSizeSpec)
    
    footerBoxTop = prs.slide_height - numbersHeight / 2 - footerFontSize
    footerBoxHeight = footerFontSize * 2

    if slideInfo.blockType in ["title", "section"]:
        if sectionFooters == "yes":
            wantFooters = True
        else:
            wantFooters = False

        if slideInfo.blockType == "section":
            prs.lastSectionTitle = slideInfo.titleText.strip()
            prs.lastSectionSlide = slide
        elif slideInfo.blockType == "title":
            prs.lastPresTitle = slideInfo.titleText.strip()
            prs.lastPresSubtitle = slideInfo.subtitleText.strip()

    else:
        wantFooters = True

    if wantFooters:
        # Left pseudo-footer
        if leftFooterText != "":
            leftFooterMargin = Inches(0.5)
            leftFooterBoxLeft = leftFooterMargin
            leftFooterBoxWidth = prs.slide_width / 3 - leftFooterMargin
            leftFooter = slide.shapes.add_textbox(
                leftFooterBoxLeft, footerBoxTop, leftFooterBoxWidth, footerBoxHeight
            )

            leftFooter.text, wantHyperLink = substituteFooterVariables(
                leftFooterText, liveFooters
            )

            if wantHyperLink:
                createShapeHyperlinkAndTooltip(leftFooter, prs.lastSectionSlide, "")

            for fp in leftFooter.text_frame.paragraphs:
                fp.alignment = PP_ALIGN.LEFT
                fp.font.size = footerFontSize

        # Middle pseudo-footer
        if middleFooterText != "":
            middleFooterBoxLeft = prs.slide_width / 3
            middleFooterBoxWidth = prs.slide_width / 3
            middleFooter = slide.shapes.add_textbox(
                middleFooterBoxLeft, footerBoxTop, middleFooterBoxWidth, footerBoxHeight
            )

            middleFooter.text, wantHyperLink = substituteFooterVariables(
                middleFooterText, liveFooters
            )

            if wantHyperLink:
                createShapeHyperlinkAndTooltip(middleFooter, prs.lastSectionSlide, "")

            for fp in middleFooter.text_frame.paragraphs:
                fp.alignment = PP_ALIGN.CENTER
                fp.font.size = footerFontSize

        # Right pseudo-footer
        if rightFooterText != "":
            rightFooterMargin = Inches(0.25)
            rightFooterBoxLeft = prs.slide_width * 2 / 3
            rightFooterBoxWidth = prs.slide_width / 3 - rightFooterMargin
            rightFooter = slide.shapes.add_textbox(
                rightFooterBoxLeft, footerBoxTop, rightFooterBoxWidth, footerBoxHeight
            )

            rightFooter.text, wantHyperLink = substituteFooterVariables(
                rightFooterText, liveFooters
            )

            if wantHyperLink:
                createShapeHyperlinkAndTooltip(rightFooter, prs.lastSectionSlide, "")

            for fp in rightFooter.text_frame.paragraphs:
                fp.alignment = PP_ALIGN.RIGHT
                fp.font.size = footerFontSize

    slideNumber = slideNumber + 1

    sequence = []

    addSlideTransition(slide, transition)

    return [slideNumber, slide, sequence]


# Add a transition effect - for transitioning INTO the slide.
def addSlideTransition(slide, transitionType):
    # Handle "no transition" case
    if (transitionType == "none") | (transitionType == ""):
        return

    if transitionType in [
        "fracture",
    ]:
        choiceNS = "p15"
    else:
        choiceNS = "p14"

    # Construct first boilerplate XML fragment
    xml = '      <mc:AlternateContent xmlns:mc="' + namespaceURL["mc"] + '">\n'
    xml += (
        "    <mc:Choice xmlns:"
        + choiceNS
        + '="'
        + namespaceURL[choiceNS]
        + '" Requires="'
        + choiceNS
        + '">\n'
    )
    xml += (
        '<p:transition xmlns:p="'
        + namespaceURL["p"]
        + '" xmlns:p14="'
        + namespaceURL["p14"]
        + '" spd="slow" p14:dur="3400">\n'
    )

    # Add in transition element
    if transitionType in [
        "wipe",
    ]:
        xml += "         <p:" + transitionType + " />\n"

    elif transitionType in [
        "push",
    ]:
        xml += "         <p:" + transitionType + ' dir="u"/>\n'

    elif transitionType in [
        "vortex",
    ]:
        xml += "         <p14:" + transitionType + ' dir="r"/>\n'

    elif transitionType in [
        "split",
    ]:
        xml += "         <p:" + transitionType + ' orient="vert"/>\n'

    elif transitionType in [
        "fracture",
    ]:
        xml += '                 <p15:prstTrans prst="fracture" />\n'

    else:
        xml += "         <p14:" + transitionType + " />\n"

    # Construct last boilerplate XML fragment

    xml += """
      </p:transition>
    </mc:Choice>
    <mc:Fallback>
    """

    xml += '      <p:transition xmlns:p="' + namespaceURL["p"] + '" spd="slow">\n'

    if transitionType in [
        "split",
    ]:
        xml += (
            "        <p:"
            + transitionType
            + ' orient="vert" xmlns:p="'
            + namespaceURL["p"]
            + '"/>\n'
        )

    else:
        xml += "        <p:fade />\n"

    xml += """
      </p:transition>
    </mc:Fallback>
  </mc:AlternateContent>
   """

    # Turn this into an XML fragment
    xmlFragment = parse_xml(xml)

    # Add to slide's XML
    slide.element.insert(-1, xmlFragment)


def createTaskSlides(prs, slideNumber, tasks, titleStem):
    tasksPerPage = processingOptions.getCurrentOption("tasksPerPage")

    taskSlideNumber = 0

    taskCount = len(tasks)
    for taskNumber, task in enumerate(tasks):
        if taskNumber % tasksPerPage == 0:
            # Is first task in a page
            if taskNumber > 0:
                # Print a "tasks" slide - as we have one to print out
                taskSlideNumber += 1
                if taskCount > tasksPerPage:
                    # More than one task page
                    title = titleStem + " - " + str(taskSlideNumber)
                else:
                    # Only one task page
                    title = titleStem

                taskBlock = [taskRows]

                slideInfo = SlideInfo(
                    title, "", "table", [], taskBlock, [], [], ["table"]
                )
                slide = createContentSlide(prs, slideNumber, slideInfo)

                # Fix up references to be active links to the slide where the task
                # was declared
                table = findBodyShape(slide).table
                for row in table.rows:
                    cell0Text = row.cells[0].text
                    if cell0Text not in ["Slide", ""]:
                        # First cell refers to a specific slide number - so link to it
                        run = row.cells[0].text_frame.paragraphs[0].runs[0]
                        createRunHyperlinkOrTooltip(
                            run, prs.slides[int(cell0Text) - 2 + templateSlideCount], ""
                        )

                slideNumber += 1

            taskRows = [["Slide", "Due", "Task", "Tags", "Done"]]
            taskRows.append(["-:", ":--:", ":----", ":----", ":--:"])
            old_sNum = 0

        sNum, taskText, dueDate, tags, done = task

        if tags != "":
            # Sort tags - if there are any
            tagList = re.split("[, ]", tags)
            sortedTagList = sorted(tagList)
            tags = str.join(",", sortedTagList)

        if sNum != old_sNum:
            taskRows.append([str(sNum), dueDate, taskText, tags, done])
        else:
            taskRows.append(["", dueDate, taskText, tags, done])
        old_sNum = sNum

    # Print a final "tasks" slide
    taskSlideNumber += 1
    if taskCount > tasksPerPage:
        title = titleStem + " - " + str(taskSlideNumber)
    else:
        title = titleStem

    taskBlock = [taskRows]
    slideInfo = SlideInfo(title, "", "table", [], taskBlock, [], [], ["table"])
    slide = createContentSlide(prs, slideNumber, slideInfo)

    # Fix up references to be active links to the slide where the task
    # was declared
    table = findBodyShape(slide).table
    for row in table.rows:
        cell0Text = row.cells[0].text
        if cell0Text not in ["Slide", ""]:
            # First cell refers to a specific slide number - so link to it
            run = row.cells[0].text_frame.paragraphs[0].runs[0]
            createRunHyperlinkOrTooltip(
                run, prs.slides[int(cell0Text) - 2 + templateSlideCount], ""
            )

    slideNumber += 1


def createGlossarySlides(prs, slideNumber, abbrevDictionary):
    termSlideNumber = 0
    glossarySlides = []

    glossaryTitle = processingOptions.getCurrentOption("glossaryTitle")
    glossaryTerm = processingOptions.getCurrentOption("glossaryTerm")
    glossaryTermsPerPage = processingOptions.getCurrentOption("glossaryTermsPerPage")
    glossaryMeaningWidth = processingOptions.getCurrentOption("glossaryMeaningWidth")
    glossaryMeaning = processingOptions.getCurrentOption("glossaryMeaning")

    termCount = len(abbrevDictionary)

    for termNumber, term in enumerate(sorted(abbrevDictionary.keys())):
        if termNumber % glossaryTermsPerPage == 0:
            # Is first glossary term in a page
            if termNumber > 0:
                # Print a "glossary" slide - as we have one to print out
                termSlideNumber += 1
                if termCount > glossaryTermsPerPage:
                    # More than one glossary page
                    title = glossaryTitle + " - " + str(termSlideNumber)
                else:
                    # Only one glossary page
                    title = glossaryTerm

                glossaryBlock = [glossaryRows]
                slideInfo = SlideInfo(
                    title, "", "table", [], glossaryBlock, [], [], ["table"]
                )
                slide = createContentSlide(prs, slideNumber, slideInfo)

                glossarySlides.append(slide)
                slideNumber += 1

            glossaryRows = [[glossaryTerm, glossaryMeaning]]
            glossaryRows.append([":-", ":" + ("-" * glossaryMeaningWidth)])
            old_sNum = 0

        meaning = abbrevDictionary.get(term)

        glossaryRows.append([term, meaning])

    # Print a final "glossary" slide
    termSlideNumber += 1
    if termCount > glossaryTermsPerPage:
        # More than one glossary page
        title = glossaryTitle + " - " + str(termSlideNumber)
    else:
        # Only one glossary page
        title = glossaryTitle

    glossaryBlock = [glossaryRows]
    slideInfo = SlideInfo(title, "", "table", [], glossaryBlock, [], [], ["table"])
    slide = createContentSlide(prs, slideNumber, slideInfo)
    glossarySlides.append(slide)
    slideNumber += 1

    return slideNumber, glossarySlides


def createSlideNotes(slide, notes_text):
    # Remove surrounding white space
    notes_text = notes_text.strip().lstrip("\n")

    if slide.notes_slide.notes_text_frame.text != "":
        # Notes already filled in
        return

    if notes_text != "":
        # There is substantive slide note text so create the note
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame

        # addFormattedText handles eg hyperlinks and entity references
        addFormattedText(text_frame.paragraphs[0], notes_text)


def createFootnoteSlides(prs, slideNumber, footnoteDefinitions):
    footnotesSlideNumber = 0
    footnoteSlides = []

    footnotesTitle = processingOptions.getCurrentOption("footnotesTitle")
    footnotesPerPage = processingOptions.getCurrentOption("footnotesPerPage")

    footnoteCount = len(footnoteDefinitions)

    for footnoteNumber, footnote in enumerate(footnoteDefinitions):
        if footnoteNumber % footnotesPerPage == 0:
            # Is first footnote in a page
            if footnoteNumber > 0:
                # Print a "footnotes" slide - as we have one to print out
                footnotesSlideNumber += 1
                if footnoteCount > footnotesPerPage:
                    # More than one footnotes page
                    title = footnotesTitle + " - " + str(footnotesSlideNumber)
                else:
                    # Only one footnotes page
                    title = footnotesTitle

                slideInfo = SlideInfo(
                    title, "", "content", bullets, [], cards, [], ["list"]
                )
                slideNumber, slide, sequence = createSlide(prs, slideNumber, slideInfo)

                footnoteSlides.append(slide)

                # Turn off bulleting
                removeBullets(findBodyShape(slide).text_frame)

                slideNumber += 1
            bullets = []
            old_sNum = 0

        bullets.append(
            [
                1,
                str(footnoteNumber + 1) + ". " + footnoteDefinitions[footnoteNumber][1],
                "bullet",
            ]
        )

    # Print a final "footnote" slide
    footnotesSlideNumber += 1
    if footnoteCount > footnotesPerPage:
        # More than one footnotes page
        title = footnotesTitle + " - " + str(footnotesSlideNumber)
    else:
        # Only one footnotes page
        title = footnotesTitle

    slideInfo = SlideInfo(title, "", "content", bullets, [], cards, [], ["list"])
    slideNumber, slide, sequence = createSlide(prs, slideNumber, slideInfo)

    footnoteSlides.append(slide)

    # Turn off bulleting
    removeBullets(findBodyShape(slide).text_frame)

    slideNumber += 1

    return slideNumber, footnoteSlides

def writeMdFile(content):
    now = time.strftime("%Y-%m-%d-%H_%M_%S", time.localtime(time.time()))
    filename =os.getcwd() + os.sep + 'md2pptx' + os.sep + 'data' + os.sep + 'md' + now + '.md'
    # if not os.path.exists(filename):
    #     os.mknod(filename)
    print(os.getcwd())
    with open(filename, 'w', encoding='GB2312') as file:
        file.write(content)
    return filename

slideHrefRegex = re.compile("(.+)\[(.+)\]$")
processingOptions = ProcessingOptions()
slideNumber = 1

bulletRegex = re.compile("^(\s)*(\*)(.*)")
numberRegex = re.compile("^(\s)*(\d+)\.(.*)")
metadataRegex = re.compile("^(.+):(.+)")

graphicRE = "!\[(.*?)\]\((.+?)\)"
graphicRegex = re.compile(graphicRE)

clickableGraphicRE = "\[" + graphicRE + "\]\((.+?)\)"
clickableGraphicRegex = re.compile(clickableGraphicRE)

videoRE = "<video (.*?)></video>"
videoRegex = re.compile(videoRE)

audioRE = "<audio (.*?)></audio>"
audioRegex = re.compile(audioRE)

linkRegex = re.compile("^\[(.+)\]\((.+)\)")
footnoteDefinitionRegex = re.compile("^\[\^(.+?)\]: (.+)")
spanClassRegex = re.compile("<span( )*class=")
spanStyleRegex = re.compile("<span( )*style=")
slideHrefRegex = re.compile("(.+)\[(.+)\]$")
anchorRegex = re.compile("^<a id=[\"'](.+)[\"']></a>")
dynamicMetadataRegex = re.compile("^<!-- md2pptx: (.+): (.+) -->")
indirectReferenceAnchorRegex = re.compile("^\[(.+?)\]: (.+)")
indirectReferenceUsageRegex = re.compile("[(.+?)]\[(.+?)]")

inBlock = False
inList = False
inTable = False
inCard = False
inTitle = False

blockType = ""
slideTitle = ""
slideSubtitle = ""
bullets = []
tableRows = []
cards = []
code = []
inCode = False
inHTMLCode = False
inFencedCode = False
notes_text = ""
slide = None
tasks = []
sequence = []

slideHrefs = {}
href_runs = {}

# Each of these is a picture, then a href, then a tooltip - as a tuple
pictureInfos = []

# Pass 2: Concatenate lines with continuations
previousLine = "\n"
linesAfterConcatenation = []

want_numbers_headings = False
want_numbers_content = False
prs = Presentation()
numbersHeight = 0
footerFontSize = 0
maxBlocks = 10

banner = (
        "md2pptx Markdown To Powerpoint Converter " + md2pptx_level + " " + md2pptx_date
    )

def genPPTX(outputfile, inputfile):
    start_time = time.time()

    banner = (
        "md2pptx Markdown To Powerpoint Converter " + md2pptx_level + " " + md2pptx_date
    )

    bannerUnderline = ""
    for i in range(len(banner)):
        bannerUnderline = bannerUnderline + "="

    print("\n" + banner + "\n" + bannerUnderline)
    print("\nOpen source project: https://github.com/MartinPacker/md2pptx")

    print("\nPython: " + platform.python_version())

    print("python-pptx: " + pptx_version)

    if have_pillow:
        print("Pillow: " + PIL.__version__)
    else:
        print("Pillow: Not Installed")

    if have_cairosvg:
        print("CairoSVG: " + cairosvg.__version__)
    else:
        print("CairoSVG: Not Installed")

    if have_graphviz:
        print("graphviz: " + graphviz.__version__)
    else:
        print("graphviz: Not Installed")

    print(f"funnel: {funnel.__version__}")

    input_file = []

    # processingOptions = ProcessingOptions()

    # if len(sys.argv) > 2:
    #     # Have input file as well as output file
    #     input_filename = sys.argv[1]
    #     output_filename = sys.argv[2]
    #
    #     if Path(input_filename).exists():
    #         input_path = Path(input_filename)
    #         with input_path.open() as file:
    #             input_file = file.readlines()
    #     else:
    #         print("Input file specified but does not exist. Terminating.")
    # elif len(sys.argv) == 1:
    #     print("No parameters. Terminating")
    #     sys.exit()
    # else:
    #     output_filename = sys.argv[1]
    #
    #     input_file = sys.stdin.readlines()
    #
    # if len(input_file) == 0:
    #     print("Empty input file. Terminating")
    #     sys.exit()

    input_filename = inputfile
    output_filename = outputfile

    if Path(input_filename).exists():
        input_path = Path(input_filename)
        with input_path.open() as file:
            input_file = file.readlines()
    else:
        print("Input file specified but does not exist. Terminating.")

    slideNumber = 1

    bulletRegex = re.compile("^(\s)*(\*)(.*)")
    numberRegex = re.compile("^(\s)*(\d+)\.(.*)")
    metadataRegex = re.compile("^(.+):(.+)")

    graphicRE = "!\[(.*?)\]\((.+?)\)"
    graphicRegex = re.compile(graphicRE)

    clickableGraphicRE = "\[" + graphicRE + "\]\((.+?)\)"
    clickableGraphicRegex = re.compile(clickableGraphicRE)

    videoRE = "<video (.*?)></video>"
    videoRegex = re.compile(videoRE)

    audioRE = "<audio (.*?)></audio>"
    audioRegex = re.compile(audioRE)

    linkRegex = re.compile("^\[(.+)\]\((.+)\)")
    footnoteDefinitionRegex = re.compile("^\[\^(.+?)\]: (.+)")
    spanClassRegex = re.compile("<span( )*class=")
    spanStyleRegex = re.compile("<span( )*style=")
    slideHrefRegex = re.compile("(.+)\[(.+)\]$")
    anchorRegex = re.compile("^<a id=[\"'](.+)[\"']></a>")
    dynamicMetadataRegex = re.compile("^<!-- md2pptx: (.+): (.+) -->")
    indirectReferenceAnchorRegex = re.compile("^\[(.+?)\]: (.+)")
    indirectReferenceUsageRegex = re.compile("[(.+?)]\[(.+?)]")

    # Default slide layout enumeration
    processingOptions.setOptionValuesArray(
        [
            ["titleSlideLayout", 0],
            ["sectionSlideLayout", 1],
            ["contentSlideLayout", 2],
            ["titleOnlyLayout", 5],
            ["blanklayout", 6],
        ]
    )

    # Abbreviation Dictionary
    abbrevDictionary = {}

    # Abbreviation Runs Dictionary
    abbrevRunsDictionary = {}

    # Footnote runs Dictionary
    footnoteRunsDictionary = {}

    # Extract metadata
    metadata_lines = []
    afterMetadataAndHTML = []


    TOCruns = []
    SectionSlides = {}

    inMetadata = True
    in_comment = False
    inHTML = False
    inCode = False

    # Pass 1: Strip out comments and metadata, storing the latter
    for line in input_file:
        if line.lstrip().startswith("<!-- md2pptx: "):
            # md2pptx dynamic metadata so keep it
            afterMetadataAndHTML.append(line)

        if line.lstrip().startswith("<!--"):
            if line.rstrip().endswith("-->"):
                # Note: Not taking text after end of comment
                continue
            else:
                in_comment = True
                continue

        elif line.rstrip().endswith("-->"):
            # Note: Not taking text after end of comment
            in_comment = False
            continue

        elif in_comment is True:
            continue

        elif (line.lstrip()[:1] == "<") & (inCode is False):
            lineLstrip = line.lstrip()
            if startswithOneOf(lineLstrip, ["<a id=", "<span "]):
                # Line goes to post-metadata array
                afterMetadataAndHTML.append(line)

            elif startswithOneOf(lineLstrip, ["<code>", "<pre>"]):
                inCode = True
                afterMetadataAndHTML.append(line)

            elif startswithOneOf(lineLstrip, ["</code>", "</pre>"]):
                inCode = False
                afterMetadataAndHTML.append(line)

            elif startswithOneOf(lineLstrip, ["<video ", "<audio "]):
                # Line goes to post-metadata array
                afterMetadataAndHTML.append(line)

            else:
                inHTML = True

            continue

        elif line.startswith("```"):
            inCode = ~inCode
            # afterMetadataAndHTML.append(line)

        elif line.lstrip()[:1] == "#":
            # Heading has triggered end of metadata and end of HTML
            inMetadata = False
            inHTML = False

        elif inHTML:
            continue

        elif inCode:
            afterMetadataAndHTML.append(line)
            continue

        elif line == "\n":
            # Blank line has triggered end of metadata
            inMetadata = False

        if inMetadata is True:
            # Line goes to metadata array
            metadata_lines.append(line)

        else:
            # Line goes to post-metadata array
            afterMetadataAndHTML.append(line)

    want_numbers_headings = False
    want_numbers_content = False

    processingOptions.setOptionValues("slideTemplateFile", "")
    processingOptions.setOptionValues("tempDir", None)

    ######################################################################################
    #                                                                                    #
    # Set default, presentation and current values for some key options                  #
    #                                                                                    #
    ######################################################################################

    processingOptions.setOptionValuesArray(
        [
            ["pageTitleSize", 30],
            ["pageSubtitleSize", "same"],
            ["pageTitleAlign", "left"],
        ]
    )

    processingOptions.setOptionValues("baseTextSize", 18)

    processingOptions.setOptionValues("baseTextDecrement", 2)

    # Code defaults
    processingOptions.setOptionValuesArray(
        [
            ["codeForeground", "000000"],
            ["codeBackground", "DFFFDF"],
            ["codeColumns", 80],
            ["fixedPitchHeightWidthRatio", 1.2],
        ]
    )

    # Text defaults
    processingOptions.setOptionValuesArray(
        [
            ["italicItalic", True],
            ["italicColour", ("None", "")],
            ["boldBold", True],
            ["boldColour", ("None", "")],
        ]
    )

    # Tables defaults
    processingOptions.setOptionValuesArray(
        [
            ["compactTables", 0],
            ["addTableLines", "no"],
            ["addTableColumnLines", []],
            ["addTableRowLines", []],
            ["addTableLineWidth", 1],
            ["addTableLineCount", 1],
            ["addTableLineColour", "000000"],
            ["tableMargin", Inches(0.2)],
            ["spanCells", "yes"],
            ["tableHeadingSize", 0],
        ]
    )

    # Cards defaults
    processingOptions.setOptionValuesArray(
        [
            ["cardPercent", 80],
            ["cardLayout", "horizontal"],
            ["cardTitleAlign", "c"],
            ["cardTitlePosition", "above"],
            ["cardShape", "rounded"],
            ["horizontalCardGap", 0.25],
            ["verticalCardGap", 0.15],
            ["cardShadow", False],
            ["cardTitleSize", 0],
            ["cardBorderWidth", 0],
            ["cardBorderColour", ("None", "")],
            ["cardTitleColour", ("None", "")],
            ["cardColour", [("None", "")]],
            ["cardTitleBackground", [("None", "")]],
            ["cardDividerColour", ("RGB", "#000000")],
            ["cardGraphicHeight", 0.75],
            ["cardGraphicPosition", "before"],
        ]
    )


    processingOptions.setOptionValues("contentSplit", [1, 1, 1, 1, 1, 1, 1, 1, 1, 1])

    processingOptions.setOptionValues("contentSplitDirection", "vertical")

    # Number of spaces a single level of indentation is
    processingOptions.setOptionValues("indentSpaces", 2)

    # Whether titles are adjusted or not
    processingOptions.setOptionValues("adjustTitles", True)

    # Title and section defaults
    processingOptions.setOptionValuesArray(
        [
            ["sectionTitleSize", 40],
            ["sectionSubtitleSize", 28],
            ["presTitleSize", 40],
            ["presSubtitleSize", 28],
            ["sectionsExpand", False],
        ]
    )

    processingOptions.setOptionValues("monoFont", "Courier")

    topHeadingLevel = 1
    titleLevel = topHeadingLevel
    sectionLevel = titleLevel + 1
    contentLevel = sectionLevel + 1
    cardLevel = contentLevel + 1

    # Abstracts defaults
    abstractTitle = processingOptions.setOptionValues("abstractTitle", "")

    # Tasks defaults
    processingOptions.setOptionValuesArray([["taskSlides", "all"], ["tasksPerPage", 20]])

    # Glossary defaults
    processingOptions.setOptionValuesArray(
        [
            ["glossaryTitle", "Title"],
            ["glossaryTerm", "Term"],
            ["glossaryMeaning", "Meaning"],
            ["glossaryMeaningWidth", 5],
            ["glossaryTermsPerPage", 20],
        ]
    )

    # Footnotes defaults
    processingOptions.setOptionValuesArray(
        [["footnotesTitle", "Footnotes"], ["footnotesPerPage", 20]]
    )

    # Table Of Contents defaults
    processingOptions.setOptionValuesArray(
        [
            ["tocTitle", "Topics"],
            ["tocStyle", ""],
            ["tocItemHeight", 0],
            ["tocItemColour", ""],
            ["tocRowGap", 0.75],
            ["tocFontSize", 0],
            ["tocLinks", False],
            ["sectionArrows", False],
            ["sectionArrowsColour", ""],
        ]
    )

    processingOptions.setOptionValues("marginBase", Inches(0.2))

    processingOptions.setOptionValues("transition", "none")

    processingOptions.setOptionValues("deleteFirstSlide", False)

    processingOptions.setOptionValuesArray(
        [
            ["leftFooterText", ""],
            ["middleFooterText", ""],
            ["rightFooterText", ""],
            ["sectionFooters", "no"],
            ["liveFooters", "no"],
            ["footerFontSize", ""],
        ]
    )

    TOCEntries = []


    metadata = []


    # Space to leave at bottom if numbers
    numbersHeight = Inches(0.4)

    # If no numbers leave all the above height anyway
    processingOptions.setOptionValuesArray(
        [
            ["numbersHeight", numbersHeight],
            ["numbersContentMargin", numbersHeight],
            ["numbersHeadingsMargin", numbersHeight],
            ["numbersFontSize", ""],
        ]
    )

    # Graphics options
    processingOptions.setOptionValuesArray(
        [
            ["exportGraphics", False],
        ]
    )

    # Funnel options
    processingOptions.setOptionValuesArray(
        [
            [
                "funnelColours",
                [
                    ("Theme", MSO_THEME_COLOR.ACCENT_1),
                    ("Theme", MSO_THEME_COLOR.ACCENT_2),
                    ("Theme", MSO_THEME_COLOR.ACCENT_3),
                    ("Theme", MSO_THEME_COLOR.ACCENT_4),
                    ("Theme", MSO_THEME_COLOR.ACCENT_5),
                    ("Theme", MSO_THEME_COLOR.ACCENT_6),
                ],
            ],
            ["funnelBorderColour", ("None", "")],
            ["funnelTitleColour", ("None", "")],
            ["funnelTextColour", ("None", "")],
            ["funnelLabelsPercent", 10],
            ["funnelLabelPosition", "before"],
            ["funnelWidest", "left"],
        ]
    )

    ######################################################################################
    #                                                                                    #
    #  Prime for style. metadata                                                         #
    #                                                                                    #
    ######################################################################################

    # Background colour class correspondence
    bgcolors = {}

    # Foreground colour class correspondence
    fgcolors = {}

    # Emphases class correspondence
    emphases = {}

    ######################################################################################
    #                                                                                    #
    #  Prime for footnotes                                                               #
    #                                                                                    #
    ######################################################################################

    # List of footnote definitions. Each is a (ref, text) pair.
    # Also array of names - for quick searching
    footnoteDefinitions = []
    footnoteReferences = []

    maxBlocks = 10

    ######################################################################################
    #                                                                                    #
    #  Parse metadata and report on the items found, setting options                     #
    #                                                                                    #
    ######################################################################################

    if len(metadata_lines) > 0:
        print("")
        print("Metadata:")
        print("=========")
        print("")
        print("Name".ljust(40) + " " + "Value")
        print("----".ljust(40) + " " + "-----")


    for line in metadata_lines:
        matchInfo = metadataRegex.match(line)

        if matchInfo is None:
            print("Ignoring invalid metadata line: " + line)
            continue

        name = matchInfo.group(1).strip()
        value = matchInfo.group(2).strip()
        metadata.append([name, value])

        # Print name as it was typed
        print(name.ljust(40) + " " + value)

        # Lower case name for checking
        name = name.lower()

        if name == "numbers":
            numbersHeight = processingOptions.getCurrentOption("numbersHeight")
            if value.lower() == "yes":
                # Want slide numbers everywhere
                want_numbers_headings = True
                processingOptions.setOptionValues("numbersHeadingsMargin", numbersHeight)

                want_numbers_content = True
                processingOptions.setOptionValues("numbersContentMargin", numbersHeight)
            elif value.lower() == "content":
                # Want slide numbers on content slides but not headings & sections
                want_numbers_headings = False
                want_numbers_content = True
                processingOptions.setOptionValues("numbersContentMargin", numbersHeight)
            else:
                # Don't want slide numbers - but they could still be added by slide master
                # (Can code any other value, including 'no' or omit this metadata type)
                want_numbers_headings = False
                want_numbers_content = False

        elif name == "numbersfontsize":
            processingOptions.setOptionValues(name, float(value))

        elif name == "pagetitlesize":
            processingOptions.setOptionValues(name, float(value))

        elif name == "pagetitlealign":
            if value in ["left", "right", "center", "centre", "l", "r", "c"]:
                if value[:1] == "l":
                    processingOptions.setOptionValues(name, "left")
                elif value[:1] == "r":
                    processingOptions.setOptionValues(name, "right")
                elif value[:1] == "c":
                    processingOptions.setOptionValues(name, "center")
            else:
                print(
                    f'PageTitleAlign value \'{value}\' unsupported. "left", "right", "centre", or "center" required.'
                )

        elif name == "pagesubtitlesize":
            if value == "same":
                processingOptions.setOptionValues(name, value)
            else:
                processingOptions.setOptionValues(name, float(value))

        elif name == "basetextsize":
            processingOptions.setOptionValues(name, float(value))

        elif name == "basetextdecrement":
            processingOptions.setOptionValues(name, float(value))

        elif name in [
            "sectiontitlesize",
            "sectionsubtitlesize",
            "prestitlesize",
            "pressubtitlesize",
        ]:
            processingOptions.setOptionValues(name, float(value))

        elif name == "deletefirstslide":
            if value.lower() == "yes":
                processingOptions.setOptionValues(name, True)
            else:
                processingOptions.setOptionValues(name, False)

        elif name == "sectionsexpand":
            if value.lower() == "yes":
                processingOptions.setOptionValues(name, True)
            else:
                processingOptions.setOptionValues(name, False)

        elif (name == "template") | (name == "master"):
            if value == "Martin Master.pptx":
                slideTemplateFile = "Martin Template.pptx"
            else:
                slideTemplateFile = value
            processingOptions.setOptionValues("slideTemplateFile", slideTemplateFile)

        elif name == "monofont":
            processingOptions.setOptionValues(name, value)

        elif name == "marginbase":
            processingOptions.setOptionValues(name, Inches(float(value)))

        elif name == "tablemargin":
            processingOptions.setOptionValues(name, Inches(float(value)))

        elif name == "tocstyle":
            if value in ["chevron", "circle", "plain"]:
                processingOptions.setOptionValues(name, value)
            else:
                print(
                    f'TOCStyle value \'{value}\' unsupported. "chevron" or "circle" required.'
                )

        elif name == "toctitle":
            processingOptions.setOptionValues(name, value)

        elif name == "tocitemheight":
            processingOptions.setOptionValues(name, float(value))

        elif (name == "tocitemcolour") | (name == "tocitemcolor"):
            processingOptions.setOptionValues("tocItemColour", value)

        elif name == "toclinks":
            if value.lower() == "yes":
                processingOptions.setOptionValues(name, True)

        elif name == "sectionarrows":
            if value.lower() == "yes":
                processingOptions.setOptionValues(name, True)

        elif (name == "sectionarrowscolour") | (name == "sectionarrowscolor"):
            processingOptions.setOptionValues("sectionArrowsColour", value)

        elif name == "tocrowgap":
            processingOptions.setOptionValues(name, float(value))

        elif name == "tocfontsize":
            processingOptions.setOptionValues(name, float(value))

        elif name == "compacttables":
            processingOptions.setOptionValues(name, float(value))

        elif name == "tableheadingsize":
            processingOptions.setOptionValues(name, float(value))

        elif name == "abstracttitle":
            processingOptions.setOptionValues(name, value)

        elif name in [
            "leftfootertext",
            "middlefootertext",
            "rightfootertext",
            "sectionfooters",
            "livefooters",
        ]:
            processingOptions.setOptionValues(name, value)

        elif name == "footerfontsize":
            processingOptions.setOptionValues(name, float(value))

        elif name == "boldbold":
            if value.lower() == "no":
                processingOptions.setOptionValues("boldBold", False)

        elif (name == "boldcolour") | (name == "boldcolor"):
            processingOptions.setOptionValues("boldColour", (parseColour(value.strip())))

        elif name == "italicitalic":
            if value == "no":
                processingOptions.setOptionValues("italicItalic", False)

        elif (name == "italiccolour") | (name == "italiccolor"):
            processingOptions.setOptionValues("italicColour", (parseColour(value.strip())))

        elif name in ["cardcolour", "cardcolor", "cardcolours", "cardcolors"]:
            valueArray2 = [parseColour(c.strip()) for c in value.split(",")]

            processingOptions.setOptionValues("cardColour", valueArray2)

        elif name in ["cardtitlebackground", "cardtitlebackgrounds"]:
            valueArray2 = [parseColour(c.strip()) for c in value.split(",")]

            processingOptions.setOptionValues("cardTitleBackground", valueArray2)

        elif (name == "cardbordercolour") | (name == "cardbordercolor"):
            processingOptions.setOptionValues(
                "cardBorderColour", parseColour(value.strip())
            )

        elif (name == "cardtitlecolour") | (name == "cardtitlecolor"):
            processingOptions.setOptionValues("cardTitleColour", parseColour(value.strip()))

        elif (name == "carddividercolour") | (name == "carddividercolor"):
            processingOptions.setOptionValues(
                "cardDividerColour", parseColour(value.strip())
            )

        elif name == "cardborderwidth":
            processingOptions.setOptionValues(name, float(value))

        elif name == "cardtitlesize":
            processingOptions.setOptionValues(name, float(value))

        elif name == "cardshadow":
            if value.lower() == "yes":
                processingOptions.setOptionValues(name, True)

        elif name in ["cardpercent", "cardgraphicheight"]:
            processingOptions.setOptionValues(name, float(value))

        elif name == "cardlayout":
            if value in ["horizontal", "vertical"]:
                processingOptions.setOptionValues(name, value)
            else:
                print(
                    f'CardLayout value \'{value}\' unsupported. "horizontal" or "vertical" required.'
                )

        elif name == "cardshape":
            if value in ["squared", "rounded", "line"]:
                processingOptions.setOptionValues(name, value)
            else:
                print(
                    f'CardShape value \'{value}\' unsupported. "squared", "rounded", or "line" required.'
                )

        elif name == "cardtitleposition":
            if value in ["above", "inside", "before", "after"]:
                processingOptions.setOptionValues(name, value)
            else:
                print(
                    f'CardTitlePosition value \'{value}\' unsupported. "inside", "above", "before", or "after" required.'
                )

        elif name == "cardgraphicposition":
            if value in ["before", "after"]:
                processingOptions.setOptionValues(name, value)
            else:
                print(
                    f'CardGraphicPosition value \'{value}\' unsupported. "before",or "after" required.'
                )

        elif name == "cardtitlealign":
            val1l = value[:1].lower()
            if val1l in ["l", "r", "c"]:
                processingOptions.setOptionValues(name, val1l)
            else:
                print(f"CardAlign value '{value}' unsupported.")

        elif name in ["horizontalcardgap", "verticalcardgap"]:
            processingOptions.setOptionValues(name, float(value))
            print(float(value))

        elif name == "contentsplit":
            splitValue = value.split()
            cs = []
            for v in splitValue:
                cs.append(int(v))

            # Extend to maximum allowed
            needMore = maxBlocks - len(cs)
            for _ in range(needMore):
                cs.append(1)

            processingOptions.setOptionValues("contentSplit", cs)

        elif name in ["contentsplitdirection", "contentsplitdirn"]:
            if value in [
                "vertical",
                "horizontal",
                "v",
                "h",
                "default",
                "pres",
                "pop",
                "prev",
            ]:
                if value in ["vertical", "horizontal"]:
                    adjustedValue = value
                elif value == "v":
                    adjustedValue = "vertical"
                elif value == "h":
                    adjustedValue = "horizontal"
                else:
                    adjustedValue = value

                processingOptions.setOptionValues("contentSplitDirection", adjustedValue)

            else:
                print(
                    f'{name} value \'{value}\' unsupported. "vertical" or "horizontal" required.'
                )

        elif name == "taskslides":
            processingOptions.setOptionValues(name, value)

        elif name == "tasksperpage":
            processingOptions.setOptionValues(name, int(value))

        elif name in [
            "titleslidelayout",
            "sectionslidelayout",
            "contentslidelayout",
            "titleonlylayout",
            "blanklayout",
        ]:
            processingOptions.setOptionValues(name, int(value))

        elif name == "numbersheight":
            numbersHeight = Inches(float(value))

            # If no numbers leave all the above height anyway
            processingOptions.setOptionValues("numbersHeight", numbersHeight)
            processingOptions.setOptionValues("numbersContentMargin", numbersHeight)
            processingOptions.setOptionValues("numbersHeadingsMargin", numbersHeight)

        elif name == "glossarytitle":
            processingOptions.setOptionValues(name, value)

        elif name == "glossaryterm":
            processingOptions.setOptionValues(name, value)

        elif name == "glossarymeaning":
            processingOptions.setOptionValues(name, value)

        elif name == "glossarymeaningwidth":
            processingOptions.setOptionValues(name, int(value))

        elif name == "glossarytermsperpage":
            processingOptions.setOptionValues(name, int(value))

        elif name == "footnotesperpage":
            processingOptions.setOptionValues(name, int(value))

        elif name == "footnotestitle":
            processingOptions.setOptionValues(name, value)

        # Following relate to styling and don't use Processing Options class

        elif name.startswith("style.bgcolor."):
            spanClass = name[14:]
            if value.startswith("#"):
                value2 = value
            else:
                value2 = "#" + value
            (check, colour) = parseRGB(value2)
            if check:
                # Valid RGB hex value
                bgcolors[spanClass] = colour
            else:
                print(f"Invalid value for {name} - {value}. Ignoring.")

        elif name.startswith("style.fgcolor."):
            spanClass = name[14:]
            if value.startswith("#"):
                value2 = value
            else:
                value2 = "#" + value
            (check, colour) = parseRGB(value2)
            if check:
                # Valid RGB hex value
                fgcolors[spanClass] = colour
            else:
                print(f"Invalid value for {name} - {value}. Ignoring.")

        elif name.startswith("style.emphasis."):
            spanClass = name[15:]
            emphases[spanClass] = value

        elif name in ["codeforeground", "codebackground"]:
            processingOptions.setOptionValues(name, value)

        elif name == "fpratio":
            processingOptions.setOptionValues("fixedPitchHeightWidthRatio", float(value))

        elif name == "codecolumns":
            processingOptions.setOptionValues(name, int(value))

        elif name == "topheadinglevel":
            titleLevel = int(value)
            sectionLevel = titleLevel + 1
            contentLevel = sectionLevel + 1
            cardLevel = contentLevel + 1

        elif name == "indentspaces":
            processingOptions.setOptionValues(name, int(value))

        elif name == "adjusttitles":
            if value == "no":
                processingOptions.setOptionValues(name, False)

        elif name in ["addtablelines", "addtablelinecolour"]:
            processingOptions.setOptionValues(name, value)

        elif name in ["addtablecolumnlines", "addtablerowlines"]:
            processingOptions.setOptionValues(name, sortedNumericList(value))

        elif name in ["addtablelinecount", "addtablelinewidth"]:
            processingOptions.setOptionValues(name, int(value))

        elif name == "spancells":
            processingOptions.setOptionValues(name, value)

        elif name == "hidemetadata":
            if value == "style":
                processingOptions.hideMetadataStyle = True

        elif name == "exportgraphics":
            if value.lower() == "yes":
                processingOptions.setOptionValues(name, True)

        elif name == "tempdir":
            processingOptions.setOptionValues(name, os.path.expanduser(value))

        elif name == "transition":
            if value.lower() in [
                "none",
                "ripple",
                "reveal",
                "honeycomb",
                "shred",
                "wipe",
                "push",
                "vortex",
                "split",
                "fracture",
            ]:
                # Valid transition name
                processingOptions.setOptionValues(name, value.lower())
            else:
                print(f"Invalid value for {name} - {value}. Ignoring.")

        elif (name == "funnelcolours") | (name == "funnelcolors"):
            valueArray2 = [parseColour(c.strip()) for c in value.split(",")]

            processingOptions.setOptionValues("funnelColours", valueArray2)

        elif (name == "funnelbordercolour") | (name == "funnelbordercolor"):
            processingOptions.setOptionValues(
                "funnelBorderColour", parseColour(value.strip())
            )

        elif (name == "funneltitlecolour") | (name == "funneltitlecolor"):
            processingOptions.setOptionValues(
                "funnelTitleColour", parseColour(value.strip())
            )

        elif (name == "funneltextcolour") | (name == "funneltextcolor"):
            processingOptions.setOptionValues(
                "funnelTextColour", parseColour(value.strip())
            )

        elif name == "funnellabelspercent":
            processingOptions.setOptionValues(name, float(value))

        elif name == "funnellabelposition":
            if value in ["before", "after"]:
                processingOptions.setOptionValues(name, value)
            else:
                print(
                    f'funnelLabelPosition value \'{value}\' unsupported. "before" or "after" required.'
                )

        elif name == "funnelwidest":
            if value in ["left", "right", "pipe", "hpipe", "vpipe", "top", "bottom"]:
                processingOptions.setOptionValues(name, value)
            else:
                print(
                    f'funnelLabelPosition value \'{value}\' unsupported. "left", "right", "pipe", or "hpipe" required.'
                )

    slideTemplateFile = processingOptions.getCurrentOption("slideTemplateFile")
    if slideTemplateFile != "":
        originalSlideTemplateFile = slideTemplateFile
        if Path(os.path.expanduser(slideTemplateFile)).exists():
            # We can successfully expand the path to pick up the file
            slideTemplateFile = os.path.expanduser(slideTemplateFile)
        else:
            # Slide template file is not present if we expand path
            script_path = os.path.dirname(__file__)
            slideTemplateFile = script_path + os.sep + slideTemplateFile
            if not Path(slideTemplateFile).exists():
                print(
                    f"\nTemplate file {originalSlideTemplateFile} does not exist. Terminating."
                )
                # sys.exit()

        print(f"\nUsing {slideTemplateFile} as base for presentation")

    if slideTemplateFile == "":
        # Use default slide deck that comes with python-pptx as base
        prs = Presentation()
        print("\nNo slide to document metadata on. Continuing without it.")

        templateSlideCount = 0
    else:
        # Use user-specified presentation as base
        prs = Presentation(slideTemplateFile)

        # If there is a slide to use fill it with metadata
        templateSlideCount = len(prs.slides)
        if templateSlideCount > 0:
            print("\nWriting processing summary slide with metadata on it.")
        else:
            print("\nNo slide to document metadata on. Continuing without it.")

        # Prime template slides with slideInfo as None
        for slide in prs.slides:
            slide.slideInfo = None

    # Following might be used in slide footers
    prs.lastSectionTitle = ""
    prs.lastSectionSlide = None
    prs.lastPresTitle = ""
    prs.lastPresSubtitle = ""

    print("")
    print("Slides:")
    print("=======")
    print("")

    inBlock = False
    inList = False
    inTable = False
    inCard = False
    inTitle = False

    blockType = ""
    slideTitle = ""
    slideSubtitle = ""
    bullets = []
    tableRows = []
    cards = []
    code = []
    inCode = False
    inHTMLCode = False
    inFencedCode = False
    notes_text = ""
    slide = None
    tasks = []
    sequence = []

    slideHrefs = {}
    href_runs = {}

    # Each of these is a picture, then a href, then a tooltip - as a tuple
    pictureInfos = []

    # Pass 2: Concatenate lines with continuations
    previousLine = "\n"
    linesAfterConcatenation = []

    for line in afterMetadataAndHTML:
        if startswithOneOf(line, ["<pre>", "<code>"]):
            # These are around code lines
            linesAfterConcatenation.append(line)
            inHTMLCode = True

        elif startswithOneOf(line, ["</pre>", "</code>"]):
            # These are around code lines
            linesAfterConcatenation.append(line)
            inHTMLCode = False

        elif line.startswith("```"):
            linesAfterConcatenation.append(line)
            inCode = not inCode

        elif line == "\n":
            # This is a blank line so copy it over
            linesAfterConcatenation.append(line)

        elif previousLine == "\n":
            # Previous line was blank so copy this one over
            linesAfterConcatenation.append(line)

        elif line.startswith("<!-- md2pptx: "):
            # This is a dynamic metadata line so keep it separate
            linesAfterConcatenation.append(line)

        elif line.startswith("<a id="):
            # This is an anchor line so keep it separate
            linesAfterConcatenation.append(line)

        elif startswithOneOf(line, ["<video ", "<audio "]):
            # This is a video / audio element line
            linesAfterConcatenation.append(line)

        elif line.lstrip() == "":
            # This is an empty line
            linesAfterConcatenation.append(line)

        else:
            # Previous line was not blank and nor is this one so consider concatenation
            if line.lstrip()[0] not in "*#\|0123456789!":
                if (
                    (previousLine[0:2] != "# ")
                    & (previousLine[0:3] != "## ")
                    & (previousLine[0:4] != "    ")
                    & (previousLine[0] != "|")
                    & (inCode is False)
                    & (inHTMLCode is False)
                ):
                    # Previous line was not Heading Level 1 or 2 and we're not in code so concatenate
                    linesAfterConcatenation[-1] = (
                        linesAfterConcatenation[-1].rstrip() + " " + line.lstrip()
                    )
                else:
                    linesAfterConcatenation.append(line)

            else:
                linesAfterConcatenation.append(line)

        # Store previous line to see if it was H1 or blank
        previousLine = line

    # Pass 3: Get footnote definitions
    metadataLinenumber = 0
    for line in linesAfterConcatenation:
        line = line.rstrip()
        if m := footnoteDefinitionRegex.match(line):
            fnRef = m.group(1).strip()
            fnText = m.group(2).strip()
            footnoteDefinitions.append([fnRef, fnText])
            footnoteReferences.append(fnRef)

            linesAfterConcatenation[metadataLinenumber] = "<ignoreme>"
        metadataLinenumber += 1

    # Pass 4: Extract any indirect reference anchors
    metadataLinenumber = 0
    indirectAnchors = []
    for line in linesAfterConcatenation:
        line = line.rstrip()
        if m := indirectReferenceAnchorRegex.match(line):
            anchorName = m.group(1).strip()
            anchorURL = m.group(2).strip()

            indirectAnchors.append([anchorName, anchorURL])

            linesAfterConcatenation[metadataLinenumber] = "<ignoreme>"
        metadataLinenumber += 1

    lastTableLine = -1
    tableCaptions = []

    # Pass 5: Main pass over the input file, now that footnote
    # references have been gathered
    for lineNumber, line in enumerate(linesAfterConcatenation):
        # Remove trailing white space
        line = line.rstrip()

        # Convert tabs to spaces
        line = line.replace("\t", " " * processingOptions.getCurrentOption("indentSpaces"))

        if line == "<ignoreme>":
            # Line was taken care of in the previous pass
            continue

        if startswithOneOf(line, ["<pre>", "<code>"]):
            code.append([])
            inCode = True
            inHTMLCode = True
            inTable = False
            inTitle = False

        if startswithOneOf(line, ["</pre>", "</code>"]):
            inCode = False
            inHTMLCode = False
            inTitle = False

        if line.startswith("```"):
            inCode = not inCode
            if inCode:
                # Just entered code
                code.append([])
                blockType = "code"
                inTable = False
            else:
                # Just exited code - but add closing line
                code[-1].append(line)
            inFencedCode = not inFencedCode
            inTitle = False

        if inCode or inHTMLCode or inFencedCode:
            if len(code) == 0:
                code.append([])
            code[-1].append(line)

            # If first line of code then mark the current sequence entry as "code"
            if len(code[-1]) == 1:
                sequence.append("code")

            inTitle = False
        if (
            (line == "")
            & (inCode is True)
            & (inHTMLCode is False)
            & (inFencedCode is False)
        ):
            inCode = False
            inTitle = False

        if (line.startswith("    ")) & (inList is False):
            # Only list items and code can be indented by 4 characters
            if inCode is False:
                code.append([])
                blockType = "code"
                code[-1].append(line[4:])
                inCode = True
            inTitle = False

        # Rewrite horizontal rule as a heading 3 with non-breaking space
        if startswithOneOf(line, ["<hr/>", "---", "***", "___"]):
            line = "### &nbsp;"
            inTitle = True

        # Taskpaper task
        if line[:1] == "-":
            # Get start of attributes
            attributesStart = line.find("@")

            # Get text up to attributes
            if attributesStart == -1:
                text = line[2:]
            else:
                text = line[2 : attributesStart - 1]

            # Attempt to extract @due information
            startDue = line.find("@due(")
            if startDue > -1:
                startDue += 5
                endDue = line.find(")", startDue)
                if endDue > -1:
                    dueDate = line[startDue:endDue]
            else:
                dueDate = ""

            # Attempt to extract @tags information
            startTags = line.find("@tags(")
            if startTags > -1:
                startTags += 6
                endTags = line.find(")", startTags)
                if endTags > -1:
                    tags = line[startTags:endTags]
            else:
                tags = ""

            # Attempt to extract @done information
            startDone = line.find("@done(")
            if startDone > -1:
                startDone += 6
                endDone = line.find(")", startDone)
                if endDone > -1:
                    done = line[startDone:endDone]
            else:
                done = ""

            tasks.append([slideNumber + 1, text, dueDate, tags, done])
            inTitle = False

        elif line.startswith("<a id="):
            # Anchor on whatever slide we're on
            if hrefMatch := anchorRegex.match(line):
                href = hrefMatch.group(1)
                if (href != "") & (href in slideHrefs.keys()):
                    print(f"Heading Reference redefined: '{href}' for slide {slideNumber}")
            inTitle = False

        elif line.startswith("<!-- md2pptx: "):
            # Dynamic metadata line
            inTitle = False
            if DMMatch := dynamicMetadataRegex.match(line):
                metadataKey = DMMatch.group(1).lower().strip()
                metadataValue = DMMatch.group(2).lower()
                if (metadataKey != "") & (metadataValue != ""):
                    # Valid metadata pair so parse key / value - and apply if valid

                    if (metadataKey == "pagesubtitlesize") & (metadataValue == "same"):
                        # Cards' layout - horizontal or vertical
                        processingOptions.dynamicallySetOption(
                            metadataKey, metadataValue, ""
                        )

                    # Floating point values where metadata key matches directly
                    elif metadataKey in [
                        # Font size for tables
                        "compacttables",
                        # Heading font size for tables
                        "tableheadingsize",
                        # Page Title Font size
                        "pagetitlesize",
                        # Page Subtitle Font size
                        "pagesubtitlesize",
                        # Cards' vertical percent of content area
                        "cardpercent",
                        # Card graphic's vertical space
                        "cardgraphicheight",
                        # Base text size
                        "basetextsize",
                        # Base text decrement
                        "basetextdecrement",
                        # Horizontal card gap
                        "horizontalcardgap",
                        # Vertical card gap
                        "verticalcardgap",
                        # Funnel label percentage
                        "funnellabelspercent",
                    ]:
                        processingOptions.dynamicallySetOption(
                            # Use the metadata key directly
                            metadataKey,
                            metadataValue.lower(),
                            "float",
                        )

                    elif metadataKey == "cardlayout":
                        # Cards' layout - horizontal or vertical
                        processingOptions.dynamicallySetOption(
                            metadataKey, metadataValue, ""
                        )

                    elif metadataKey == "numbersheight":
                        # vertical space to reserve for slide number
                        if metadataValue == "default":
                            numbersheight = processingOptions.getDefaultOption(
                                "numbersHeight"
                            )
                        elif metadataValue == "pres":
                            numbersheight = processingOptions.getPresentationOption(
                                "numbersHeight"
                            )
                        else:
                            numbersheight = int(Inches(float(metadataValue)))

                        processingOptions.dynamicallySetOption(
                            metadataKey, numbersheight, "int"
                        )

                    elif metadataKey == "marginbase":
                        #  space to reserve as a margin
                        if metadataValue == "default":
                            marginbase = processingOptions.getDefaultOption("marginbase")
                        elif metadataValue == "pres":
                            marginbase = processingOptions.getPresentationOption(
                                "marginbase"
                            )
                        else:
                            marginbase = int(Inches(float(metadataValue)))

                        processingOptions.dynamicallySetOption(
                            metadataKey, marginbase, "int"
                        )

                    elif metadataKey == "tablemargin":
                        #  space to reserve as a table margin
                        if metadataValue == "default":
                            tablemargin = processingOptions.getDefaultOption("tablemargin")
                        elif metadataValue == "pres":
                            tablemargin = processingOptions.getPresentationOption(
                                "tablemargin"
                            )
                        else:
                            tablemargin = int(Inches(float(metadataValue)))

                        processingOptions.dynamicallySetOption(
                            metadataKey, tablemargin, "int"
                        )

                    elif metadataKey == "spancells":
                        #  whether table cells can span more than one column
                        if metadataValue == "default":
                            spancells = processingOptions.getDefaultOption("spancells")
                        elif metadataValue == "pres":
                            spancells = processingOptions.getPresentationOption("spancells")
                        else:
                            spancells = metadataValue

                        processingOptions.dynamicallySetOption(metadataKey, spancells, "")

                    # Note: Actual value handling prevents using dynamicallySetOption
                    elif metadataKey == "cardtitlealign":
                        # Card title alignment
                        processingOptions.dynamicallySetOption(
                            metadataKey, metadataValue.lower(), ""
                        )

                    elif metadataKey == "cardtitleposition":
                        # Card title position - above or inside
                        processingOptions.dynamicallySetOption(
                            metadataKey, metadataValue, ""
                        )

                    elif metadataKey == "cardgraphicposition":
                        # Card graphic position - before or after
                        processingOptions.dynamicallySetOption(
                            metadataKey, metadataValue, ""
                        )

                    elif metadataKey == "cardshape":
                        # Card shape
                        if metadataValue in ["squared", "rounded", "line"]:
                            processingOptions.dynamicallySetOption(
                                metadataKey, metadataValue, ""
                            )
                        else:
                            print(
                                f'CardShape value \'{metadataValue}\' unsupported. "squared", "rounded", or "line" required.'
                            )

                    elif metadataKey in [
                        "cardcolour",
                        "cardcolor",
                        "cardcolours",
                        "cardcolors",
                    ]:
                        valueArray2 = [
                            parseColour(c.strip()) for c in metadataValue.split(",")
                        ]

                        processingOptions.dynamicallySetOption(
                            "cardColour", valueArray2, ""
                        )

                    elif metadataKey in [
                        "cardtitlebackground",
                        "cardtitlebackgrounds",
                    ]:
                        valueArray2 = [
                            parseColour(c.strip()) for c in metadataValue.split(",")
                        ]

                        processingOptions.dynamicallySetOption(
                            "cardTitleBackground", valueArray2, ""
                        )

                    elif metadataKey in ["funnelcolours", "funnelcolors"]:
                        valueArray2 = [
                            parseColour(c.strip()) for c in metadataValue.split(",")
                        ]

                        processingOptions.dynamicallySetOption(
                            "funnelColours", valueArray2, ""
                        )

                    elif metadataKey in [
                        "funneltitlecolour",
                        "funnelbordercolour",
                        "funneltextcolour",
                    ]:
                        # Funnel single colour options
                        processingOptions.dynamicallySetOption(
                            metadataKey, metadataValue, ""
                        )

                    # Note: Actual value handling prevents using dynamicallySetOption
                    elif metadataKey == "contentsplit":
                        # Proportions for each content element on a slide
                        if metadataValue == "default":
                            contentSplit = processingOptions.getDefaultOption(
                                "contentSplit"
                            )
                        elif metadataValue == "pres":
                            contentSplit = processingOptions.getPresentationOption(
                                "contentSplit"
                            )
                        elif metadataValue in ["pop", "prev"]:
                            processingOptions.popCurrentOption("contentSplit")
                            contentSplit = processingOptions.getPresentationOption(
                                "contentSplit"
                            )
                        else:
                            splitValue = metadataValue.split()
                            contentSplit = []
                            for v in splitValue:
                                contentSplit.append(int(v))

                            # Extend to maximum allowed
                            needMore = maxBlocks - len(contentSplit)
                            for _ in range(needMore):
                                contentSplit.append(1)
                        processingOptions.dynamicallySetOption(
                            "contentSplit", contentSplit, ""
                        )

                    elif metadataKey in ["contentsplitdirection", "contentsplitdirn"]:
                        if metadataValue in [
                            "vertical",
                            "horizontal",
                            "v",
                            "h",
                            "pres",
                            "default",
                            "pop",
                            "prev",
                        ]:
                            if metadataValue == "v":
                                adjustedValue = "vertical"
                            elif metadataValue == "h":
                                adjustedValue = "horizontal"
                            else:
                                adjustedValue = metadataValue

                            processingOptions.dynamicallySetOption(
                                "contentSplitDirection",
                                adjustedValue,
                                "",
                            )

                        else:
                            print(
                                f'{metadataKey} value \'{metadataValue}\' unsupported. "vertical" or "horizontal" required.'
                            )

                    elif metadataKey in ["codeforeground", "codebackground"]:
                        processingOptions.dynamicallySetOption(
                            metadataKey,
                            metadataValue,
                            "",
                        )

                    elif metadataKey == "fpratio":
                        # Fixed Pitch font height to width ratio
                        processingOptions.dynamicallySetOption(
                            "fixedPitchHeightWidthRatio",
                            metadataValue,
                            "float",
                        )

                    elif metadataKey == "codecolumns":
                        processingOptions.dynamicallySetOption(
                            metadataKey,
                            metadataValue,
                            "int",
                        )

                    elif metadataKey == "indentspaces":
                        # Spaces representing a single indentation level
                        processingOptions.dynamicallySetOption(
                            metadataKey,
                            metadataValue,
                            "int",
                        )

                    elif metadataKey in ["addtablecolumnlines", "addtablerowlines"]:
                        processingOptions.dynamicallySetOption(
                            metadataKey,
                            metadataValue,
                            "sortednumericlist",
                        )

                    elif metadataKey in ["addtablelinecount", "addtablelinewidth"]:
                        processingOptions.dynamicallySetOption(
                            metadataKey,
                            metadataValue,
                            "int",
                        )

                    elif metadataKey in ["addtablelines", "addtablelinecolour"]:
                        processingOptions.dynamicallySetOption(
                            metadataKey,
                            metadataValue,
                            "",
                        )

                    elif metadataKey == "transition":
                        if metadataValue.lower() in [
                            "none",
                            "ripple",
                            "reveal",
                            "honeycomb",
                            "shred",
                            "wipe",
                            "push",
                            "vortex",
                            "split",
                            "fracture",
                            "pres",
                            "default",
                            "pop",
                            "prev",
                        ]:
                            processingOptions.dynamicallySetOption(
                                metadataKey,
                                metadataValue,
                                "",
                            )
                        else:
                            print(
                                f"Invalid value for {metadataKey}: {metadataValue} in '{line}"
                            )

                    elif metadataKey == "funnellabelposition":
                        if metadataValue.lower() in [
                            "before",
                            "after",
                        ]:
                            processingOptions.dynamicallySetOption(
                                metadataKey,
                                metadataValue,
                                "",
                            )
                        else:
                            print(
                                f"Invalid value for {metadataKey}: {metadataValue} in '{line}"
                            )

                    elif metadataKey == "funnelwidest":
                        if metadataValue.lower() in [
                            "left",
                            "right",
                            "pipe",
                            "hpipe",
                            "vpipe",
                            "top",
                            "bottom",
                        ]:
                            processingOptions.dynamicallySetOption(
                                metadataKey,
                                metadataValue,
                                "",
                            )
                        else:
                            print(
                                f"Invalid value for {metadataKey}: {metadataValue} in '{line}"
                            )

                    else:
                        # Invalid dynamic metadata specification
                        print(f"Invalid dynamic metadata key: '{metadataKey}' in '{line}'")

        elif line.startswith("#"):
            if line[:cardLevel] == "#" * cardLevel:
                # Card on an existing slide
                inCard = True

                # If there is an outstanding href then store it
                if href != "":
                    slideHrefs[href] = slideNumber + 1

                # Get card title text and any href
                cardName, href = parseTitleText(line[cardLevel:])

                # No bullets for the card yet
                cardBullets = []

                # No graphic for the card yet
                cardGraphic = ""

                # Create information about a card_title
                cards.append([cardName, cardGraphic, cardBullets])

                # Might have octothorpes but isn't a title
                inTitle = False


            else:
                # One of Content, Section, or Title slide
                if inBlock is True:
                    # Create the previous slide
                    slideInfo = SlideInfo(
                        slideTitle,
                        slideSubtitle,
                        blockType,
                        bullets,
                        tableRows,
                        cards,
                        code,
                        sequence,
                    )
                    slideNumber, slide, sequence = createSlide(prs, slideNumber, slideInfo)

                    # Register the previous slide's href - if there was one
                    if href != "":
                        slideHrefs[href] = slideNumber

                if line[:contentLevel] == "#" * contentLevel:
                    # Heading Level 3 - slide
                    thisLevel = contentLevel
                    blockType = "content"

                elif line[:sectionLevel] == "#" * sectionLevel:
                    # Heading Level 2 - section
                    thisLevel = sectionLevel
                    blockType = "section"
                    inTitle = True

                elif line[:titleLevel] == "#" * titleLevel:
                    # Heading Level 1 - slide title
                    thisLevel = titleLevel
                    blockType = "title"
                    inTitle = True
                else:
                    inTitle = False

                # Get slide title text and any href
                slideTitle, href = parseTitleText(line[thisLevel:])

                inBlock = True
                inList = False
                inTable = False
                inCard = False
                slideSubtitle = ""
                bullets = []
                tableRows = []
                cards = []
                code = []

                if (notes_text != "") & (slide != None):
                    createSlideNotes(slide, notes_text)

                notes_text = ""

            # Check whether any heading reference is a duplicate
            if (href != "") & (href in slideHrefs.keys()):
                print(f"Heading Reference redefined: '{href}' for slide {slideNumber}")

        elif match := bulletRegex.match(line):
            # Bulleted list
            bulletLine = match.group(3).lstrip()

            bulletLevel = calculateIndentationLevel(
                match.start(2), processingOptions.getCurrentOption("indentSpaces")
            )

            bulletType = "bulleted"

            if inCard:
                # Bullet is on a card
                cardBullets.append([bulletLevel, bulletLine, bulletType])

                # Update bullet list for the card
                cards[-1][2] = cardBullets
            else:
                # Bullet is not on a card
                bullets.append([bulletLevel, bulletLine, bulletType])

            if inList is False:
                sequence.append("list")

            inList = True
            inTable = False
            inTitle = False

        elif match := numberRegex.match(line):
            # Numbered list
            bulletLine = match.group(3).lstrip()

            bulletLevel = calculateIndentationLevel(
                match.start(2), processingOptions.getCurrentOption("indentSpaces")
            )

            bulletType = "numbered"

            if inCard:
                # Bullet is on a card
                cardBullets.append([bulletLevel, bulletLine, bulletType])

                # Update bullet list for the card
                cards[-1][2] = cardBullets
            else:
                # Bullet is not on a card
                bullets.append([bulletLevel, bulletLine, bulletType])

            if inList is False:
                sequence.append("list")

            inList = True
            inTable = False
            inTitle = False

        # Following marshals media into one or two table cells and adds the new row
        # to either a new or existing table. It also says we're in a table
        elif startswithOneOf(line, ["<video ", "<audio ", "[![", "!["]):
            if inCard:
                # Graphic in a card
                if startswithOneOf(line, ["<video ", "<audio ", "[!["]):
                    # Might implement these later on
                    print("video, audio and clickable graphic not supported")
                else:
                    # graphicCount
                    for m in re.finditer(graphicRegex, line):
                        # A non-clickable graphic
                        cardGraphic = m.group(0)[4 : -1]

                        # Add graphic to latest card
                        cards[-1][1] = cardGraphic
            else:
                # There is at least one media item in the line and not in a card
                # - so set up table
                if inTable is False:
                    # Switch to being in a table
                    blockType = "table"
                    sequence.append("table")
                    inTable = True
                    inList = False
                    inCard = False

                    # Create a new empty table
                    # tableRows = []

                # Start the creation of a new row
                tableRow = []

                # Collect the media items
                mediaItems = []
                for m in re.finditer(videoRegex, line):
                    # A video
                    mediaItem = [m.start(0), m.end(0), m.group(0)]
                    mediaItems.append(mediaItem)

                for m in re.finditer(audioRegex, line):
                    # An audio file
                    mediaItem = [m.start(0), m.end(0), m.group(0)]
                    mediaItems.append(mediaItem)

                for m in re.finditer(clickableGraphicRegex, line):
                    # A clickable graphic
                    mediaItem = [m.start(0), m.end(0), m.group(0)]
                    mediaItems.append(mediaItem)

                for m in re.finditer(graphicRegex, line):
                    # A non-clickable graphic
                    start = m.start(0)
                    end = m.end(0)
                    mediaItem = [start, end, m.group(0)]

                    rematch = False
                    for existingMediaItem in mediaItems:
                        if (start >= existingMediaItem[0]) & (end <= existingMediaItem[1]):
                            # This graphic already found as part of a clickable graphic link
                            rematch = True
                            break

                    if rematch is False:
                        # This graphic not already found as part of a clickable graphic link
                        mediaItems.append(mediaItem)

                # compose table row - based on start offset
                tableRow = []
                for mediaItem in sorted(mediaItems):
                    tableRow.append(mediaItem[2])

                # Add table row to saved table rows
                if len(tableRows) == 0:
                    tableRows.append([])
                tableRows[-1].append(tableRow)

                inTitle = False

        elif line[:1] == "|":
            # Table or side-by-side
            lastTableLine = lineNumber

            # As we're in a table we can't have a table caption yet
            tableCaption = ""

            if inTable is False:
                blockType = "table"
                tableRows.append([])
                inTable = True
                sequence.append("table")
                inList = False
                inCard = False

            # Create a table row - but with (maybe empty) junk before and after
            words = line.split("|")
            tableRow = []
            for cell in words:
                tableRow.append(cell)

            # Remove first element
            tableRow.pop(0)

            # Remove last element - if blank
            if cell == "":
                tableRow.pop()

            # Add clean table row to saved table rows
            tableRows[-1].append(tableRow)
            inTitle = False

        elif (
            (line.startswith("[")) & line.endswith("]") & (lineNumber == lastTableLine + 1)
        ):
            tableCaption = line[1:-1]
            tableCaptions.append(tableCaption)
        else:
            # Not in a table
            inTable = False

            if len(tableCaptions) < len(tableRows):
                tableCaptions.append("")

            if not inCode:
                # Must be a slide note line or a subtitle line
                if line == "":
                    inTitle = False
                if inTitle:
                    slideSubtitle = slideSubtitle + "\n" + line
                elif startswithOneOf(line, ["</pre>", "</code>"]) is False:
                    notes_text = notes_text + "\n" + line

    # Ensure there's a blank table caption - for the case the table ended in the final line
    if len(tableCaptions) < len(tableRows):
        tableCaptions.append(tableCaption)

    ######################################################################################
    #                                                                                    #
    # Finish off last slide                                                              #
    #                                                                                    #
    ######################################################################################
    if (inBlock is True) | (inCode is True) | (inTable is True):
        slideInfo = SlideInfo(
            slideTitle, slideSubtitle, blockType, bullets, tableRows, cards, code, sequence
        )
        slideNumber, slide, sequence = createSlide(prs, slideNumber, slideInfo)

        if href != "":
            slideHrefs[href] = slideNumber

        if (notes_text != "") & (slide != None):
            createSlideNotes(slide, notes_text)

        notes_text = ""

    ######################################################################################
    #                                                                                    #
    # Add a footnotes slide - if there were any footnote definitions                     #
    #                                                                                    #
    ######################################################################################
    if len(footnoteDefinitions) > 0:
        slideNumber, footnoteSlides = createFootnoteSlides(
            prs, slideNumber, footnoteDefinitions
        )

        footnotesPerPage = processingOptions.getCurrentOption("footnotesPerPage")
        # Fix up any footnote slide hyperlinks
        footnoteNumber = -1
        for footnoteRun in footnoteRunsDictionary.keys():
            footnoteNumber += 1
            run = footnoteRunsDictionary[footnoteRun]

            footnoteSlideNumber = int(footnoteNumber / footnotesPerPage)
            createRunHyperlinkOrTooltip(run, footnoteSlides[footnoteSlideNumber], "")

    ######################################################################################
    #                                                                                    #
    # Add a dictionary slide - if there were any abbr elements encountered               #
    #                                                                                    #
    ######################################################################################
    if len(abbrevDictionary) > 0:
        glossaryTermsPerPage = processingOptions.getCurrentOption("glossaryTermsPerPage")
        slideNumber, glossarySlides = createGlossarySlides(
            prs, slideNumber, abbrevDictionary
        )
        # Fix up internal glossary hyperlinks
        abbrevNumber = -1
        for abbreviation in sorted(abbrevRunsDictionary.keys()):
            abbrevNumber += 1
            runs = abbrevRunsDictionary[abbreviation]
            for run in runs:
                # Add tooltip for glossary definition
                glossarySlideNumber = int(abbrevNumber / glossaryTermsPerPage)
                createRunHyperlinkOrTooltip(
                    run, glossarySlides[glossarySlideNumber], abbrevDictionary[abbreviation]
                )

    ######################################################################################
    #                                                                                    #
    # Add final slide - or more than one - with any Taskpaper tasks in                   #
    #                                                                                    #
    ######################################################################################
    taskSlides = processingOptions.getCurrentOption("taskSlides")

    if (len(tasks) > 0) & (taskSlides != "none"):
        # Turn tasks into a table slide

        # Might need to winnow slides
        if taskSlides != "all":
            complete = []
            incomplete = []
            for task in tasks:
                sNum, taskText, dueDate, tags, done = task
                if done == "":
                    incomplete.append(task)
                else:
                    complete.append(task)

            if (taskSlides == "separate") & (len(tasks) > 0):
                want_task_slides = True
            elif (taskSlides == "remaining") & (len(incomplete) > 0):
                want_task_slides = True
            elif (taskSlides == "done") & (len(complete) > 0):
                want_task_slides = True
            else:
                want_task_slides = False
        else:
            want_task_slides = True

        if want_task_slides:
            if taskSlides != "separate":
                createTaskSlides(prs, slideNumber, tasks, "Tasks")
            else:
                createTaskSlides(prs, slideNumber, complete, "Completed Tasks")
                createTaskSlides(prs, slideNumber, incomplete, "Incomplete Tasks")

    ######################################################################################
    #                                                                                    #
    # Make any TOC / Section-related links                                               #
    #                                                                                    #
    ######################################################################################
    if processingOptions.getCurrentOption("tocLinks"):
        # Linkify section items
        for run in TOCruns:
            createRunHyperlinkOrTooltip(run, SectionSlides[run.text])

    if processingOptions.getCurrentOption("sectionArrows"):
        # Add navigation arrows between section slides
        sectionArrowsColour = processingOptions.getCurrentOption("sectionArrowsColour")

        buttonTop = prs.slide_height - Inches(2 / 3)
        forwShape = None

        previousSection = None
        for sectionNumber, sectionSlide in enumerate(SectionSlides):
            slide = SectionSlides[sectionSlide]
            if sectionNumber == 0:
                TOCslide = slide

            if forwShape != None:
                createShapeHyperlinkAndTooltip(forwShape, slide, "Next Section")

            buttonShapes = []
            if sectionNumber > 1:
                # Need backwards arrow
                backShape = slide.shapes.add_shape(
                    MSO_SHAPE.ACTION_BUTTON_BACK_OR_PREVIOUS,
                    prs.slide_width / 2 - Inches(2 / 3),
                    buttonTop,
                    Inches(1 / 3),
                    Inches(1 / 3),
                )

                createShapeHyperlinkAndTooltip(backShape, previousSlide, "Previous Section")

                buttonShapes.append(backShape)

            # Always need home arrow - except for TOC
            if sectionNumber > 0:
                homeShape = slide.shapes.add_shape(
                    MSO_SHAPE.ACTION_BUTTON_HOME,
                    prs.slide_width / 2 - Inches(1 / 6),
                    buttonTop,
                    Inches(1 / 3),
                    Inches(1 / 3),
                )

                createShapeHyperlinkAndTooltip(homeShape, TOCslide, "Table Of Contents")

                buttonShapes.append(homeShape)

            if (sectionNumber < len(SectionSlides) - 1) & (sectionNumber > 0):
                # Need forwards
                forwShape = slide.shapes.add_shape(
                    MSO_SHAPE.ACTION_BUTTON_FORWARD_OR_NEXT,
                    prs.slide_width / 2 + Inches(1 / 3),
                    buttonTop,
                    Inches(1 / 3),
                    Inches(1 / 3),
                )

                buttonShapes.append(forwShape)

            else:
                forwShape = None

            # Fix background colour of the buttons
            if sectionArrowsColour != "":
                for buttonShape in buttonShapes:
                    buttonShape.fill.solid()
                    buttonShape.fill.fore_color.rgb = RGBColor.from_string(
                        sectionArrowsColour
                    )

            previousSlide = slide


    ######################################################################################
    #                                                                                    #
    # Fix up any internal links                                                          #
    #                                                                                    #
    ######################################################################################
    xrefCheck_errors = False
    for href in href_runs.keys():
        run = href_runs[href]
        if href in slideHrefs.keys():
            createRunHyperlinkOrTooltip(
                run, prs.slides[slideHrefs[href] - 2 + templateSlideCount], ""
            )
        else:
            # No id defined with that name
            if not xrefCheck_errors:
                # First time in this run a cross reference error occurred
                xrefCheck_errors = True
                print("\nHyperlink Reference Errors")
                print("--------------------------")

            print(
                "Slide "
                + str(prs.slides.index(SlideFromRun(run)) + 1 - templateSlideCount)
                + f": '{href}'"
            )

    # Each picture appears in pictures
    # There's a corresponding entry in picture_Hrefs
    # There's a corresponding entry in picture_tooltips

    # fix up any clickable picture links
    picture_count = len(pictureInfos)
    for p in range(picture_count):
        picture, href, tooltip = pictureInfos[p]

        # Pick up link target - if any
        if href == "#XYZZY-None":
            target = None
        else:
            if href[1:] in slideHrefs.keys():
                # Is an internal link
                target = prs.slides[slideHrefs[href[1:]] - 2 + templateSlideCount]
            else:
                # Is an external Link
                target = href

        createPictureHyperlinkOrTooltip(picture, target, tooltip)

    # if templateSlideCount > 0:
        # createProcessingSummarySlide(prs, metadata)

    try:
        if processingOptions.getCurrentOption("deletefirstslide"):
            # Remove first slide
            deleteSlide(prs, 0)

    except:
        pass

    if processingOptions.getCurrentOption("SectionsExpand"):
        createExpandingSections(prs)

    prs.save(output_filename)

    elapsed_time = time.time() - start_time

    print(
        "\nProcessing complete.\nElapsed time: "
        + str(int(1000 * elapsed_time) / 1000)
        + "s"
    )

    # Run a script against every slide
    script = "if slide.slideInfo is not None:\n  print(slide.slideInfo)"
    script = ""
    for slide in prs.slides:
        exec(script)
    return outputfile
